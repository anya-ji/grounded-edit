from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5] # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = "Global and Local Cultural Products"
title_frame.paragraphs[0].font.size = Pt(40)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(0, 0, 255) # Blue background
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255) # White text

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(15), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Cultural Products"
subtitle_frame.paragraphs[0].font.size = Pt(28)
subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0) # Black text

# Bullet Points
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(15), Inches(2))
content_frame = content_box.text_frame
content_frame.word_wrap = True

p = content_frame.add_paragraph()
p.text = "✓ are goods and services such as arts, architectures, museums etc. that showcase the history and information of certain which belong to the country’s cultural heritage."
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(0, 0, 0) # Black text
p.space_after = Pt(14)  # Adds space after the paragraph for better readability
content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(255, 255, 255) # White background

# Images
image_paths = [
    "../../media/image_0.jpg",
    "../../media/image_1.jpg",
    "../../media/image_2.jpg"
]

x_positions = [Inches(1), Inches(5.5), Inches(10)]
for i, image_path in enumerate(image_paths):
    slide.shapes.add_picture(image_path, x_positions[i], Inches(5), Inches(4), Inches(3))

presentation.save("render.pptx")