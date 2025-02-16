from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set lighter background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 240, 240)  # Lightened background color

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(5), Inches(1))
title_frame = title_box.text_frame
title_p = title_frame.add_paragraph()
title_p.text = "PART 3"
title_p.font.bold = True
title_p.font.size = Pt(48)
title_p.font.color.rgb = RGBColor(0, 0, 0)
title_p.alignment = PP_ALIGN.LEFT

# Add instructions with check logo
instruction_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(12), Inches(1))
instruction_frame = instruction_box.text_frame
instruction_frame.word_wrap = True
instruction_p = instruction_frame.add_paragraph()
instruction_p.text = "Choose true (T), false (F), or not given (N) according to the information in the video."
instruction_p.font.size = Pt(24)
instruction_p.font.color.rgb = RGBColor(0, 0, 0)

# Add green check image
check_image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(check_image_path, Inches(11), Inches(1.5), height=Inches(0.5))

# Add table for statements with checkboxes
table = slide.shapes.add_table(rows=3, cols=4, left=Inches(1), top=Inches(3), width=Inches(14), height=Inches(2)).table

# Set column widths
table.columns[0].width = Inches(10)
table.columns[1].width = Inches(1.5)
table.columns[2].width = Inches(1.5)
table.columns[3].width = Inches(1.5)

# Add statements
statements = [
    "Phoebe sarcastically says that Sarah is a monster for her etiquette.",
    "Joey prefers eating French fries with his fingers.",
    "Rachel is surprised to hear about Joey’s food sharing rule."
]

for i, statement in enumerate(statements):
    cell = table.cell(i, 0)
    cell.text = statement
    cell.text_frame.paragraphs[0].font.size = Pt(20)
    cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

# Add checkboxes headers: T, F, N
for i in range(1, 4):
    table.cell(0, i).text = '✓' if i == 1 else ''
    table.cell(1, i).text = '✓' if i == 2 else ''
    table.cell(2, i).text = '✓' if i == 3 else ''

# Update the first statement to have a checkmark
table.cell(0, 1).text = '✓'  # Checkmark for the first statement

# Add Joey's image
joey_image_path = "../../media/image_1.jpg"
joey_image = slide.shapes.add_picture(joey_image_path, Inches(11), Inches(5.5), height=Inches(2))

# Add buttons below Joey's image
button_check = slide.shapes.add_textbox(Inches(11), Inches(7.7), Inches(2.5), Inches(0.7))
button_check.fill.solid()
button_check.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Transparent background
check_p = button_check.text_frame.add_paragraph()
check_p.text = "Check your answers"
check_p.font.size = Pt(18)
check_p.font.color.rgb = RGBColor(255, 20, 147)  # Pink
check_p.alignment = PP_ALIGN.CENTER
button_check.line.color.rgb = RGBColor(255, 20, 147)  # Pink border

button_view = slide.shapes.add_textbox(Inches(13.7), Inches(7.7), Inches(2.5), Inches(0.7))
button_view.fill.solid()
button_view.fill.fore_color.rgb = RGBColor(0, 0, 139)  # Dark blue background
view_p = button_view.text_frame.add_paragraph()
view_p.text = "Viewing Activity"
view_p.font.size = Pt(18)
view_p.font.color.rgb = RGBColor(255, 255, 255)  # White
view_p.alignment = PP_ALIGN.CENTER

# Save presentation
presentation.save("render.pptx")