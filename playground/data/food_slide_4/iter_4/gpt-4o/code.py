from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation and slide
presentation = Presentation()
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)  # Changed to white background

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
p = title_frame.add_paragraph()
p.text = "PART 2"
p.font.bold = True
p.font.size = Pt(60)  # Increased font size for better visibility
p.font.color.rgb = RGBColor(0, 0, 0)  # Dark font

# Add subtitle
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(1))
subtitle_frame = subtitle_box.text_frame
p = subtitle_frame.add_paragraph()
p.text = "Discuss the questions below about dating."
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(0, 0, 0)

# List the questions with numbers in circles
questions = [
    "What do you think are some other good (or bad) date ideas not listed on the previous slide?",
    "What are some things that could go wrong on a first date?",
    "What are some common etiquette mistakes to avoid on a first date?"
]

for i, question in enumerate(questions, start=1):
    # Add circle with number
    left = Inches(1)
    top = Inches(2.5 + i * 0.8)
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(0.5), Inches(0.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 102, 204)  # Blue color
    text_frame = shape.text_frame
    p = text_frame.paragraphs[0]
    p.text = str(i)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)  # White font

    # Add question text
    left = Inches(1.8)
    question_box = slide.shapes.add_textbox(left, top, Inches(13), Inches(0.5))
    question_frame = question_box.text_frame
    p = question_frame.add_paragraph()
    p.text = question
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)

# Add "PREVIEW ACTIVITY" button
button_left = Inches(12)
button_top = Inches(7.5)
button = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, button_left, button_top, Inches(3), Inches(0.8))
button.fill.solid()
button.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue
button_text_frame = button.text_frame
p = button_text_frame.add_paragraph()
p.text = "PREVIEW ACTIVITY"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(255, 255, 255)

# Include a communication graphic
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(14), Inches(3), Inches(2), Inches(2))

# Add a decorative element on the right side
decorative_shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(14), Inches(4), Inches(1.5), Inches(1))
decorative_shape.fill.solid()
decorative_shape.fill.fore_color.rgb = RGBColor(255, 223, 186)  # Light color for decorative element

# Save presentation
presentation.save("render.pptx")