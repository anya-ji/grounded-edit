from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Choose a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 240, 240)

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "STEP 2:"
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.size = Pt(32)

# Add bullet points
bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(15), Inches(3))
bullet_frame = bullet_box.text_frame
bullet_points = [
    "Imagine you actually went on the date you chose on the previous slide and something went wrong.",
    "Describe what happened on the date using the past simple and past continuous.",
    "Take a few minutes to prepare your story about the date below. Write 3 - 5 sentences."
]
for point in bullet_points:
    p = bullet_frame.add_paragraph()
    p.text = point
    p.space_after = Pt(10)
    p.level = 0
bullet_frame.word_wrap = True
bullet_frame.paragraphs[0].font.size = Pt(20)

# Add text box for story
story_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(10), Inches(2))
story_frame = story_box.text_frame
story_frame.text = "YOUR ‘DATE GONE WRONG’ STORY\n\n"
story_frame.paragraphs[0].font.bold = True
story_frame.paragraphs[0].font.size = Pt(20)
story_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Set padding (margins)
story_frame.margin_left = Pt(20)  # Left padding
story_frame.margin_right = Pt(20)  # Right padding
story_frame.margin_top = Pt(20)    # Top padding
story_frame.margin_bottom = Pt(20)  # Bottom padding

# Add image beside text box
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(11), Inches(5), height=Inches(2))

# Add additional section
step3_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(15), Inches(1.5))
step3_frame = step3_box.text_frame
step3 = "STEP 3: Retell the story to your class or teacher. Use stress/intonation in the appropriate places when retelling the story."
p = step3_frame.add_paragraph()
p.text = step3
p.font.bold = True
p.font.size = Pt(20)
step3_frame.word_wrap = True

# Save the presentation
presentation.save("render.pptx")