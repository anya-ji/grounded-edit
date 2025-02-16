from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

# Initialize presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5]  
slide = presentation.slides.add_slide(slide_layout)

# Add background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add white overlay
overlay = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(15), Inches(5))
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
overlay.fill.transparency = 0.3
overlay.line.color.rgb = RGBColor(255, 255, 255)

# Add additional title
additional_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(5), Inches(0.5))
additional_title_frame = additional_title_box.text_frame
additional_title = additional_title_frame.add_paragraph()
additional_title.text = "The Subject and Content of Art"
additional_title.font.size = Pt(20)
additional_title.font.bold = True
additional_title.font.color.rgb = RGBColor(0, 0, 0)
additional_title.alignment = PP_ALIGN.LEFT

# Add main title
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1))
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(255, 255, 0)
title_frame = title_box.text_frame
title_frame.word_wrap = True
title = title_frame.add_paragraph()
title.text = "Two kinds of Art as to Subject"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title.alignment = PP_ALIGN.CENTER

# Add main point
content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(6))
text_frame = content_box.text_frame

# Add main point
title_paragraph = text_frame.add_paragraph()
title_paragraph.text = "2. Non-representational or Non-objective Art"
title_paragraph.font.size = Pt(24)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(0, 0, 0)

# Add bullet points
bullet_points = [
    "✓ Also known as non-objective art",
    "✓ The subject is not obviously or directly represented",
    "✓ Example: Total Abstractions- non-representational or non-objective because they stray away from the reality. Not all abstracts are non-representational."
]

for point in bullet_points:
    paragraph = text_frame.add_paragraph()
    paragraph.text = point
    paragraph.font.size = Pt(18)
    paragraph.font.color.rgb = RGBColor(0, 0, 0)
    paragraph.space_before = Pt(10)

# Save the presentation
presentation.save("render.pptx")