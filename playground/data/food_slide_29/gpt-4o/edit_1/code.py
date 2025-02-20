from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide layout
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 240, 255)  # Adjusted light background color

# Add title "PART 1"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "PART 1"
title.font.size = Pt(32)
title.font.bold = True

# Add chat icon
icon_path = "../../media/image_0.jpg"  # Use the provided image path
slide.shapes.add_picture(icon_path, Inches(0.5), Inches(1.5), Inches(2), Inches(2))

# Add instruction text
instruction_box = slide.shapes.add_textbox(Inches(3), Inches(1.5), Inches(10), Inches(1))
instruction_frame = instruction_box.text_frame
instruction = instruction_frame.add_paragraph()
instruction.text = "Answer the questions below about the video."
instruction.font.size = Pt(16)
instruction.font.bold = True

# Add questions
questions = [
    "1. What happened on Joey’s first date with Sarah?",
    "2. How did Joey’s friends, Phoebe and Rachel, react or respond to his date story?",
    "3. How did Joey feel about his first date with Sarah?",
    "4. What happened on Joey’s second date with Sarah?"
]

top_position = 3
for question in questions:
    question_box = slide.shapes.add_textbox(Inches(3), Inches(top_position), Inches(12), Inches(1))
    question_frame = question_box.text_frame
    question_text = question_frame.add_paragraph()
    question_text.text = question
    question_text.font.size = Pt(14)
    
    # Add a border around the text box
    sp = question_box
    sp_element = sp._element
    sp_element.get_or_add_ln().append(parse_xml(
        '<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:srgbClr val="B0B0B0"/>'
        '</a:solidFill>'
    ))

    top_position += 1.2  # Move down for the next question

# Add button "QUIZ & REVIEW ACTIVITY"
button_box = slide.shapes.add_textbox(Inches(12.5), Inches(7.5), Inches(3), Inches(1))
button_frame = button_box.text_frame
button_paragraph = button_frame.add_paragraph()
button_paragraph.text = "QUIZ & REVIEW ACTIVITY"
button_paragraph.font.size = Pt(14)
button_paragraph.font.bold = True
button_box.fill.solid()
button_box.fill.fore_color.rgb = RGBColor(100, 149, 237)  # Cornflower Blue

presentation.save("render.pptx")