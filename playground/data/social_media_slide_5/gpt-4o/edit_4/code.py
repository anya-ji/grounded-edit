from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add title box with dark blue background
title_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(1.5))
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue color
title_box.line.color.rgb = RGBColor(0, 51, 102)  # Match line color to title background

title_tf = title_box.text_frame
title_tf.text = "5 Stages of Development of Media"
title_tf.paragraphs[0].font.size = Pt(44)
title_tf.paragraphs[0].font.bold = True
title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text

# Add subtitle below title with increased spacing
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(15), Inches(1))  # Changed Y position to Inches(2)
subtitle_tf = subtitle_box.text_frame
subtitle_p = subtitle_tf.add_paragraph()
subtitle_p.text = "2. SCRIPT"
subtitle_p.font.size = Pt(32)
subtitle_p.font.bold = True
subtitle_p.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add bullet points with increased spacing
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(10), Inches(4))  # Changed Y position to Inches(3.2)
content_tf = content_box.text_frame
content_tf.word_wrap = True

bullet_points = [
    "It allowed humans to communicate over a larger space and for a much longer duration.",
    "It allowed the permanent codification of economic, cultural, and political practice."
]

for bullet in bullet_points:
    p = content_tf.add_paragraph()
    p.text = bullet
    p.font.size = Pt(24)
    p.space_after = Pt(12)

# Add image on the right side
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(11), Inches(2), Inches(5), Inches(5))

# Save the presentation
presentation.save("render.pptx")