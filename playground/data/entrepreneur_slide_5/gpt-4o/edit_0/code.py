from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Create a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add a white background for the main content area
background = slide.shapes.add_shape(
    autoshape_type_id=1, 
    left=Inches(0.5), 
    top=Inches(0.5), 
    width=Inches(15), 
    height=Inches(8)
)
background.fill.solid()
background.fill.fore_color.rgb = RGBColor(255, 255, 255)

# Add a vertical yellow bar on the left side
yellow_bar = slide.shapes.add_shape(
    autoshape_type_id=1, 
    left=Inches(0), 
    top=Inches(0), 
    width=Inches(0.5), 
    height=Inches(9)
)
yellow_bar.fill.solid()
yellow_bar.fill.fore_color.rgb = RGBColor(255, 223, 0)

# Add the title "Early life of Elon Musk"
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_text_frame = title_box.text_frame
title_text_frame.word_wrap = True

title_paragraph = title_text_frame.add_paragraph()
title_paragraph.text = "Early life of Elon Musk"
title_paragraph.font.size = Pt(44)
title_paragraph.font.bold = True
title_paragraph.alignment = PP_ALIGN.LEFT

# Add the bullet points on the right side with solid circles
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(7))
content_text = (
    "• Elon Reeve Musk was born in 1971 in Pretoria, South Africa.\n"
    "• He is the Errol Musk and Maye Musk's oldest son of three children.\n"
    "• Elon was constantly buried in Encyclopedia Britannica at the age of four.\n"
    "• At the age of ten, he developed an interest in computing with the Commodore VIC-20\n"
    "  and taught himself computer programming.\n"
    "• He used his skills to create a code of a BASIC-BASED video game called Blaster,\n"
    "  “a trivial game” and he sold it.\n"
    "• In 1992, he went to the University of Pennsylvania.\n"
    "• He got two degrees in Business and Physics Major."
)
content_text_frame = content_box.text_frame
content_text_frame.word_wrap = True
content_paragraph = content_text_frame.add_paragraph()
content_paragraph.text = content_text
content_paragraph.font.size = Pt(20)
content_paragraph.alignment = PP_ALIGN.LEFT

# Save the presentation
presentation.save("render.pptx")