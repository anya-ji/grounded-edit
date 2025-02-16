from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()

# Set slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add title and yellow horizontal line
title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title.text_frame
title_frame.clear()
p = title_frame.add_paragraph()
p.text = "Elon's Future Plans"
p.font.size = Pt(44)
p.font.bold = True

# Add a yellow horizontal line below the title
line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1), Inches(1.5), Inches(14), Inches(0.1)
)
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(255, 223, 0)  # Yellow color
line.line.fill.background()

# Add bullet points
left_content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
left_content_frame = left_content_box.text_frame
left_content_frame.word_wrap = True
left_content_frame.clear()

bullet_points = [
    "By 2027, Musk wants to launch 12,000 broadband satellites into orbit.",
    "HYPERLOOP trains start running in 2025.",
    "Propagation of Open-source technology.",
    "Mars colonization by 2030.",
]

for point in bullet_points:
    p = left_content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(24)
    p.space_after = Pt(10)

# Add Elon Musk image on the right
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(10), Inches(2), Inches(5), Inches(6))

# Add yellow accent on the right side
accent = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(15.5), Inches(0), Inches(0.5), Inches(9)
)
accent.fill.solid()
accent.fill.fore_color.rgb = RGBColor(255, 223, 0)  # Yellow color
accent.line.fill.background()

# Save the presentation
presentation.save("render.pptx")