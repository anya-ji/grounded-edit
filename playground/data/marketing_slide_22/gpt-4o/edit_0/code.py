from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation and slide
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)
slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set background color to a dark gradient
background = slide.background
fill = background.fill
fill.gradient()  # Use gradient method for a gradient background
fill.gradient_stops[0].color.rgb = RGBColor(0, 128, 128)  # dark teal
fill.gradient_stops[1].color.rgb = RGBColor(0, 64, 64)  # a darker shade of teal

# Add title on the right side
text_box = slide.shapes.add_textbox(Inches(10.5), Inches(1), Inches(5), Inches(1))
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = "After Sales Services"
p.font.size = Pt(44)
p.font.color.rgb = RGBColor(255, 255, 255)  # white font

# Insert images on the left side
image1_path = "../../media/image_0.jpg"
image1 = slide.shapes.add_picture(image1_path, Inches(0.5), Inches(0.5), width=Inches(5))

image2_path = "../../media/image_1.jpg"
image2 = slide.shapes.add_picture(image2_path, Inches(0.5), Inches(3.25), width=Inches(5))

image3_path = "../../media/image_2.jpg"
image3 = slide.shapes.add_picture(image3_path, Inches(0.5), Inches(6), width=Inches(5))

# Add captions below each image
caption_text = [
    "This Photo by Unknown author is licensed under CC BY-SA-NC",
    "This Photo by Unknown author is licensed under CC BY-SA-NC",
    "This Photo by Unknown author is licensed under CC BY"
]

for i, img in enumerate([image1, image2, image3]):
    left = img.left
    top = img.top + img.height
    width = img.width
    caption_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = caption_box.text_frame
    cap = tf.add_paragraph()
    cap.text = caption_text[i]
    cap.font.size = Pt(12)

# Save the presentation
presentation.save("render.pptx")