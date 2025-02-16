from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
presentation = Presentation()

# Define slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add the title "PART 2"
title_box = slide.shapes.add_textbox(Inches(6.5), Inches(0.5), Inches(3), Inches(0.8))  # Centered
title_frame = title_box.text_frame
title_p = title_frame.add_paragraph()
title_p.text = "PART 2"
title_p.font.bold = True
title_p.font.size = Pt(32)

# Add instructions below the title
instructions_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(15), Inches(1))
instructions_frame = instructions_box.text_frame
instructions_p = instructions_frame.add_paragraph()
instructions_p.text = ("Write the correct form of the verbs in the dialogue below, "
                       "depending on whether the past simple or past continuous should be used.")
instructions_p.font.size = Pt(18)
instructions_p.space_after = Pt(12)  # Adjust space after the paragraph
instructions_frame.word_wrap = True

# Add dialogue with blanks
dialogue_text = (
    "JOEY: We __________ (be) out to dinner. We __________ (have) a really nice time, "
    "I __________ (think) she was really cool and then, out of nowhere, she __________ (reach) over "
    "and __________ (take) some of my fries!\n\n"
    "PHOEBE: So she __________ (take) some fries, big deal!\n\n"
    "RACHEL: Oh yeah, Joey doesn’t share food. I mean, just last week, we __________ (have) breakfast, "
    "and he __________ (have) a couple of grapes on his plate..."
)

dialogue_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(3))  # Centered
dialogue_frame = dialogue_box.text_frame
dialogue_p = dialogue_frame.add_paragraph()
dialogue_p.text = dialogue_text
dialogue_p.font.size = Pt(16)
dialogue_frame.word_wrap = True
# Adjusting line spacing
for paragraph in dialogue_frame.paragraphs:
    paragraph.space_after = Pt(10)  # Set the space after each paragraph

# Add "PART 3"
part3_box = slide.shapes.add_textbox(Inches(6.5), Inches(5.5), Inches(3), Inches(0.8))  # Adjusted position
part3_frame = part3_box.text_frame
part3_p = part3_frame.add_paragraph()
part3_p.text = "PART 3"
part3_p.font.bold = True
part3_p.font.size = Pt(32)
part3_p.space_after = Pt(12)  # Adjust space after the paragraph

# Add instructions for PART 3
part3_instr_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(15), Inches(1))  # Adjusted position
part3_instr_frame = part3_instr_box.text_frame
part3_instr_p = part3_instr_frame.add_paragraph()
part3_instr_p.text = ("Retell Joey’s part of the dialogue (above) out loud using stress/intonation "
                      "in the appropriate places.")
part3_instr_p.font.size = Pt(18)
part3_instr_frame.word_wrap = True

# Save the presentation
presentation.save("render.pptx")