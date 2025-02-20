from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation object and set the dimensions
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background color to white
background = slide.background
background.fill.solid()
background.fill.fore_color.rgb = RGBColor(255, 255, 255)

# Add a dark blue rectangle for the title background
title_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), presentation.slide_width, Inches(1)
)
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue
title_shape.line.color.rgb = RGBColor(0, 51, 102)  # Same as fill to avoid lines

# Add title text
title = "3 Factors that have affected the process of Economic Globalization"
title_box = slide.shapes.add_textbox(Inches(0), Inches(0), presentation.slide_width, Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_p = title_frame.add_paragraph()
title_p.text = title
title_p.font.size = Pt(32)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(255, 255, 255)  # White

# Add main content text
content = (
    "3. The character and pace of economic integration have been significantly "
    "influenced by public policies, although it is not always in the direction "
    "of increasing economic integration."
)
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), presentation.slide_width - Inches(2), Inches(5))
content_frame = content_box.text_frame
content_frame.word_wrap = True
content_p = content_frame.add_paragraph()
content_p.text = content
content_p.font.size = Pt(24)
content_p.font.bold = False
content_p.font.name = 'Calibri'  # Sans-serif font
content_p.font.color.rgb = RGBColor(0, 0, 0)  # Black

# Insert bullet point
bullet_p = content_frame.add_paragraph()
bullet_p.text = "• "  # Bullet character
bullet_p.font.size = Pt(24)
bullet_p.font.bold = False
bullet_p.font.name = 'Calibri'  # Sans-serif font
bullet_p.font.color.rgb = RGBColor(0, 0, 0)  # Black
bullet_p.space_before = Pt(10)  # Add space before to separate from previous content

# Save the presentation
presentation.save("render.pptx")