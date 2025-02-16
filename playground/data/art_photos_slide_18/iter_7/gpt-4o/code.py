from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide with custom dimensions
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background to a solid bright color (e.g., bright yellow)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 0)  # Bright yellow background

# Create a thin yellow border
border_shape = slide.shapes.add_shape(
    1,  # 1 is the shape type for rectangle
    0,  # Left
    0,  # Top
    presentation.slide_width, 
    presentation.slide_height
)
border_shape.fill.solid()
border_shape.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow color
border_shape.line.color.rgb = RGBColor(255, 255, 0)  # Yellow border
border_shape.line.width = Pt(1)  # Thin line

# Add the main title
main_title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1.5))
main_title_frame = main_title_box.text_frame
main_title_frame.text = "THE CONTENT OF ART"
main_title_frame.paragraphs[0].font.size = Pt(44)
main_title_frame.paragraphs[0].font.bold = True
main_title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
main_title_box.fill.solid()
main_title_box.fill.fore_color.rgb = RGBColor(211, 211, 211)  # Light gray background

# Add the subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(10), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "The Subject and Content of Art"
subtitle_frame.paragraphs[0].font.size = Pt(20)
subtitle_frame.paragraphs[0].font.bold = True
subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
subtitle_box.fill.solid()
subtitle_box.fill.fore_color.rgb = RGBColor(211, 211, 211)  # Light gray background

# Add the content
content_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(5))
content_frame = content_box.text_frame
content_frame.text = ("✓ It is the mass of ideas associated with each artwork "
                      "and communicated through the following:")
content_frame.paragraphs[0].font.size = Pt(20)
content_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

bullet_points = [
    "✓ 1. The art’s imagery",
    "✓ 2. The symbolic meaning",
    "✓ 3. Its surroundings where it is used or displayed",
    "✓ 4. The customs, beliefs and values of the culture that uses it",
]

# Black checkmarks for bullet points
for point in bullet_points:
    p = content_frame.add_paragraph()
    p.text = point.replace('✓', '', 1).strip()  # Remove the existing checkmark
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_before = Pt(10)  # Space before each bullet point

# Add a black checkmark in front of each listed item
for idx, point in enumerate(bullet_points):
    p = content_frame.paragraphs[idx + 1]  # Access the newly added paragraphs
    p.text = '✓ ' + p.text  # Add checkmark at the front
    p.font.color.rgb = RGBColor(0, 0, 0)  # Ensure checkmarks are black

content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

# Save the presentation
presentation.save("render.pptx")