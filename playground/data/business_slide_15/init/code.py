from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5] # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add a bright orange rectangle for contrast behind the title
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(1.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(255, 165, 0)  # Bright orange
shape.line.fill.background()

# Add title text
title_text = "Innovations (Future Products or Services)"
title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(16), Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = title_text
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.size = Pt(40)
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add content text
content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(6))
content_frame = content_box.text_frame

content_text = [
    "Airbnb 2021",
    "Release: 100+ innovations and upgrades across our entire service",
    "“We are seeing three fundamental shifts in travel as people become less tethered and more flexible.” — Brian Chesky, Co-Founder and CEO of Airbnb"
]

for line in content_text:
    p = content_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black

# Save the presentation
presentation.save("render.pptx")