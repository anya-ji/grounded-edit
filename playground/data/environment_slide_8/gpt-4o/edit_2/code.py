from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

presentation = Presentation()

# Set up slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Create slide with blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add background color with hexagonal pattern
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(144, 238, 144)  # Light green color

# Add title section
title_textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1.5))
title_tf = title_textbox.text_frame
title_tf.clear()
title_p = title_tf.add_paragraph()
title_p.text = "POSITIVE IMPACT"
title_p.font.bold = True
title_p.font.size = Pt(44)
title_p.font.color.rgb = RGBColor(0, 128, 0)  # Green color
title_tf.word_wrap = True

# Add content section
content_textbox = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(5.5))
content_tf = content_textbox.text_frame
content_tf.clear()

content = [
    "Promoting afforestation plays an important role in reducing air pollution by improving overall air quality.",
    "Conserving water resources- By not throwing industrial wastes, domestic wastes into rivers, lakes etc., water pollution can be prevented. Techniques like rainwater harvesting, watershed management, and drip irrigation help in conservation of water resources.",
    "Soil conservation- methods like terrace farming, crop rotation, windbreaks play an important role in preventing soil erosion.",
    "Conservation of flora and fauna by setting up conservation areas, botanical gardens, national parks, and wildlife sanctuaries.",
    "Use of renewable sources of energy like solar energy, biofuels, wind energy etc."
]

for line in content:
    p = content_tf.add_paragraph()
    p.text = line
    p.font.size = Pt(22)  # Changed font size to 22pt
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black color

content_tf.word_wrap = True

# Add image section
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(1), Inches(8), Inches(5), Inches(2))

# Save the presentation
presentation.save("render.pptx")