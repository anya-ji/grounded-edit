from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Choose a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set a light green background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(204, 255, 204)  # Light Green

# Add a large white textbox in the middle
left = Inches(1.5)
top = Inches(1)
width = Inches(13)
height = Inches(7)
textbox = slide.shapes.add_textbox(left, top, width, height)
textbox_fill = textbox.fill
textbox_fill.solid()
textbox_fill.fore_color.rgb = RGBColor(255, 255, 255)

text_frame = textbox.text_frame
text_frame.word_wrap = True

# Adding padding around the text
text_frame.margin_top = Pt(20)
text_frame.margin_bottom = Pt(20)
text_frame.margin_left = Pt(20)
text_frame.margin_right = Pt(20)

# Add title
title_p = text_frame.add_paragraph()
title_p.text = "Further types of natural environment"
title_p.font.bold = True
title_p.font.size = Pt(44)
title_p.font.color.rgb = RGBColor(0, 102, 0)  # Bold Green
title_p.alignment = PP_ALIGN.LEFT

# Add content with bullet points
bullet_points = [
    "Lithosphere- It is defined as the rock and crust surface that covers the earth. An example of lithosphere is the rocky mountain range.",
    "Hydrosphere- all the water on the earth's surface such as lakes and seas is known as hydrosphere.",
    "Atmosphere- the mixture of gases that surrounds the earth is known as atmosphere.",
    "Biosphere- the part of earth's surface and atmosphere which supports life and where plants and animals live is known as biosphere."
]

for point in bullet_points:
    p = text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black color for content text
    p.space_after = Pt(10)
    p.alignment = PP_ALIGN.LEFT
    p.level = 0

# Save the presentation to the specified path
presentation.save("render.pptx")