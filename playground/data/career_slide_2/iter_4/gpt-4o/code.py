from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a Title and Content layout
slide_layout = presentation.slide_layouts[5]  # Choosing a blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set the background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(10, 10, 10)  # Dark background

# Add title text
title_text = "WHAT ARE 'TOP GRADES'?"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1.5))
tf = title_box.text_frame
p = tf.add_paragraph()
p.text = title_text
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)  # White text for contrast

# Add initial bulleted text
content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(6))
tf = content_box.text_frame

# Add main bullet points
main_points = [
    "Highest possible grades or marks that a student can achieve in a particular educational system or institution.",
    "The specific grading system and what constitutes top grades can vary from one country, educational level, or institution to another.",
    "However, in many cases, top grades are associated with excellence in academic performance and are often represented by letters (e.g., A+, A, A-, etc.) or numerical scores (e.g., 90-100%)."
]

for point in main_points:
    p = tf.add_paragraph()
    p.text = point
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)  # White text for contrast
    p.space_after = Pt(14)  # Spacing after the point

# Add the new bullet points back to the text without the removed phrases
grading_system_points = [
    "Can vary from country to country.",
    "Generally associated with academic excellence and achievement."
]

for point in grading_system_points:
    p = tf.add_paragraph()
    p.text = point
    p.level = 1  # This makes it a sub-bullet
    p.font.size = Pt(20)  # Slightly smaller for sub-bullets
    p.font.color.rgb = RGBColor(255, 255, 255)  # White text for contrast
    p.space_after = Pt(10)  # Spacing between sub-points

# Save the presentation
presentation.save("render.pptx")