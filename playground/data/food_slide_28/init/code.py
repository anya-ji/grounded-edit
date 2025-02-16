from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background color to dark blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(10, 51, 102)  # dark blue RGB

# Add title text
title_box = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(12), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True

p = title_frame.add_paragraph()
p.text = "QUIZ & REVIEW ACTIVITY"
p.font.size = Pt(56)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)  # white RGB
p.alignment = PP_ALIGN.CENTER

# Add instruction text
instruction_box = slide.shapes.add_textbox(Inches(4.5), Inches(3.5), Inches(7), Inches(1))
instruction_frame = instruction_box.text_frame
instruction_frame.word_wrap = True

p = instruction_frame.add_paragraph()
p.text = "Work with a partner."
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(255, 255, 255)  # white RGB
p.alignment = PP_ALIGN.CENTER

# Add relevant icon (magnifying glass) image
icon_path = "icon_magnifying_glass.png"
slide.shapes.add_picture(icon_path, Inches(0.5), Inches(1.3), height=Inches(1.5))

# Add "FLUENT" logo to the top right corner
logo_path = "logo_fluent.png"
slide.shapes.add_picture(logo_path, Inches(12.5), Inches(0.5), height=Inches(1))

# Save the presentation
presentation.save("render.pptx")