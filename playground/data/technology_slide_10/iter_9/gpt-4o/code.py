from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Slide layout
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Title
title_text = "Advantages of Blockchain"
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title_frame.text = title_text
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(44)
title_paragraph.font.bold = True

# Checklist
checklist_items = [
    "Greater Transparency",
    "Highly secure",
    "Easily traceable",
    "High efficiency and speed",
    "Low cost",
    "Zero percentage of fraud",
    "Extremely volatile"
]

left = Inches(1)
top = Inches(1.5)
width = Inches(5.5)  # Adjusted width for better visibility
height = Inches(6)   # Adjusted height for all items to be visible

checklist_box = slide.shapes.add_textbox(left, top, width, height)
checklist_frame = checklist_box.text_frame
checklist_frame.word_wrap = True

# Adjust line spacing and alignment for uniformity
for item in checklist_items:
    p = checklist_frame.add_paragraph()
    p.text = f"☐ {item}"  # Use checkbox character
    p.font.size = Pt(28)  # Changed font size for better legibility
    p.font.bold = True  # Change to bold font
    p.space_after = Pt(10)  # Add space after each item for uniformity
    p.alignment = 0  # Left alignment (0 = LEFT)

# Image (Benefits of Blockchain Technology diagram)
image_path = "../../media/image_0.jpg"
# Center the image by adjusting the left position
image_left = (presentation.slide_width - Inches(8)) / 2  # Center horizontally
slide.shapes.add_picture(image_path, image_left, Inches(2), Inches(6), Inches(8))  # Increased height and width

# Save presentation
presentation.save("render.pptx")