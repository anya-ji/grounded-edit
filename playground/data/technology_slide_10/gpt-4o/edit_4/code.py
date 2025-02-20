from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Slide layout
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Title
title_text = "Advantages of Blockchain"
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title_frame.text = title_text
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(44)
title_paragraph.font.bold = True

# Checklist
checklist_items = [
    "Greater Transparency",
    "Highly secure",
    "Easily traceable",
    "High efficiency and speed",
    "Low cost",
    "Zero percentage of fraud",
    "Extremely volatile"
]

left = Inches(1)
top = Inches(1.5)
width = Inches(5)
height = Inches(5)

checklist_box = slide.shapes.add_textbox(left, top, width, height)
checklist_frame = checklist_box.text_frame
checklist_frame.word_wrap = True
for item in checklist_items:
    p = checklist_frame.add_paragraph()
    p.text = f"• {item}"
    p.font.size = Pt(24)

# Image (Benefits of Blockchain Technology diagram)
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(9), Inches(2), Inches(6), Inches(4))

# Caption for the image
caption_text = "Benefits of Blockchain Technology"
caption_box = slide.shapes.add_textbox(Inches(9), Inches(6), Inches(6), Inches(0.5))
caption_frame = caption_box.text_frame
caption_frame.text = caption_text
caption_paragraph = caption_frame.paragraphs[0]
caption_paragraph.font.size = Pt(18)
caption_paragraph.font.bold = True
caption_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black color

# Save presentation
presentation.save("render.pptx")