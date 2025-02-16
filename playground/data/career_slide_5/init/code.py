from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object and set slide size
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color to dark
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# Add title
title = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(16), Inches(1))
title.text_frame.text = "KEY REASONS WHY ACADEMIC SUCCESS IS IMPORTANT IN SOCIETY"
title.text_frame.paragraphs[0].font.size = Pt(34)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
title.text_frame.paragraphs[0].alignment = 1  # Center alignment

# Define layouts
layout_width = Inches(5)
layout_height = Inches(7.5)
image_height = Inches(3)
text_box_height = Inches(4.5)

# Section details
sections = [
    # Left Section
    {
        'image_path': "../../media/image_0.jpg",
        'title': "Innovation and Technological Advancement",
        'content': "Academic success fuels innovation, driving progress in various fields, including healthcare, communication, and transportation.",
        'left': Inches(0.5)
    },
    # Middle Section
    {
        'image_path': "../../media/image_1.jpg",
        'title': "Social Mobility and Equality",
        'content': "Education can promote social mobility and reduce disparities by providing opportunities for academic success.",
        'left': Inches(5.5)
    },
    # Right Section
    {
        'image_path': "../../media/image_2.jpg",
        'title': "Civic Engagement and Informed Citizenship",
        'content': "Education promotes critical thinking, enabling individuals to engage in democracy, contribute to community development, and solve social issues.",
        'left': Inches(10.5)
    },
]

# Add content for each section
for section in sections:
    # Add image
    slide.shapes.add_picture(section['image_path'], section['left'], Inches(1.5), width=layout_width, height=image_height)

    # Add title
    title_box = slide.shapes.add_textbox(section['left'], Inches(4.2), layout_width, Inches(1))
    tf_title = title_box.text_frame
    p = tf_title.add_paragraph()
    p.text = section['title']
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 0)

    # Add content
    text_box = slide.shapes.add_textbox(section['left'], Inches(5), layout_width, text_box_height)
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = section['content']
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)

# Save the presentation
presentation.save("render.pptx")