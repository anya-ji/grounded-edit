from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image  # Import Pillow for image processing

def enhance_image(image_path):
    img = Image.open(image_path)
    img = img.convert('RGB')
    data = img.getdata()
    enhanced_data = []
    for item in data:
        enhanced_data.append((
            min(int(item[0] * 1.2), 255),
            min(int(item[1] * 1.2), 255),
            min(int(item[2] * 1.2), 255)
        ))
    img.putdata(enhanced_data)
    enhanced_image_path = "enhanced_" + image_path.split("/")[-1]
    img.save(enhanced_image_path)
    return enhanced_image_path

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 240, 240)

title = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(16), Inches(1))
title.text_frame.text = "KEY REASONS WHY ACADEMIC SUCCESS IS IMPORTANT IN SOCIETY"
title.text_frame.paragraphs[0].font.size = Pt(34)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
title.text_frame.paragraphs[0].alignment = 1

layout_width = Inches(5)
layout_height = Inches(7.5)
image_height = Inches(3)
text_box_height = Inches(4.5)

sections = [
    {'image_path': "../../media/image_0.jpg", 'title': "Innovation and Technological Advancement",
     'content': "Academic success fuels innovation, driving progress in various fields, including healthcare, communication, and transportation.",
     'left': Inches(4.5)},
    {'image_path': "../../media/image_1.jpg", 'title': "Social Mobility and Equality",
     'content': "Education can promote social mobility and reduce disparities by providing opportunities for academic success.",
     'left': Inches(4.5), 'top_spacing': Inches(5)},
    {'image_path': "../../media/image_2.jpg", 'title': "Civic Engagement and Informed Citizenship",
     'content': "Education promotes critical thinking, enabling individuals to engage in democracy, contribute to community development, and solve social issues.",
     'left': Inches(4.5), 'top_spacing': Inches(10)},
]

for section in sections:
    enhanced_image_path = enhance_image(section['image_path'])
    slide.shapes.add_picture(enhanced_image_path, section['left'], Inches(1.5 + section.get('top_spacing', 0)), width=layout_width, height=image_height)

    title_box = slide.shapes.add_textbox(section['left'], Inches(1.5 + section.get('top_spacing', 0)), layout_width, Inches(1))
    tf_title = title_box.text_frame
    p = tf_title.add_paragraph()
    p.text = section['title']
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    text_box = slide.shapes.add_textbox(section['left'], Inches(2.5 + section.get('top_spacing', 0)), layout_width, text_box_height)
    tf = text_box.text_frame
    p = tf.add_paragraph()
    p.text = section['content']
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0, 0, 0)

# Save the presentation
presentation.save("render.pptx")