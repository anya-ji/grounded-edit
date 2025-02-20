from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
presentation = Presentation()

# Set slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide
slide_layout = presentation.slide_layouts[5]  # Blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Add title text box with dark blue background
title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(16), Inches(1.5))
title_text_frame = title_box.text_frame
title_text_frame.clear()

p = title_text_frame.add_paragraph()
p.text = "Various Media drive Various forms of Global Integration"
p.font.size = Pt(32)  # Changed font size to 32pt
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(0, 0, 139)  # Dark blue background
p.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Add main content text box with white background
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(15), Inches(6.5))
content_text_frame = content_box.text_frame
content_text_frame.clear()

# Add "Various Media" section
p = content_text_frame.add_paragraph()
p.text = "Various Media"
p.font.size = Pt(28)
p.font.bold = True
p.space_after = Pt(10)

bullet_points = [
    "are used for globalization to work all over the world",
    "Media plays a major role in globalization"
]

for point in bullet_points:
    bullet = content_text_frame.add_paragraph()
    bullet.text = point
    bullet.level = 1
    bullet.font.size = Pt(20)

# Add "Examples:" subheading and bullet points
p = content_text_frame.add_paragraph()
p.text = "Examples:"
p.font.size = Pt(28)
p.font.bold = True
p.space_after = Pt(10)

examples_points = [
    "Television, Internet, Computers etc.",
    "Considered to have a significant influence in globalization"
]

for example in examples_points:
    bullet = content_text_frame.add_paragraph()
    bullet.text = example
    bullet.level = 1
    bullet.font.size = Pt(20)

# Set the background of the content box to white
content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

# Save the presentation
presentation.save("render.pptx")