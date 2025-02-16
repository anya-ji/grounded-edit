from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Insert the warehouse image
warehouse_img_path = "../../media/image_1.jpg"
slide.shapes.add_picture(warehouse_img_path, Inches(0), Inches(0), Inches(8), Inches(9))

# Insert the ship image
ship_img_path = "../../media/image_0.jpg"
left = Inches(8)
top = Inches(0)
width = Inches(8)
height = Inches(9)
slide.shapes.add_picture(ship_img_path, left, top, width, height)

# Add semi-transparent overlay to enhance text visibility
overlay = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, left, top, width, height
)
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
overlay.fill.transparency = 0.3
overlay.line.fill.background()

# Remove the curved shape to allow a full view of the forklift
# Remove the following lines:
# left_part_width = Inches(8)
# curved_shape = slide.shapes.add_shape(
#     MSO_SHAPE.WAVE, left_part_width - Inches(1), Inches(0), Inches(2), Inches(9)
# )
# curved_shape.fill.solid()
# curved_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
# curved_shape.line.fill.background()

# Add the title text
title_text_box = slide.shapes.add_textbox(left, Inches(3.5), Inches(7), Inches(2))
text_frame = title_text_box.text_frame
p = text_frame.add_paragraph()
p.text = "Transportation and Storage"
p.font.bold = True
p.font.size = Pt(44)
p.font.color.rgb = RGBColor(0, 0, 0)  # Black color for contrast

# Save the presentation
presentation.save("render.pptx")