from pptx import Presentation
from pptx.util import Inches, Pt

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add title
title_text = "Disadvantages of Blockchain"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_frame.text = title_text
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.bold = True
title_paragraph.font.size = Pt(40)

# Add bullet points
points = [
    "Blockchain programming is complex and challenging",
    "It is international so it will have to face a variety of global laws",
    "Blockchain relies on a network effect",
    "More power consumption",
    "Immutability"
]

left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7), Inches(5))
left_frame = left_box.text_frame
left_frame.word_wrap = True

for point in points:
    p = left_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(24)
    p.space_after = Pt(10)  # Adjusted spacing for uniformity
    p.space_before = Pt(0)   # Ensure no space before for uniformity

# Add image on the right
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(9), Inches(1.5), Inches(6), Inches(5))

# Save presentation
presentation.save("render.pptx")