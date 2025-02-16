from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide
slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Title setup
title = "ACADEMIC WORLD AND PROFESSIONAL WORLD: DEMANDS AND CHARACTERISTICS"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_p = title_frame.add_paragraph()
title_p.text = title
title_p.font.size = Pt(28)
title_p.font.bold = True

# Left Column: ACADEMIC WORLD
academic_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7), Inches(0.5))
academic_title_frame = academic_title_box.text_frame
academic_title_p = academic_title_frame.add_paragraph()
academic_title_p.text = "ACADEMIC WORLD"
academic_title_p.font.size = Pt(24)
academic_title_p.font.bold = True
academic_title_p.font.color.rgb = RGBColor(0, 0, 255)  # Blue

academic_content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.2), Inches(7), Inches(5))
academic_content_frame = academic_content_box.text_frame

# Hierarchy
hierarchy = academic_content_frame.add_paragraph()
hierarchy.text = "Hierarchy"
hierarchy.font.size = Pt(18)
hierarchy.font.bold = True
hierarchy.space_after = Pt(12)

instructors = academic_content_frame.add_paragraph()
instructors.text = "Instructors"
instructors.font.size = Pt(16)

bullet1 = academic_content_frame.add_paragraph()
bullet1.text = "• Professors and teachers are the primary authority figures."
bullet1.space_before = Pt(0)
bullet1.level = 1
bullet1.font.size = Pt(14)

# Peer Interaction
peer_interaction = academic_content_frame.add_paragraph()
peer_interaction.text = "Peer Interaction"
peer_interaction.font.size = Pt(18)
peer_interaction.font.bold = True
peer_interaction.space_after = Pt(12)

bullet2 = academic_content_frame.add_paragraph()
bullet2.text = "• Students interact mainly with their peers in a relatively closed academic community."
bullet2.space_before = Pt(0)
bullet2.font.size = Pt(14)

# Learning Approach
learning_approach = academic_content_frame.add_paragraph()
learning_approach.text = "Learning Approach"
learning_approach.font.size = Pt(18)
learning_approach.font.bold = True  # Ensure it is bold
learning_approach.space_after = Pt(12)

theoretical = academic_content_frame.add_paragraph()
theoretical.text = "Theoretical"
theoretical.font.size = Pt(16)

bullet3 = academic_content_frame.add_paragraph()
bullet3.text = "• The academic world emphasizes theoretical knowledge and academic theories."
bullet3.space_before = Pt(0)
bullet3.level = 1
bullet3.font.size = Pt(14)

research = academic_content_frame.add_paragraph()
research.text = "Emphasis on Research"
research.font.size = Pt(16)

bullet4 = academic_content_frame.add_paragraph()
bullet4.text = "• Research is a fundamental aspect, particularly at the university level."
bullet4.space_before = Pt(0)
bullet4.level = 1
bullet4.font.size = Pt(14)

# Right Column: PROFESSIONAL WORLD
professional_title_box = slide.shapes.add_textbox(Inches(8.5), Inches(1.5), Inches(7), Inches(0.5))
professional_title_frame = professional_title_box.text_frame
professional_title_p = professional_title_frame.add_paragraph()
professional_title_p.text = "PROFESSIONAL WORLD"
professional_title_p.font.size = Pt(24)
professional_title_p.font.bold = True
professional_title_p.font.color.rgb = RGBColor(0, 255, 0)  # Changed to green

professional_content_box = slide.shapes.add_textbox(Inches(8.5), Inches(2.2), Inches(7), Inches(5))
professional_content_frame = professional_content_box.text_frame

# Hierarchy
prof_hierarchy = professional_content_frame.add_paragraph()
prof_hierarchy.text = "Hierarchy"
prof_hierarchy.font.size = Pt(18)
prof_hierarchy.font.bold = True
prof_hierarchy.space_after = Pt(12)

supervisors = professional_content_frame.add_paragraph()
supervisors.text = "Supervisors and Managers"
supervisors.font.size = Pt(16)

bullet5 = professional_content_frame.add_paragraph()
bullet5.text = "• In the professional world, supervisors, managers, and clients are the authority figures."
bullet5.space_before = Pt(0)
bullet5.level = 1
bullet5.font.size = Pt(14)

# Collaboration
collaboration = professional_content_frame.add_paragraph()
collaboration.text = "Collaboration"
collaboration.font.size = Pt(18)
collaboration.font.bold = True
collaboration.space_after = Pt(12)

bullet6 = professional_content_frame.add_paragraph()
bullet6.text = "• Collaboration with colleagues, clients, and stakeholders is common."
bullet6.space_before = Pt(0)
bullet6.font.size = Pt(14)

# Learning Approach
prof_learning_approach = professional_content_frame.add_paragraph()
prof_learning_approach.text = "Learning Approach"
prof_learning_approach.font.size = Pt(18)
prof_learning_approach.font.bold = True
prof_learning_approach.space_after = Pt(12)
prof_learning_approach.font.color.rgb = RGBColor(0, 255, 0)  # Updated to green

practical = professional_content_frame.add_paragraph()
practical.text = "Practical"
practical.font.size = Pt(16)

bullet7 = professional_content_frame.add_paragraph()
bullet7.text = "• Professionals apply theoretical knowledge acquired in academia to real-world problems and situations."
bullet7.space_before = Pt(0)
bullet7.level = 1
bullet7.font.size = Pt(14)

continuous_learning = professional_content_frame.add_paragraph()
continuous_learning.text = "Continuous Learning"
continuous_learning.font.size = Pt(16)

bullet8 = professional_content_frame.add_paragraph()
bullet8.text = "• Lifelong learning and professional development are essential to stay competitive."
bullet8.space_before = Pt(0)
bullet8.level = 1
bullet8.font.size = Pt(14)

# Save the presentation
presentation.save("render.pptx")