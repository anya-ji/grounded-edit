from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_slide():
    # Create presentation object
    presentation = Presentation()
    presentation.slide_width = Inches(16)
    presentation.slide_height = Inches(9)

    # Choose a blank slide layout
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)

    # Set slide background color to light grey for contrast
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light grey color

    # Add slide title
    title_text = "Purpose of Blockchain"
    title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True

    title_p = title_frame.add_paragraph()
    title_p.text = title_text
    title_p.font.bold = True
    title_p.font.size = Pt(44)
    title_p.alignment = PP_ALIGN.CENTER

    # Text for bullet points
    bullet_points = [
        "Blockchain uses an online ledger which is very secure.",
        "There is no third-party interference.",
        "It used distributed ledger which is very transparent.",
        "It is a digital world offering many new tools and there are many centralized administrators.",
        "Since there is no third-party guarantee cost is very low."
    ]

    # Add bullet points to the slide
    left_margin = Inches(1.5)
    top_margin = Inches(2)

    bullet_box = slide.shapes.add_textbox(left_margin, top_margin, Inches(13), Inches(6))
    bullet_frame = bullet_box.text_frame
    bullet_frame.word_wrap = True

    # Adding bullet points with checkboxes
    for point in bullet_points:
        p = bullet_frame.add_paragraph()
        p.text = point
        p.font.size = Pt(24)
        p.font.color.rgb = RGBColor(0, 0, 0)  # setting text color to black
        
        # Add a bullet (checkbox styled) prefix
        p.text = "\u2611 " + point  # ☑ is unicode character for checked box

    # Save the presentation
    presentation.save("render.pptx")

# Run the function to create the slide
create_slide()