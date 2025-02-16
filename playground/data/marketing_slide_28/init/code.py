from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Initialize presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Create slide layout
slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Define colors
red_color = RGBColor(255, 0, 0)
white_color = RGBColor(255, 255, 255)
blue_gradient = RGBColor(0, 128, 128)  # Alternative: use a real gradient if your library version supports it

# Set slide background color to blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = blue_gradient

# Add double-arrow shape with title in the center
arrow_shape = slide.shapes.add_shape(
    MSO_SHAPE.CHEVRON, Inches(4.5), Inches(0.5), Inches(7), Inches(1)
)
arrow_shape.fill.solid()
arrow_shape.fill.fore_color.rgb = red_color
arrow_shape_text_frame = arrow_shape.text_frame
arrow_shape_text_frame.text = "Types of Market"
p = arrow_shape_text_frame.paragraphs[0]
p.font.bold = True
p.font.size = Pt(44)
p.font.color.rgb = white_color
p.alignment = PP_ALIGN.CENTER

# Add two side arrows
left_arrow = slide.shapes.add_shape(
    MSO_SHAPE.RIGHT_ARROW, Inches(1), Inches(0.8), Inches(3), Inches(0.6)
)
left_arrow.flip_horizontal()
left_arrow.fill.solid()
left_arrow.fill.fore_color.rgb = red_color

right_arrow = slide.shapes.add_shape(
    MSO_SHAPE.RIGHT_ARROW, Inches(12), Inches(0.8), Inches(3), Inches(0.6)
)
right_arrow.fill.solid()
right_arrow.fill.fore_color.rgb = red_color

# Add text boxes for markets
market_texts = [
    "Weekly market",
    "Neighborhood market",
    "Shopping Complexes",
    "Malls",
]

market_boxes = []
for idx, text in enumerate(market_texts):
    box = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(2.5 + (idx % 2) * 7), Inches(2 + (idx // 2) * 3), Inches(6), Inches(1)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = red_color
    text_frame = box.text_frame
    text_frame.text = text
    p = text_frame.paragraphs[0]
    p.font.bold = True
    p.font.size = Pt(32)
    p.font.color.rgb = white_color
    p.alignment = PP_ALIGN.CENTER
    market_boxes.append(box)

# Add arrows connecting title to text boxes
for idx, box in enumerate(market_boxes):
    start_x = Inches(8)
    end_x = Inches(5.5 + (idx % 2) * 8.5)
    arrow = slide.shapes.add_connector(
        MSO_SHAPE.LINE_ARROW, start_x, Inches(1.5), end_x, Inches(2.5 + (idx // 2) * 3)
    )
    arrow.line.color.rgb = red_color

# Save presentation
presentation.save("render.pptx")