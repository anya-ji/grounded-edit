from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5] 
slide = presentation.slides.add_slide(slide_layout)

# Set the background color to white for the slide
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Add a textbox for the title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(10), Inches(1))
title = title_box.text_frame.add_paragraph()
title.text = "How he came up with his ideas?"
title.font.size = Pt(44)
title.font.bold = True

# Add the main text content
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(10), Inches(4))
content = content_box.text_frame.add_paragraph()
content.text = (
    "When he was in college, he considered what he wanted to accomplish with his life "
    "and came up with a list of five ideas."
)
content.font.size = Pt(28)
content.font.bold = False

# List of five ideas as bullet points
bullet_points = [
    "The internet",
    "Sustainable energy",
    "Space exploration",
    "Life beyond earth",
    "Artificial intelligence",
    "Reprogramming the human genetic code"
]

for point in bullet_points:
    bullet = content_box.text_frame.add_paragraph()
    bullet.text = point
    bullet.font.size = Pt(24)
    bullet.space_after = Pt(10)

# Add a yellow accent on the right side
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(10.5), Inches(0), Inches(5.5), Inches(9)
)
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 223, 34)  # Light yellow color

# Place the image on the right side of the slide
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(11), Inches(1), Inches(4), Inches(6))

# Save the presentation
presentation.save("render.pptx")