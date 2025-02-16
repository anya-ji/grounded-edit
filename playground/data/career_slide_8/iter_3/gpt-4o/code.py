from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(14), Inches(1.5))
title = title_box.text_frame.add_paragraph()
title.text = "SUCCESSFUL INDIVIDUALS WITH TOP ACADEMIC RECORDS"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(255, 255, 255)  # Changed title text color to white for contrast

# Background color (dark gray)
fill = slide.background.fill
fill.solid()  # Changed to solid fill
fill.fore_color.rgb = RGBColor(50, 50, 50)  # Dark gray

# Function to add individual section
def add_individual_section(left, image_path, name, title, education_list):
    # Add circular picture
    picture = slide.shapes.add_picture(image_path, left, Inches(2), width=Inches(3))
    picture.left = left + (picture.width - picture.height) // 2
    picture.height = picture.width

    # Create a circle shape to mask the picture
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, picture.left, picture.top, picture.height, picture.height)
    slide.shapes._spTree.remove(circle._element)
    slide.shapes._spTree.insert(2, circle._element)

    # Name
    text_box = slide.shapes.add_textbox(left, Inches(5), Inches(3), Inches(1))
    text_frame = text_box.text_frame
    name_p = text_frame.add_paragraph()
    name_p.text = name
    name_p.font.bold = True
    name_p.font.size = Pt(18)
    name_p.font.color.rgb = RGBColor(255, 255, 255)  # Name text color changed to white for contrast

    # Title
    title_p = text_frame.add_paragraph()
    title_p.text = title
    title_p.font.size = Pt(14)
    title_p.font.color.rgb = RGBColor(255, 255, 255)  # Title text color changed to white for contrast
    
    # Education
    edu_p = text_frame.add_paragraph()
    edu_p.text = "Education"
    edu_p.font.size = Pt(14)
    edu_p.font.bold = True
    edu_p.font.color.rgb = RGBColor(255, 255, 255)  # Education header color changed to white for contrast

    for edu in education_list:
        edu_bullet = text_frame.add_paragraph()
        edu_bullet.text = edu
        edu_bullet.font.size = Pt(12)
        edu_bullet.font.color.rgb = RGBColor(255, 255, 255)  # Education bullet color changed to white for contrast
        edu_bullet.level = 1

# Paths to images
elon_image_path = "../../media/image_0.jpg"
warren_image_path = "../../media/image_2.jpg"
sundar_image_path = "../../media/image_3.jpg"
condoleezza_image_path = "../../media/image_1.jpg"

# Add sections with adjusted spacing
spacing = 4.0  # Adjust the spacing as necessary
add_individual_section(Inches(0.5), elon_image_path, "ELON MUSK", "CEO of SpaceX and Tesla", 
                       ["University of Pennsylvania", "Bachelor of Science in Physics", "Bachelor of Arts in Economics"])

add_individual_section(Inches(0.5 + spacing), warren_image_path, "WARREN BUFFETT", "Business magnate, investor, and philanthropist", 
                       ["University of Pennsylvania", "University of Nebraska-Lincoln (BS)", "Columbia University"])

add_individual_section(Inches(0.5 + 2 * spacing), sundar_image_path, "SUNDAR PICHAI", "CEO of Alphabet and Google", 
                       ["IIT Kharagpur", "Stanford University (MS)", "University of Pennsylvania (MBA)"])

add_individual_section(Inches(0.5 + 3 * spacing), condoleezza_image_path, "CONDOLEEZZA RICE", "Former U.S. Secretary of State", 
                       ["University of Denver (BA, PhD)", "University of Notre Dame (MA)"])

# Save presentation
presentation.save("render.pptx")