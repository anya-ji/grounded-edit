from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a new presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5] 
slide = presentation.slides.add_slide(slide_layout)

# Set the background color to light green
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(204, 255, 204) # light green

# Adding a textbox on the right side for title and content
left = Inches(10)
top = Inches(1)
width = Inches(5.5)
height = Inches(3)

textbox = slide.shapes.add_textbox(left, top, width, height)

# Add title
title_tf = textbox.text_frame
title = title_tf.add_paragraph()
title.text = "NATURAL ENVIRONMENT"
title.font.bold = True
title.font.size = Pt(44)
title.font.color.rgb = RGBColor(0, 102, 0)  # green color
title.space_after = Pt(20)

# Add content
content = "Name- Abhra Pawan Sharma\nClass- 9H\nRoll no- 4"
content_ph = title_tf.add_paragraph()
content_ph.text = content
content_ph.font.size = Pt(26)
content_ph.font.color.rgb = RGBColor(0, 0, 0)  # black color
content_ph.space_after = Pt(20)

# Save presentation
presentation.save("render.pptx")