from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add gradient background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 64, 0)
fill_color = slide.shapes.add_shape(
    8, Inches(0), Inches(0), presentation.slide_width, presentation.slide_height
)
fill_color.fill.solid()
fill_color.fill.fore_color.rgb = RGBColor(34, 139, 34)
fill_color.fill.gradient()
fill_color.fill.gradient_stops[0].position = 0.0
fill_color.fill.gradient_stops[1].position = 1.0

# Title text box
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "LOOK AROUND:"
title_frame.paragraphs[0].font.size = Pt(44)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.underline = True
# Change text color to white
title_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
title_frame.paragraphs[0].alignment = PP_ALIGN.LEFT

# Content text box
content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(4))
content_frame = content_box.text_frame
p = content_frame.add_paragraph()
p.text = "Consider yourself as a producer, list the various things you must be taking care till your product reaches the consumers."
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(255, 255, 255)
p.space_after = Pt(14)
content_frame.word_wrap = True

# Insert person/group icon
image_path = "../../media/image_0.jpg"
icon = slide.shapes.add_picture(image_path, Inches(12), Inches(2), Inches(3), Inches(3))

# Save the presentation
presentation.save("render.pptx")