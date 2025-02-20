from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]  # Blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Add the background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add title textbox
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(5), Inches(1))
title_frame = title_box.text_frame
title_text = title_frame.add_paragraph()
title_text.text = "The Subject and Content of Art"
title_text.font.size = Pt(24)
title_text.font.bold = True
title_text.font.color.rgb = RGBColor(0, 0, 0)

# Set title background color
fill = title_box.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 0)

# Add quote text box with white background
quote_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(4))
quote_frame = quote_box.text_frame
quote_frame.clear()  # Clear default placeholder content
quote_text = quote_frame.add_paragraph()
quote_text.text = (
    "In an artwork, the subject matter is what the image literally depicts. "
    "The content of the artwork is what the image means."
)
quote_text.font.size = Pt(32)
quote_text.font.color.rgb = RGBColor(0, 0, 0)  # Change text color to black

# Set white background for the quote text area
fill = quote_box.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

presentation.save("render.pptx")