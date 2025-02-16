from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]  # Blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Add the background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add a colorful abstract background
abstract_box = slide.shapes.add_shape(
    1,  # Shape type: 1 corresponds to a rectangle
    Inches(1), Inches(1), Inches(14), Inches(4)  # Position and size to cover the quote box
)
abstract_fill = abstract_box.fill
abstract_fill.solid()
abstract_fill.fore_color.rgb = RGBColor(0, 128, 255)  # Example color change

# Add a semi-transparent overlay with 70% opacity
overlay_box = slide.shapes.add_shape(
    1,  # Shape type: 1 corresponds to a rectangle
    Inches(0), Inches(0), Inches(16), Inches(9)  # Same size as slide
)
overlay_fill = overlay_box.fill
overlay_fill.solid()
overlay_fill.fore_color.rgb = RGBColor(255, 255, 255)  # White color
overlay_fill.alpha = 178  # 70% opacity (255 * 0.7 = 178.5, rounded down)

# Add a solid yellow header bar with a narrower width
header_box = slide.shapes.add_shape(
    1,  # Shape type: 1 corresponds to a rectangle
    Inches(0), Inches(0), Inches(12), Inches(1)  # Adjusted width to 12 inches
)
header_fill = header_box.fill
header_fill.solid()
header_fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow color

# Add title textbox
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(5), Inches(1))
title_frame = title_box.text_frame
title_text = title_frame.add_paragraph()
title_text.text = "The Subject and Content of Art"
title_text.font.size = Pt(24)
title_text.font.bold = True
title_text.font.color.rgb = RGBColor(255, 255, 0)  # Change text color to yellow

# Set title background color
fill = title_box.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 0)  # Also make background yellow to match

# Add quote text box with yellow background and center alignment
quote_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(4))  # Adjusted vertical position
quote_frame = quote_box.text_frame
quote_frame.clear()  # Clear default placeholder content
quote_text = quote_frame.add_paragraph()
quote_text.text = (
    "In an artwork, the subject matter is what the image literally depicts. "
    "The content of the artwork is what the image means."
)
quote_text.font.size = Pt(32)
quote_text.font.color.rgb = RGBColor(0, 0, 0)  # Change text color to black
quote_frame.text = ""  # Clear any previously set text to reset layout.
quote_frame.add_paragraph().text = quote_text.text  # Ensure paragraph is added to preserve formatting

# Set yellow background for the quote text area
fill = quote_box.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 0)  # Change background color to yellow

# Align text in the middle
for paragraph in quote_frame.paragraphs:
    paragraph.alignment = 1  # 1 corresponds to center alignment

presentation.save("render.pptx")