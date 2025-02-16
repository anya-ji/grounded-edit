from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide layout
slide_layout = presentation.slide_layouts[5]  # Use blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set colors
black_rgb = RGBColor(0, 0, 0)
white_rgb = RGBColor(255, 255, 255)
gray_rgb = RGBColor(128, 128, 128)  # Original gray color
light_gray_rgb = RGBColor(192, 192, 192)  # Lighter gray color for better contrast

# Add black rectangle on the left side
left_box = slide.shapes.add_shape(
    autoshape_type_id=1,  # Rectangle
    left=Inches(0),
    top=Inches(0),
    width=Inches(8),
    height=Inches(9)
)
left_box.fill.solid()
left_box.fill.fore_color.rgb = black_rgb

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(7), Inches(1))
title_frame = title_box.text_frame
title_frame.clear()  # Clear any existing text

title = title_frame.add_paragraph()
title.text = "ELO's"
# Removed bold setting for title
title.font.size = Pt(150)  # Increased font size to 150
title.font.color.rgb = light_gray_rgb  # Changed font color to lighter gray

# Add bullet points
bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(7), Inches(5))
bullet_frame = bullet_box.text_frame
bullet_frame.clear()

for bullet_text in ["Market and Opportunity", "Market and Equality"]:
    bullet = bullet_frame.add_paragraph()
    bullet.text = bullet_text
    bullet.font.size = Pt(32)
    bullet.font.bold = True  # Make bullet points bold
    bullet.font.color.rgb = white_rgb

# Insert image on the right side
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(8), Inches(0), Inches(8), Inches(9))

# Save presentation
presentation.save("render.pptx")