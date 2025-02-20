from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Define slide layout
slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add title textbox
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
p = title_frame.add_paragraph()
p.text = "Global and Local Cultural Products"
p.font.bold = True
p.font.size = Pt(44)
p.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Add subtitle
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(1))
subtitle_frame = subtitle_box.text_frame
p_subtitle = subtitle_frame.add_paragraph()
p_subtitle.text = "Global Product"
p_subtitle.font.size = Pt(32)
p_subtitle.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Change title background color
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark Blue

# Add main content with padding
content_box = slide.shapes.add_textbox(Inches(1 + (20 / 96)), Inches(3), Inches(14 - (40 / 96)), Inches(2))
content_frame = content_box.text_frame
p_content = content_frame.add_paragraph()
p_content.text = "• Those products that are marketed internationally under the same brand name, features, and specifications across countries."
p_content.font.size = Pt(28)
content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

# Add images
image_paths = [
    "../../media/image_0.jpg",  # Coca-Cola
    "../../media/image_1.jpg",  # McDonald's
    "../../media/image_2.jpg",  # Apple
    "../../media/image_3.jpg"   # Adidas
]

# Calculate spacing
image_width = Inches(2)
image_height = Inches(2)
space_between = Inches(0.5)
start_x = (presentation.slide_width - (4 * image_width + 3 * space_between)) / 2

# Add each image
for i, image_path in enumerate(image_paths):
    slide.shapes.add_picture(image_path, start_x + i * (image_width + space_between), Inches(6), image_width, image_height)

# Save the presentation
presentation.save("render.pptx")