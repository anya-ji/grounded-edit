from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)  # white background

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(6), Inches(1))
title = title_box.text_frame.add_paragraph()
title.text = "ELON MUSK"
title.font.size = Pt(44)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)  # black

# Add subtitle
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(6), Inches(1))
subtitle = subtitle_box.text_frame.add_paragraph()
subtitle.text = "Lahiru Herath"
subtitle.font.size = Pt(24)
subtitle.font.color.rgb = RGBColor(0, 0, 0)  # black

# Insert image with white border
image_path = "../../media/image_0.jpg"
pic = slide.shapes.add_picture(image_path, Inches(8), Inches(1), Inches(6), Inches(7))
line = pic.line
line.color.rgb = RGBColor(255, 255, 255)  # white border
line.width = Pt(4)

# Add yellow corner border
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(13.5), Inches(0), Inches(2.5), Inches(0.5)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(255, 215, 0)  # yellow
shape.line.fill.background()

# Save presentation
presentation.save("render.pptx")