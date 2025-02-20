from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation and a blank slide
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)
slide_layout = presentation.slide_layouts[5]  # use a blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set the background image with numbers and character
bg_image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(bg_image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add Title text box
title_box = slide.shapes.add_textbox(Inches(8), Inches(3), Inches(8), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_p = title_frame.add_paragraph()
title_p.text = "Are Top Grades Enough?"

# Format Title
title_p.font.bold = True
title_p.font.size = Pt(80)  # Increased font size

# Apply color to "Top Grades"
start_pos = title_p.text.find("Top Grades")
end_pos = start_pos + len("Top Grades")
title_run = title_p.runs[0]
title_run.font.color.rgb = RGBColor(0x00, 0x00, 0xFF)  # Blue
title_run = title_p.add_run()
title_run.text = title_p.text[start_pos:end_pos]
title_run.font.bold = True
title_run.font.size = Pt(80)  # Increased font size
title_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)  # Yellow

# Format subtitle
subtitle_box = slide.shapes.add_textbox(Inches(8), Inches(4.5), Inches(8), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.word_wrap = True
subtitle_p = subtitle_frame.add_paragraph()
subtitle_p.text = "A presentation by TopicsForSeminar.com"
subtitle_p.font.color.rgb = RGBColor(0x00, 0x00, 0xFF)  # Blue
subtitle_p.font.size = Pt(36)

# Add another run for "TopicsForSeminar.com" to make it yellow and underlined
subtitle_run = subtitle_p.add_run()
subtitle_run.text = "TopicsForSeminar.com"
subtitle_run.font.color.rgb = RGBColor(0xFF, 0xFF, 0x00)  # Yellow
subtitle_run.font.size = Pt(36)
subtitle_run.font.underline = True

# Add logo "Topics For Seminar"
logo_path = "../../media/image_1.jpg"
logo = slide.shapes.add_picture(logo_path, Inches(0.5), Inches(7.5), Inches(2), Inches(2))

# Add confused boy image on top of the logo
boy_image_path = "../../media/image_2.jpg"
slide.shapes.add_picture(boy_image_path, Inches(0.5), Inches(5.5), Inches(2), Inches(2))

# Save presentation
presentation.save("render.pptx")