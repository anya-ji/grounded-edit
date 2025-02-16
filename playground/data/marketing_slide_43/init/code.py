from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background to a gradient of dark teal to lighter teal
fill = slide.background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 128, 128)  # Dark teal color

# Add a rectangle with gradient
shapes = slide.shapes
left = top = Inches(0)
width = presentation.slide_width
height = presentation.slide_height

shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
fill = shape.fill
fill.gradient_stops[0].position = 0.0
fill.gradient_stops[0].color.rgb = RGBColor(0, 77, 77)  # Darker teal
fill.gradient_stops[1].position = 1.0
fill.gradient_stops[1].color.rgb = RGBColor(172, 216, 230)  # Lighter teal
shape.line.color.rgb = RGBColor(255, 255, 255) # White border

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5))
tf = title_box.text_frame
p = tf.add_paragraph()
p.text = "Home Fun:"
p.font.bold = True
p.font.size = Pt(44)
p.font.color.rgb = RGBColor(255, 255, 0)  # Yellow color
tf.word_wrap = True

# Add bullet points
left = Inches(1)
top = Inches(2.5)
width = Inches(14)
height = Inches(3)

text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame
p = tf.add_paragraph()
p.text = "Explain how market create opportunities."
p.level = 0

p = tf.add_paragraph()
p.text = "Do you think market create inequality among the buyers? If yes, how?"
p.level = 0

tf.word_wrap = True
for paragraph in tf.paragraphs:
    paragraph.font.size = Pt(24)
    paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White color

# Add a small red rectangle in the top right corner
red_rect_x = presentation.slide_width - Inches(1.5)
red_rect_y = Inches(0.5)
red_rect_width = Inches(1)
red_rect_height = Inches(0.5)
red_rect = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, red_rect_x, red_rect_y, red_rect_width, red_rect_height
)
red_fill = red_rect.fill
red_fill.solid()
red_fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red color
red_rect.line.color.rgb = RGBColor(255, 255, 255)  # White border

# Save presentation
presentation.save("render.pptx")