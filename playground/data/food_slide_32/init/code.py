from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add blank slide layout
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set background color to dark blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue

# Add FLUENTIZE logo on the left side (assuming logo is in 'fluentize_logo.png')
# Note: Replace 'fluentize_logo.png' with the actual path of your logo image
logo_path = 'fluentize_logo.png'
logo_left = Inches(1)
logo_top = Inches(3)
logo_width = Inches(3)
slide.shapes.add_picture(logo_path, logo_left, logo_top, width=logo_width)

# Add "THANK YOU!" text on the right side
thank_you_left = Inches(5)
thank_you_top = Inches(3.5)
thank_you_width = Inches(8)
thank_you_height = Inches(2)
thank_you_textbox = slide.shapes.add_textbox(thank_you_left, thank_you_top, thank_you_width, thank_you_height)
thank_you_tf = thank_you_textbox.text_frame
thank_you_p = thank_you_tf.add_paragraph()
thank_you_p.text = "THANK YOU!"
thank_you_p.font.size = Pt(60)
thank_you_p.font.bold = True
thank_you_p.font.color.rgb = RGBColor(255, 255, 255)  # White color
thank_you_p.alignment = PP_ALIGN.CENTER

# Add copyright text at the bottom
copyright_text = "Fluentize, LLC. Copyright 2022. For use only with license."
copyright_left = Inches(1)
copyright_top = Inches(7.5)
copyright_width = Inches(14)
copyright_height = Inches(1)
copyright_textbox = slide.shapes.add_textbox(copyright_left, copyright_top, copyright_width, copyright_height)
copyright_tf = copyright_textbox.text_frame
copyright_p = copyright_tf.add_paragraph()
copyright_p.text = copyright_text
copyright_p.font.size = Pt(12)
copyright_p.font.color.rgb = RGBColor(255, 255, 255)  # White color
copyright_p.alignment = PP_ALIGN.CENTER

# Save the presentation
presentation.save("render.pptx")