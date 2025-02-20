from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set slide background to black
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
title_frame = title_box.text_frame
p = title_frame.add_paragraph()
p.text = "Now: I, CAN"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)

# Add first bullet point with gradient orange box
left = Inches(1)
top = Inches(2)
width = Inches(6)
height = Inches(1.5)
orange_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
fill = orange_box.fill
fill.gradient()
stop1 = fill.gradient_stops[0]
stop1.position = 0.0
stop1.color.rgb = RGBColor(255, 165, 0)  # orange
stop2 = fill.gradient_stops[1]
stop2.position = 1.0
stop2.color.rgb = RGBColor(255, 140, 0)  # dark orange
orange_box.line.color.rgb = RGBColor(255, 165, 0)

text_frame1 = orange_box.text_frame
p1 = text_frame1.add_paragraph()
p1.text = "Define the term Market"
p1.font.size = Pt(24)
p1.font.color.rgb = RGBColor(0, 0, 0)
text_frame1.word_wrap = True

# Add second bullet point with yellow box
top = Inches(4)
yellow_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
yellow_box.fill.solid()
yellow_box.fill.fore_color.rgb = RGBColor(255, 255, 0)
yellow_box.line.color.rgb = RGBColor(255, 255, 0)

text_frame2 = yellow_box.text_frame
p2 = text_frame2.add_paragraph()
p2.text = "Explain how products reach to market"
p2.font.size = Pt(24)
p2.font.color.rgb = RGBColor(0, 0, 0)
text_frame2.word_wrap = True

# Add the image as a background on the right side
img_path = "../../media/image_0.jpg"
slide.shapes.add_picture(img_path, Inches(8), Inches(0), width=Inches(8), height=Inches(9))

# Save the presentation
presentation.save("render.pptx")