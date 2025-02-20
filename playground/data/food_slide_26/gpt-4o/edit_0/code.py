from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background color to light blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(164, 210, 225)  # Updated to #A4D2E1

# Add a title at the top
textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
text_frame = textbox.text_frame
title = text_frame.add_paragraph()
title.text = "STEP 1: Choose one of the dates below by ticking (✔) the box."
title.font.bold = True
title.font.size = Pt(32)
title.alignment = PP_ALIGN.CENTER

# Icons and images
check_icon_path = "../../media/image_0.jpg"
image_paths = [
    "../../media/image_1.jpg",  # Example image used multiple times
    "../../media/image_1.jpg",
    "../../media/image_1.jpg",
    "../../media/image_1.jpg",
    "../../media/image_1.jpg",
    "../../media/image_1.jpg",
]

# Adding checkbox and images in grid
start_x, start_y = Inches(1), Inches(1.5)
image_width, image_height = Inches(4), Inches(3)
checkbox_offset_x, checkbox_offset_y = Inches(-0.5), Inches(-0.5)

for row in range(2):
    for col in range(3):
        # Calculate position
        x = start_x + (col * (image_width + Inches(0.5)))
        y = start_y + (row * (image_height + Inches(0.5)))

        # Add image
        slide.shapes.add_picture(image_paths[col], x, y, width=image_width, height=image_height)
        
        # Add checkbox
        slide.shapes.add_picture(check_icon_path, x + checkbox_offset_x, y + checkbox_offset_y, width=Inches(0.5), height=Inches(0.5))

# Add a footer
footer_textbox = slide.shapes.add_textbox(Inches(1), Inches(7.5), Inches(14), Inches(1))
footer_text_frame = footer_textbox.text_frame
footer = footer_text_frame.add_paragraph()
footer.text = "‘DATE GONE WRONG’ STORY"
footer.font.size = Pt(18)
footer.alignment = PP_ALIGN.CENTER

# Save the presentation
presentation.save("render.pptx")