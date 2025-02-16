from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
presentation = Presentation()

# Set slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)  # White

# Add a white textbox in the middle
left = Inches(1)
top = Inches(1)
width = Inches(14)
height = Inches(5)
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
text_frame.word_wrap = True

# Add title
title = text_frame.add_paragraph()
title.text = "Human impact or influence on environment."
title.font.bold = True
title.font.color.rgb = RGBColor(0, 102, 0)  # Green font
title.font.size = Pt(32)
title.alignment = PP_ALIGN.LEFT

# Add content text
content = text_frame.add_paragraph()
content.text = (
    "We are living in this environment so we have definitely gave an impact "
    "on environment. So lets see some positive and negative impact of humans "
    "on environment."
)
content.font.color.rgb = RGBColor(0, 0, 0)  # Black font
content.font.size = Pt(20)
content.alignment = PP_ALIGN.LEFT

# Add the sunset image below the texts
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(3), Inches(5.5), Inches(10), Inches(3))

# Save the presentation
presentation.save("render.pptx")