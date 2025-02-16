from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Create a blank slide layout
slide_layout = presentation.slide_layouts[5]  
slide = presentation.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 105, 180)  # Pink color

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_tf = title_box.text_frame
title_tf.paragraphs[0].text = "Misión"
title_tf.paragraphs[0].font.bold = True
title_tf.paragraphs[0].font.size = Pt(44)
title_tf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White
title_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

# Add main text
main_text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(5))
main_tf = main_text_box.text_frame
p1 = main_tf.add_paragraph()
p1.text = "Create a world where anyone can belong anywhere"
p1.font.size = Pt(32)
p1.font.bold = True
p1.font.color.rgb = RGBColor(255, 255, 255)  # White

p2 = main_tf.add_paragraph()
p2.text = "Airbnb conecta a las personas con actividades y alojamientos en todo el mundo. El motor de la comunidad son nuestros anfitriones, que ofrecen a los viajeros la oportunidad única de descubrir su destino como un habitante más."
p2.font.size = Pt(20)
p2.font.bold = False
p2.font.color.rgb = RGBColor(255, 255, 255)  # White

# Add Airbnb logo
logo_path = "../../media/image_1.jpg"
slide.shapes.add_picture(logo_path, Inches(1), Inches(7), Inches(2), Inches(2))

# Add image of people dining together
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(12), Inches(1.5), Inches(3), Inches(6))

# Save presentation
presentation.save("render.pptx")