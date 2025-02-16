from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation and set slide dimensions
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set a lighter gradient background to improve visibility
background = slide.background
fill = background.fill
fill.gradient()
fill.gradient_stops[0].color.rgb = RGBColor(173, 216, 230)  # light blue
fill.gradient_stops[1].color.rgb = RGBColor(0, 191, 255)   # deeper light blue

# Insert the graphic on the left side
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0.5), Inches(2), Inches(3), Inches(3))

# Add text on the right side
text_box = slide.shapes.add_textbox(Inches(4), Inches(2), Inches(11.5), Inches(5))
text_frame = text_box.text_frame
p = text_frame.add_paragraph()

p.text = "Q. What is that one quality your friend has that you would like to have?"
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(255, 255, 255)  # white
p.font.bold = False  # Changed to False for a softer typeface

# Save the presentation
presentation.save("render.pptx")