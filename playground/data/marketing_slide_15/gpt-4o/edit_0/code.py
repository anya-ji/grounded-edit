from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation and set slide dimensions
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set a gradient background with shades of teal
background = slide.background
fill = background.fill
fill.gradient()
stop1 = fill.gradient_stops[0]
stop1.position = 0.0
stop1.color.rgb = RGBColor(0, 128, 128)  # teal
stop2 = fill.gradient_stops[1]
stop2.position = 1.0
stop2.color.rgb = RGBColor(0, 255, 255)  # lighter teal

# Insert the graphic on the left side
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0.5), Inches(2), Inches(3), Inches(3))

# Add text on the right side
text_box = slide.shapes.add_textbox(Inches(4), Inches(2), Inches(11.5), Inches(5))
text_frame = text_box.text_frame
p = text_frame.add_paragraph()

p.text = "Q. What is that one quality your friend has that you would like to have?"
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(255, 255, 255)  # white
p.font.bold = True

# Save the presentation
presentation.save("render.pptx")