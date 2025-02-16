from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the slide background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Title
title_text = "Content"
title_box = slide.shapes.add_textbox(Inches(1.5), Inches(0.5), Inches(13), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_p = title_frame.add_paragraph()
title_p.text = title_text
title_p.font.bold = True
title_p.font.italic = True  # Added italic formatting
title_p.font.size = Pt(44)
title_p.font.color.rgb = RGBColor(0, 0, 0)

# Vertical yellow sidebar on the left
left_bar_width = Inches(1)  
left_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), left_bar_width, Inches(9))
left_bar.fill.solid()
left_bar.fill.fore_color.rgb = RGBColor(255, 204, 0)  # Match the color scheme

# Bullet points
content_texts = [
    "Brief highlights of Elon Musk",
    "Biography",
    "Early life of Elon Musk",
    "How he came up with his ideas",
    "Elon’s current stage",
    "Obstacles that Elon faced",
    "Lessons that we can learn from Elon Musk's life",
    "Elon’s Future Plans"
]

content_box = slide.shapes.add_textbox(Inches(2.5), Inches(1.5), Inches(13), Inches(6))
content_frame = content_box.text_frame
content_frame.word_wrap = True

for item in content_texts:
    p = content_frame.add_paragraph()
    p.text = item
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_before = Pt(10)

# Save the presentation
presentation.save("render.pptx")