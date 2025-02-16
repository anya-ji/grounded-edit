from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide
slide = presentation.slides.add_slide(presentation.slide_layouts[5])

# Set the background color to dark gray
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(64, 64, 64)  # Changed to dark gray

# Add title with a dark gray background
title_height = Inches(1.5)
textbox = slide.shapes.add_textbox(0, 0, presentation.slide_width, title_height)
text_frame = textbox.text_frame
text_frame.clear()

title_fill = textbox.fill
title_fill.solid()
title_fill.fore_color.rgb = RGBColor(64, 64, 64)  # Dark gray

title_paragraph = text_frame.add_paragraph()
title_paragraph.text = "References"
title_paragraph.font.bold = True  # Ensure the title is bold
title_paragraph.font.size = Pt(44)
title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # Changed to white
title_paragraph.font.name = 'Arial'  # Adjusted font type (example)

# Add URLs below the title
content_top = title_height + Inches(0.5)
content_width = presentation.slide_width - Inches(2)
content_left = Inches(1)

textbox_content = slide.shapes.add_textbox(content_left, content_top, content_width, Inches(6))
text_frame_content = textbox_content.text_frame
text_frame_content.clear()

# Updated URLs for references
urls = [
    "https://new-url.com/article1",
    "https://new-url.com/article2",
    "https://new-url.com/article3",
    "https://new-url.com/article4",
    "https://new-url.com/article5"
]

for url in urls:
    p = text_frame_content.add_paragraph()
    p.text = url
    p.font.size = Pt(20)  # Adjusted font size to be slightly smaller than the title
    p.font.color.rgb = RGBColor(0, 0, 255)  # Change links' color to a brighter blue
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(10)  # Increased line height by adding space after each paragraph

presentation.save("render.pptx")