from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide
slide = presentation.slides.add_slide(presentation.slide_layouts[5])

# Add title with a dark gray background
title_height = Inches(1.5)
textbox = slide.shapes.add_textbox(0, 0, presentation.slide_width, title_height)
text_frame = textbox.text_frame
text_frame.clear()

title_fill = textbox.fill
title_fill.solid()
title_fill.fore_color.rgb = RGBColor(64, 64, 64)  # Dark gray

title_paragraph = text_frame.add_paragraph()
title_paragraph.text = "References"
title_paragraph.font.size = Pt(44)
title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
title_paragraph.alignment = PP_ALIGN.CENTER

# Ensure the title is bold
title_paragraph.font.bold = True

# Add URLs below the title
content_top = title_height + Inches(0.5)
content_width = presentation.slide_width - Inches(2)
content_left = Inches(1)

textbox_content = slide.shapes.add_textbox(content_left, content_top, content_width, Inches(6))
text_frame_content = textbox_content.text_frame
text_frame_content.clear()

# Example URLs for references
urls = [
    "https://example.com/article1",
    "https://example.com/article2",
    "https://example.com/article3",
    "https://example.com/article4",
    "https://example.com/article5"
]

for url in urls:
    p = text_frame_content.add_paragraph()
    p.text = url
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 0, 255)  # Blue
    p.alignment = PP_ALIGN.LEFT

presentation.save("render.pptx")