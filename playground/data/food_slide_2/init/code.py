from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a Presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with blank layout
slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color to dark blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 51, 102)  # dark blue color

# Add title text
title_textbox = slide.shapes.add_textbox(Inches(4), Inches(1), Inches(8), Inches(2))
text_frame = title_textbox.text_frame
text_frame.clear()
p = text_frame.add_paragraph()
p.text = "PREVIEW ACTIVITY"
p.font.bold = True
p.font.size = Pt(60)
p.font.color.rgb = RGBColor(255, 255, 255) # white color
p.alignment = 1  # center alignment

# Add a graphic of a light bulb with a brain inside it
# This would be a placeholder for the graphic. In practice, you'd replace with an actual graphic resource.
graphic_lb_brain = slide.shapes.add_shape(MSO_SHAPE.LIGHTBULB, Inches(3), Inches(1), Inches(1.5), Inches(3))
graphic_lb_brain.fill.solid()
graphic_lb_brain.fill.fore_color.rgb = RGBColor(255, 192, 0)  # bright color for visibility

# Add a horizontal line below the title
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1), Inches(3), Inches(14), Inches(0.1))
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(255, 255, 255)  # white color

# Placeholder for the "FLUENT" logo in the top right corner
# Replace with actual logo resource
fluent_logo = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(14.5), Inches(0.5), Inches(1), Inches(1))
fluent_logo.fill.solid()
fluent_logo.fill.fore_color.rgb = RGBColor(255, 255, 255)  # white color

# This would be the part where you add an actual logo image if available.
# e.g. slide.shapes.add_picture('fluent_logo.png', Inches(14.5), Inches(0.5), Inches(1), Inches(1))

# Save the presentation
presentation.save("render.pptx")