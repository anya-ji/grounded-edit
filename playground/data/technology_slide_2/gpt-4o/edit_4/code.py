from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Choose a blank slide layout
slide_layout = presentation.slide_layouts[5]  # slide with blank layout
slide = presentation.slides.add_slide(slide_layout)

# Title: "What is Blockchain?"
title_text = "What is Blockchain?"
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title_frame.clear()  # clear default text before we set new one

title_p = title_frame.add_paragraph()
title_p.text = title_text
title_p.font.bold = True
title_p.font.size = Pt(48)
title_p.alignment = PP_ALIGN.CENTER

# Add bullet points for blockchain content
content = [
    "Blockchain is a type of decentralized, distributed database or ledger that is used to record transactions across many computers.",
    "In blockchain, a list of records is stored as blocks and they are linked by using cryptography.",
    "Hence, we can say that blockchain is a digital record and is used for recording transactions made with cryptocurrencies such as bitcoins.",
    "If we add a new block, it can be linked with the previous block with the help of a cryptographic block generated from the previous block. By doing this, the block is permanently recorded and it is not broken.",
    "In blockchain, previous transactions' alteration is very difficult because all subsequent blocks also need to be altered. We can say that blockchain is the purest peer-to-peer database that is immutable."
]

content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(7))
content_frame = content_box.text_frame
content_frame.clear()

for point in content:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(28)
    p.space_after = Pt(10)
    p.level = 0
    p.alignment = PP_ALIGN.LEFT  # Set alignment to left

# Set a clean white background for readability
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Save to the specified path
presentation.save("render.pptx")