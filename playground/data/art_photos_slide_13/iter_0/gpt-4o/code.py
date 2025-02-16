from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

# Initialize presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5] # BLANK layout
slide = presentation.slides.add_slide(slide_layout)

# Set background image
background_img_path = "../../media/image_0.jpg"
slide.shapes.add_picture(background_img_path, Inches(0), Inches(0), width=Inches(16), height=Inches(9))

# Add white overlay for text readability
overlay = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(9))
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
overlay.fill.transparency = 0.3  # 30% transparency
overlay.line.fill.background()  # No outline

# Add main title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "The Subject and Content of Art"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(44)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text
title_frame.text_anchor = "middle"

# Add subtitle
subtitle_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(5), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Animals"
subtitle_paragraph = subtitle_frame.paragraphs[0]
subtitle_paragraph.font.size = Pt(36)
subtitle_paragraph.font.bold = True
subtitle_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text
subtitle_frame.text_anchor = "middle"

# Image paths
rabbit_img_path = "../../media/high_res_hare.jpg"  # Updated to high-resolution image
animals_person_img_path = "../../media/image_2.jpg"

# Add images side by side
left_img = slide.shapes.add_picture(rabbit_img_path, Inches(2), Inches(3), height=Inches(4))
right_img = slide.shapes.add_picture(animals_person_img_path, Inches(9), Inches(3), height=Inches(4))

# Add white borders (10 pixels)
for img in [left_img, right_img]:
    sp = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        img.left - Inches(0.1),
        img.top - Inches(0.1),
        img.width + Inches(0.2),
        img.height + Inches(0.2)
    )
    sp.fill.background()
    sp.line.color.rgb = RGBColor(255, 255, 255)
    sp.line.width = Pt(2)

# Save presentation
presentation.save("render.pptx")