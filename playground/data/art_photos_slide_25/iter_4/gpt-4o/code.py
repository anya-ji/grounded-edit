from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add the background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add title with pastel gradient background
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(14), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True

title_fill = title_box.fill
title_fill.gradient()
stops = title_fill.gradient_stops
stops[0].position = 0.0
stops[0].color.rgb = RGBColor(255, 182, 193)  # Pastel Pink
stops[1].position = 1.0
stops[1].color.rgb = RGBColor(255, 218, 185)  # Pastel Orange

p = title_frame.add_paragraph()
p.text = "References:"
p.font.bold = True
p.font.size = Pt(44)
p.font.color.rgb = RGBColor(0, 0, 0)  # Black

# Add references with white background
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(6))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content_fill = content_box.fill
content_fill.solid()
content_fill.fore_color.rgb = RGBColor(255, 255, 255)  # White

references = (
    "1. Garces, J.C., Inoc, D.S., Inocian, R.B., Labajo, O.P., Rama, F., Tiempo, A.Y., & Un, J.L. 2021. "
    "Modular Approach to Art Appreciation. Lorimar Publishing Inc.\n"
    "2. Casaul, J.A., Caslib Jr. B.N., & Garing, D.C. 2018. Art Appreciation. Rex Bookstore, Inc."
)

p = content_frame.add_paragraph()
p.text = references
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(0, 0, 0)  # Black

presentation.save("render.pptx")