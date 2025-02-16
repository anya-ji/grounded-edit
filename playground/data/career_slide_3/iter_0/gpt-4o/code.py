from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the slide background color to black
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# Add a red speech bubble shape
left = Inches(3)
top = Inches(2)
width = Inches(10)
height = Inches(5)

shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(178, 34, 34)  # Red color

# Add drop shadow effect to the shape
shape.shadow.draw_opaque = True
shape.shadow.blur_radius = Inches(0.2)
shape.shadow.distance = Inches(0.2)
shape.shadow.angle = 45
shape.shadow.color.rgb = RGBColor(0, 0, 0)  # Black shadow

# Add decorative quotation marks
quotation_top = slide.shapes.add_textbox(Inches(7.5), Inches(1.5), Inches(1), Inches(1))
quotation_p = quotation_top.text_frame.add_paragraph()
quotation_p.text = "“"
quotation_p.font.size = Pt(48)
quotation_p.font.color.rgb = RGBColor(64, 64, 64)  # Dark gray

quotation_bottom = slide.shapes.add_textbox(Inches(7.5), Inches(6.5), Inches(1), Inches(1))
quotation_p = quotation_bottom.text_frame.add_paragraph()
quotation_p.text = "”"
quotation_p.font.size = Pt(48)
quotation_p.font.color.rgb = RGBColor(64, 64, 64)  # Dark gray

# Add the main text content to the speech bubble
text_box = slide.shapes.add_textbox(left + Inches(0.5), top + Inches(0.5), width - Inches(1), height - Inches(1))
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.add_paragraph()
p.text = ("Academic success plays a significant role in society for several reasons, "
          "and its importance is reflected in various aspects of an individual's life and the broader community.")
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(255, 255, 255)  # White

# Save the presentation
presentation.save("render.pptx")