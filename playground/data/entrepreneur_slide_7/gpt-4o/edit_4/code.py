from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Initialize the presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Create a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color to white
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

# Add vertical yellow line for separation
left_line = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(2), Inches(0), Inches(0.1), Inches(9))
left_line.fill.solid()
left_line.fill.fore_color.rgb = RGBColor(255, 215, 0)  # Yellow
left_line.line.fill.background()

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Elon Musk Current Stage"
title.font.bold = True
title.font.size = Pt(36)
title.font.color.rgb = RGBColor(255, 215, 0)  # Yellow

# Add bullet points with additional space
bullet_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(7), Inches(5))
bullet_frame = bullet_box.text_frame

bullets = [
    "CEO and Chief Engineer at SpaceX",
    "CEO and Product Architect of Tesla",
    "Founder of The Boring Company",
    "Cofounder of Neuralink",
    "Cofounder of OpenAI"
]

for bullet in bullets:
    p = bullet_frame.add_paragraph()
    p.text = bullet
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black
    bullet_frame.add_paragraph()  # Add an empty paragraph for additional space

# Insert image on the right side
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(8), Inches(1), Inches(7), Inches(7))

# Save the presentation
presentation.save("render.pptx")