from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()

# Set slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add the blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)  # White color

# Add title textbox
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(2.5))
title_frame = title_box.text_frame
title_frame.text = "ENVIRONMENT"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.bold = True
title_paragraph.font.size = Pt(24)
title_paragraph.font.color.rgb = RGBColor(0, 204, 0)  # Brighter green font

# Add content textbox
content_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(5))
content_frame = content_box.text_frame

# Add bullet points
bullet_points = [
    "Everything around us is known as environment.",
    "It includes every living and non living thing.",
    "Everything which we see in our surroundings such as sunlight, atmosphere, land, water, minerals, humans, plants, animals, insects, air etc. comes in environment.",
    "It is the sum of all surroundings which we see."
]

for point in bullet_points:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black font for contrast
    p.space_after = Pt(10)  # Space between bullet points

# Add image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(11), Inches(0.5), Inches(4), Inches(3))

# Save the presentation
presentation.save("render.pptx")