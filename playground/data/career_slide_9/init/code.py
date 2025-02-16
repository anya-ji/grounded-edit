from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color to black
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# Add title text box
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
p = title_frame.add_paragraph()
p.text = "TOP GRADES ARE NOT THE SOLE DETERMINANT OF SUCCESS"
p.alignment = PP_ALIGN.CENTER
p.font.size = Pt(44)
p.font.color.rgb = RGBColor(211, 211, 211) # light gray

# Add main content text box
content_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Add formatted text
paragraph = content_frame.add_paragraph()
paragraph.text = "Success "
paragraph.font.bold = True
paragraph.font.size = Pt(36)
paragraph.font.color.rgb = RGBColor(255, 255, 0) # yellow
paragraph.alignment = PP_ALIGN.CENTER

paragraph = content_frame.add_paragraph()
paragraph.text = "is multifaceted and depends on a variety of factors, including"
paragraph.font.size = Pt(28)
paragraph.font.color.rgb = RGBColor(255, 255, 255) # white
paragraph.alignment = PP_ALIGN.CENTER

paragraph = content_frame.add_paragraph()
paragraph.text = "Skills, "
paragraph.font.bold = True
paragraph.font.size = Pt(36)
paragraph.font.color.rgb = RGBColor(0, 255, 0) # green
paragraph.alignment = PP_ALIGN.CENTER

paragraph = content_frame.add_paragraph()
paragraph.text = "Experiences, "
paragraph.font.bold = True
paragraph.font.size = Pt(36)
paragraph.font.color.rgb = RGBColor(0, 0, 255) # blue
paragraph.alignment = PP_ALIGN.CENTER

paragraph = content_frame.add_paragraph()
paragraph.text = "and Personal Attributes."
paragraph.font.bold = True
paragraph.font.size = Pt(36)
paragraph.font.color.rgb = RGBColor(255, 0, 0) # red
paragraph.alignment = PP_ALIGN.CENTER

# Add image of trophy and dartboard
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(12), Inches(6), Inches(3.5), Inches(2.5))

# Save the presentation
presentation.save("render.pptx")