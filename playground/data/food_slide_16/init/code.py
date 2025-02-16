from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Set light background color for the slide
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(230, 240, 255)  # Light blue background

# Add title in a blue textbox
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
textbox.fill.solid()
textbox.fill.fore_color.rgb = RGBColor(0, 112, 192)  # Blue color
text_frame = textbox.text_frame
p = text_frame.add_paragraph()
p.text = "PART 1"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)  # White font

# Add instruction text
instruction_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1))
text_frame = instruction_textbox.text_frame
p = text_frame.add_paragraph()
p.text = "Discuss the questions below."
p.font.size = Pt(24)

# Add conversation icon
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0.5), Inches(2.5), Inches(1.5), Inches(1.5))

# Add questions
questions_text = (
    "1. What do you think about Sarah’s etiquette on the first date? "
    "What about Joey’s etiquette on the second date? Whose etiquette was worse and why?\n"
    "2. What parts of the video do you think were the funniest and why?"
)
questions_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(10), Inches(3))
text_frame = questions_textbox.text_frame
p = text_frame.add_paragraph()
p.text = questions_text
p.font.size = Pt(20)
p.word_wrap = True

# Add a circular image of Joey from Friends
joey_image_path = "../../media/image_1.jpg"
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(12), Inches(1), Inches(3), Inches(3))
fill = circle.fill
fill.solid()
fill.user_picture(joey_image_path)

# Add play button icon labeled "VIEWING FOLLOW-UP"
play_button = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(0.5), Inches(7.5), Inches(1), Inches(1))
fill = play_button.fill
play_button.rotation = 90
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red color
play_button.text = "▶"
play_button.text_frame.paragraphs[0].font.size = Pt(32)
play_textbox = slide.shapes.add_textbox(Inches(1.2), Inches(8), Inches(3), Inches(1))
p = play_textbox.text_frame.add_paragraph()
p.text = "VIEWING FOLLOW-UP"
p.font.size = Pt(16)

# Save presentation
presentation.save("render.pptx")