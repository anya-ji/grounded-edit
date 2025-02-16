from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

# Create a presentation with a specified slide size
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set slide background to black
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# Add the main title in the center
title = slide.shapes.add_textbox(Inches(4), Inches(4), Inches(8), Inches(1))
title_frame = title.text_frame
title_frame.text = "THE ADVANTAGES OF ACHIEVING TOP GRADES"
p = title_frame.paragraphs[0]
p.font.bold = True
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(0, 255, 0)

# Add dashed lines around the title
line_xml = f'<a:ln {nsdecls("a")} w="5000"><a:prstDash val="dash"/></a:ln>'
line = slide.shapes.add_line(Inches(1), Inches(4.5), Inches(15), Inches(4.5))
line.width = Inches(0.2)
line.fill.solid()
line.fill.fore_color.rgb = RGBColor(255, 255, 0)
line.line_format = parse_xml(line_xml)

# Add light blue circular nodes along the line
node_positions = [2.5, 11.5] # Example positions
for pos in node_positions:
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(pos), Inches(4.45), Inches(0.3), Inches(0.3))
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(173, 216, 230)

# Define positions and texts for the text boxes
texts_top = [
    ("Educational Opportunities", RGBColor(255, 0, 0)),
    ("Career Advancement", RGBColor(255, 255, 0)),
    ("Increased Earning Potential", RGBColor(0, 0, 255)),
    ("Personal Fulfillment", RGBColor(255, 192, 203)),
    ("Intellectual Development", RGBColor(0, 255, 0)),
    ("Access to Competitive Programs", RGBColor(255, 165, 0))
]

texts_bottom = [
    ("Networking Opportunities", RGBColor(255, 165, 0)),
    ("Scholarship Eligibility", RGBColor(144, 238, 144)),
    ("Research and Leadership Opportunities", RGBColor(255, 192, 203)),
    ("Graduate School Options", RGBColor(0, 0, 255)),
    ("Recognition and Awards", RGBColor(255, 255, 0)),
    ("Personal Growth and Discipline", RGBColor(255, 0, 0))
]

# Set positions for text boxes top and bottom
positions_top = [(2, 1), (5.5, 1), (9, 1), (12, 1), (3.5, 2.5), (10, 2.5)]
positions_bottom = [(2, 6), (5.5, 6), (9, 6), (12, 6), (3.5, 7.5), (10, 7.5)]

# Add text boxes for top items
for i, (text, color) in enumerate(texts_top):
    tx_box = slide.shapes.add_textbox(Inches(positions_top[i][0]), Inches(positions_top[i][1]), Inches(3), Inches(1))
    tx_box.fill.solid()
    tx_box.fill.fore_color.rgb = color
    tx_frame = tx_box.text_frame
    tx_frame.text = text
    p = tx_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True

# Add text boxes for bottom items
for i, (text, color) in enumerate(texts_bottom):
    tx_box = slide.shapes.add_textbox(Inches(positions_bottom[i][0]), Inches(positions_bottom[i][1]), Inches(3), Inches(1))
    tx_box.fill.solid()
    tx_box.fill.fore_color.rgb = color
    tx_frame = tx_box.text_frame
    tx_frame.text = text
    p = tx_frame.paragraphs[0]
    p.font.size = Pt(14)
    p.font.bold = True

# Save the presentation
presentation.save("render.pptx")