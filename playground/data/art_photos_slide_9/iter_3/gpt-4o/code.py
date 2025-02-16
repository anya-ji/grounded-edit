from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Change slide background to solid light gray color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(211, 211, 211)  # Light gray color

# Add background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Title text box
title_box = slide.shapes.add_textbox(Inches(2), Inches(1), Inches(12), Inches(1))
title_frame = title_box.text_frame
title_paragraph = title_frame.add_paragraph()
title_paragraph.text = "Representational or Objective Art"
title_paragraph.font.bold = True
title_paragraph.font.size = Pt(40)
title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # Changed text to white for readability

# Subheading text box
subheading_box = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(12), Inches(1))
subheading_frame = subheading_box.text_frame
subheading_box.fill.solid()
subheading_box.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow background
subheading_paragraph = subheading_frame.add_paragraph()
subheading_paragraph.text = "B. Portrait"
subheading_paragraph.font.bold = True
subheading_paragraph.font.underline = True
subheading_paragraph.font.size = Pt(30)
subheading_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Content text box
content_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(12), Inches(3))
content_frame = content_box.text_frame

content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

bullet_points = [
    "✓ is a realistic presentation of a human face",
    "✓ The purpose may be to show the resemble, personality, or disposition of the individual"
]

for point in bullet_points:
    paragraph = content_frame.add_paragraph()
    paragraph.text = point
    paragraph.font.size = Pt(24)

# Smaller title on top-left
small_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(5), Inches(1))
small_title_frame = small_title_box.text_frame
small_title_paragraph = small_title_frame.add_paragraph()
small_title_paragraph.text = "The Subject and Content of Art"
small_title_paragraph.font.size = Pt(18)
small_title_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Save presentation
presentation.save("render.pptx")