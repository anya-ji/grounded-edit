from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Title section
title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(16), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(211, 211, 211)  # Light gray background

title = title_frame.add_paragraph()
title.text = "Branding"
title.font.size = Pt(44)
title.font.bold = True

# Year and paragraph
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(2))
content_frame = content_box.text_frame
content_frame.word_wrap = True

year_paragraph = content_frame.add_paragraph()
year_paragraph.text = "2013"
year_paragraph.font.size = Pt(36)

main_paragraph = content_frame.add_paragraph()
main_paragraph.text = (
    "A brand that could express the vision of its founders and guide and support "
    "its business through exponential growth and whatever the future may hold."
)
main_paragraph.font.size = Pt(20)

# Symbols paragraph (changed color to green)
symbols_paragraph = content_frame.add_paragraph()
symbols_paragraph.text = "PEOPLE + PLACES + LOVE + AIRBNB = Bélo"
symbols_paragraph.font.size = Pt(20)
symbols_paragraph.font.color.rgb = RGBColor(0, 128, 0)  # Changed to green

# Visual Elements
slide.shapes.add_picture("../../media/image_0.jpg", Inches(1), Inches(3.5), Inches(10), Inches(2))

# Additional paragraphs
additional_box = slide.shapes.add_textbox(Inches(1), Inches(6), Inches(14), Inches(2))
additional_frame = additional_box.text_frame
additional_frame.word_wrap = True

additional_paragraph = additional_frame.add_paragraph()
additional_paragraph.text = (
    "This higher purpose shifted Airbnb’s narrative to focus on belonging, with language "
    "and messaging elevating the conversation away from insurance and bookings to warmth and welcome."
)
additional_paragraph.font.size = Pt(18)

spanish_paragraph = additional_frame.add_paragraph()
spanish_paragraph.text = (
    "Este propósito superior cambió la narrativa de Airbnb para centrarse en la pertenencia, con "
    "un lenguaje y mensajes que elevaban la conversación más allá del seguro y las reservas a calidez y bienvenida."
)
spanish_paragraph.font.size = Pt(18)

# Highlighted Text Box
highlight_box = slide.shapes.add_textbox(Inches(12), Inches(4), Inches(3), Inches(1))
highlight_box.fill.solid()
highlight_box.fill.fore_color.rgb = RGBColor(255, 165, 0)  # Bright orange background
highlight_frame = highlight_box.text_frame
highlight_frame.word_wrap = True

highlight_paragraph = highlight_frame.add_paragraph()
highlight_paragraph.text = "The Bélo tells the story of the company in one symbol."
highlight_paragraph.font.size = Pt(18)
highlight_paragraph.font.bold = True
highlight_paragraph.font.color.rgb = RGBColor(0, 128, 0)  # Changed to green

# Color Palette
slide.shapes.add_picture("../../media/image_1.jpg", Inches(12), Inches(7), Inches(3), Inches(2))

# Save the presentation
presentation.save("render.pptx")