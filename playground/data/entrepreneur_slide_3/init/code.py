from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add the title to the slide
title_text = "Brief highlights of Elon Musk"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_paragraph = title_frame.add_paragraph()
title_paragraph.text = title_text
title_paragraph.font.size = Pt(44)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
title_frame.word_wrap = True

# Key events for the timeline
events = [
    ("1971", "Born in South Africa"),
    ("1983", "Sold his first video game"),
    ("1999", "Sold Zip2, first company to Compaq"),
    ("2002", "eBay acquired PayPal"),
    ("2002", "Founded SpaceX"),
    ("2004", "Join with Tesla"),
    ("2006", "Cofounded Solarcity"),
    ("2013", "Develop the Hyperloop concept"),
    ("2016", "Cofounded Neuralink"),
    ("2021", "Become world’s richest man"),
]

# Space out timeline elements evenly
x_offset = 0.75
y_position = 2.5
shape_radius = 0.25
space_between = (presentation.slide_width.inches - x_offset * 2) / (len(events) - 1)

# Create the timeline
for i, (year, description) in enumerate(events):
    x_position = x_offset + i * space_between

    # Draw yellow circles for each year
    circle = slide.shapes.add_shape(
        auto_shape_type_id=1, # 1 for elliptical shape
        left=Inches(x_position),
        top=Inches(y_position),
        width=Inches(shape_radius),
        height=Inches(shape_radius)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow color
    circle.line.color.rgb = RGBColor(0, 0, 0)  # Black outline

    # Add year text inside the circle
    year_box = slide.shapes.add_textbox(
        left=Inches(x_position),
        top=Inches(y_position + 0.15),
        width=Inches(shape_radius),
        height=Inches(0.5)
    )
    year_frame = year_box.text_frame
    year_paragraph = year_frame.add_paragraph()
    year_paragraph.text = year
    year_paragraph.font.size = Pt(14)
    year_paragraph.font.bold = True
    year_paragraph.alignment = 1  # Center alignment

    # Add event description below each year
    desc_box = slide.shapes.add_textbox(
        left=Inches(x_position - 0.5),
        top=Inches(y_position + 0.7),
        width=Inches(1.5),
        height=Inches(0.5)
    )
    desc_frame = desc_box.text_frame
    desc_paragraph = desc_frame.add_paragraph()
    desc_paragraph.text = description
    desc_paragraph.font.size = Pt(10)
    desc_paragraph.alignment = 1  # Center alignment
    desc_frame.word_wrap = True

# Save the presentation
presentation.save("render.pptx")