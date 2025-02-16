from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set light blue background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light Blue

# Add the slide title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
title = title_box.text_frame.add_paragraph()
title.text = "Three perspectives on global cultural flows:"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)  # Black
title.alignment = PP_ALIGN.LEFT

# Add the subheading
subheading_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(1))
subheading = subheading_box.text_frame.add_paragraph()
subheading.text = "2. CULTURAL HYBRIDIZATION"
subheading.font.size = Pt(28)
subheading.font.bold = True
subheading.font.color.rgb = RGBColor(0, 0, 0)  # Black
subheading.alignment = PP_ALIGN.LEFT

# Bullet points
bullet_points = [
    "a process by which a cultural element blends into another culture by modifying the element to fit cultural norms.",
    "It is actually an integration of local and global cultures.",
    "A key concept is “glocalization” or the interpenetration of the global and local resulting in unique outcomes in different geographic areas."
]

# Add bullets with checkmarks
left_inset = 1.5
top_inset = 3
for i, point in enumerate(bullet_points):
    checkmark_box = slide.shapes.add_textbox(Inches(left_inset), Inches(top_inset + i * 1), Inches(0.5), Inches(0.5))
    checkmark = checkmark_box.text_frame.add_paragraph()
    checkmark.text = "✓"
    checkmark.font.size = Pt(20)
    checkmark.font.color.rgb = RGBColor(0, 0, 0)  # Black
    
    text_box = slide.shapes.add_textbox(Inches(left_inset + 0.5), Inches(top_inset + i * 1), Inches(13), Inches(1))
    text = text_box.text_frame.add_paragraph()
    text.text = point
    text.font.size = Pt(20)
    text.font.color.rgb = RGBColor(0, 0, 0)  # Black
    text.alignment = PP_ALIGN.LEFT

# Save presentation
presentation.save("render.pptx")