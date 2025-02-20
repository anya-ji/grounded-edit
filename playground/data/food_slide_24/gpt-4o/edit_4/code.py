from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(211, 211, 211)  # Light gray background

# Add title text
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
title = title_box.text_frame.add_paragraph()
title.text = "PART 4"
title.font.size = Pt(40)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 139)  # Blue color

# Add instruction text
instruction_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(10), Inches(1))
instruction = instruction_box.text_frame.add_paragraph()
instruction.text = "Tell the short story below out loud using the pronunciation rule you studied on the previous slide."
instruction.font.size = Pt(20)
instruction.font.bold = True

# Add speaking icon
icon_path = "../../media/image_1.jpg"
slide.shapes.add_picture(icon_path, Inches(4), Inches(0.5), Inches(1), Inches(1))

# Add story text
story_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(10), Inches(2))
story = story_box.text_frame.add_paragraph()
story.text = (
    "Last week, I went on a date. We were skating in the park, we were laughing, "
    "and the guy was making some really funny jokes...and then, out of nowhere, "
    "I fell down and I broke my leg!"
)
story.font.size = Pt(20)

# Add skating image
skating_image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(skating_image_path, Inches(11), Inches(1.5), Inches(4), Inches(3))

# Add audio instruction text to bottom-left corner
audio_instruction_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(10), Inches(1))
audio_instruction = audio_instruction_box.text_frame.add_paragraph()
audio_instruction.text = "Play and listen to the audio clip to check your answers. Repeat if necessary."
audio_instruction.font.size = Pt(18)

# Add audio icon at bottom-left corner
audio_icon_path = "../../media/image_2.jpg"
slide.shapes.add_picture(audio_icon_path, Inches(10.5), Inches(7), Inches(1), Inches(1))

# Save presentation
presentation.save("render.pptx")