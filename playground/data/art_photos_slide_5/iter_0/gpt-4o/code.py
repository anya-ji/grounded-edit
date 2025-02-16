from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add the image as the slide background
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, 0, 0, presentation.slide_width, presentation.slide_height)

# Add a smaller title on the top-left corner
left_corner_textbox = slide.shapes.add_textbox(Inches(0.3), Inches(0.3), Inches(5), Inches(1))
left_corner_textbox.fill.solid()
left_corner_textbox.fill.fore_color.rgb = RGBColor(255, 255, 102)  # Yellow
tf_left = left_corner_textbox.text_frame
p_left = tf_left.add_paragraph()
p_left.text = "The Subject and Content of Art"
p_left.font.size = Pt(20)
p_left.font.bold = True
p_left.font.color.rgb = RGBColor(0, 0, 0)  # Black

# Add the main title with modified color
title_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(1.5))
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Brightened Yellow
tf_title = title_box.text_frame
p_title = tf_title.add_paragraph()
p_title.text = "Two kinds of Art as to Subject"
p_title.font.size = Pt(44)
p_title.font.bold = True
p_title.font.color.rgb = RGBColor(0, 0, 0)  # Black
p_title.alignment = PP_ALIGN.CENTER

# Add the main content
content_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(4))
content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White
tf_content = content_box.text_frame
p_content1 = tf_content.add_paragraph()
p_content1.text = "1. Representational or Objective Art"
p_content1.font.size = Pt(28)
p_content1.font.bold = True
p_content1.font.color.rgb = RGBColor(0, 0, 0)  # Black

p_content2 = tf_content.add_paragraph()
p_content2.text = "   - The subject is directly represented in the artwork like the subjects seen in portraits and other realistic visual art presentations."
p_content2.font.size = Pt(24)
p_content2.font.color.rgb = RGBColor(0, 0, 0)  # Black

# Save the presentation
presentation.save("render.pptx")