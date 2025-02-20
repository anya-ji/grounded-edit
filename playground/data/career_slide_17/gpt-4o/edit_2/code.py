from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Define slide layout and add a new slide
slide_layout = presentation.slide_layouts[5]  # choose a blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(50, 50, 50)  # dark gray background for contrast

# Add the main title
title_text = "ACADEMIC WORLD AND PROFESSIONAL WORLD: DEMANDS AND CHARACTERISTICS"
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_p = title_frame.add_paragraph()
title_p.text = title_text
title_p.font.bold = True
title_p.font.size = Pt(36)  # Changed to 36pt
title_p.alignment = 1  # center align
title_p.font.color.rgb = RGBColor(255, 255, 255)  # white font

# Left section: Academic World
academic_title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(7), Inches(1))
academic_title_frame = academic_title.text_frame
academic_title_p = academic_title_frame.add_paragraph()
academic_title_p.text = "ACADEMIC WORLD"
academic_title_p.font.size = Pt(24)
academic_title_p.font.bold = True
academic_title_p.font.color.rgb = RGBColor(0, 0, 255)  # blue font

academic_outcome = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(7), Inches(1))
academic_outcome_frame = academic_outcome.text_frame
academic_outcome_p = academic_outcome_frame.add_paragraph()
academic_outcome_p.text = "Outcome"
academic_outcome_p.font.size = Pt(20)
academic_outcome_p.font.bold = True
academic_outcome_p.font.color.rgb = RGBColor(255, 255, 0)  # yellow font

academic_details = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(7), Inches(2))
academic_details_frame = academic_details.text_frame
academic_details_p = academic_details_frame.add_paragraph()
academic_details_p.text = "Degrees and Certifications"
academic_details_p.font.size = Pt(18)
academic_details_p.font.bold = True

academic_details_bullet = academic_details_frame.add_paragraph()
academic_details_bullet.text = (
    "Academic success leads to degrees and certifications (e.g., diplomas, degrees, Ph.Ds.), "
    "which serve as qualifications for future opportunities."
)
academic_details_bullet.font.size = Pt(16)
academic_details_bullet.level = 1

# Right section: Professional World
professional_title = slide.shapes.add_textbox(Inches(8), Inches(2), Inches(7), Inches(1))
professional_title_frame = professional_title.text_frame
professional_title_p = professional_title_frame.add_paragraph()
professional_title_p.text = "PROFESSIONAL WORLD"
professional_title_p.font.size = Pt(24)
professional_title_p.font.bold = True
professional_title_p.font.color.rgb = RGBColor(0, 128, 0)  # green font

professional_outcome = slide.shapes.add_textbox(Inches(8), Inches(3), Inches(7), Inches(1))
professional_outcome_frame = professional_outcome.text_frame
professional_outcome_p = professional_outcome_frame.add_paragraph()
professional_outcome_p.text = "Outcome"
professional_outcome_p.font.size = Pt(20)
professional_outcome_p.font.bold = True
professional_outcome_p.font.color.rgb = RGBColor(255, 255, 0)  # yellow font

professional_achievement = slide.shapes.add_textbox(Inches(8), Inches(4), Inches(7), Inches(2))
professional_achievement_frame = professional_achievement.text_frame
professional_achievement_p = professional_achievement_frame.add_paragraph()
professional_achievement_p.text = "Achievement and Success"
professional_achievement_p.font.size = Pt(18)
professional_achievement_p.font.bold = True

professional_achievement_bullet = professional_achievement_frame.add_paragraph()
professional_achievement_bullet.text = (
    "Professional success is measured by job performance, career advancement, and the impact of one's work on "
    "organizations and clients."
)
professional_achievement_bullet.font.size = Pt(16)
professional_achievement_bullet.level = 1

professional_compensation = professional_achievement_frame.add_paragraph()
professional_compensation.text = "Compensation"
professional_compensation.font.size = Pt(18)
professional_compensation.font.bold = True

professional_compensation_bullet = professional_achievement_frame.add_paragraph()
professional_compensation_bullet.text = (
    "Compensation and benefits are typically tied to professional success and experience."
)
professional_compensation_bullet.font.size = Pt(16)
professional_compensation_bullet.level = 1

# Save the presentation
presentation.save("render.pptx")