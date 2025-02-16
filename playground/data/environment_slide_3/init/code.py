from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml
import os

# Define the path for saving the presentation
save_path = "examples/environment/slide_3/gpt_4o.pptx"

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Select a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color (light green with a subtle hexagonal pattern)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(200, 230, 200)

# Add a large white textbox in the middle
text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(9), Inches(5))
text_frame = text_box.text_frame

# Title section
title_p = text_frame.add_paragraph()
title_p.text = "Types of environments"
title_p.font.bold = True
title_p.font.size = Pt(36)
title_p.font.color.rgb = RGBColor(34, 139, 34)  # Green color
title_p.alignment = PP_ALIGN.CENTER

# Content section
bullet_points = [
    "Natural environment",
    "Human environment",
    "Physical environment"
]

for point in bullet_points:
    p = text_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(28)
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black color
    p.alignment = PP_ALIGN.LEFT

# Image section - Place countryside image on the right side
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(10.5), Inches(2), Inches(5), Inches(5))

# Save the presentation
presentation.save(save_path)