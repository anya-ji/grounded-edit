from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color to dark blue
background_shape = slide.shapes.add_shape(
    shape_type=1,  # Rectangle
    left=0, top=0, width=presentation.slide_width, height=presentation.slide_height
)
background_shape.fill.solid()
background_shape.fill.fore_color.rgb = RGBColor(0, 0, 139)  # Dark blue color
background_shape.line.fill.background()  # No border

# Add the title
text_box = slide.shapes.add_textbox(Inches(5), Inches(4), Inches(6), Inches(1.5))
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.add_paragraph()
p.text = "PRONUNCIATION ACTIVITY"
p.font.bold = True
p.font.size = Pt(44)
p.font.color.rgb = RGBColor(255, 255, 255)  # White color

# Add a graphic of a speaking head
speaking_head_path = "speaking_head.png"  # Add your image path here
slide.shapes.add_picture(speaking_head_path, Inches(3), Inches(4), Inches(2), Inches(2))

# Add the FLUENT logo in the top right corner
fluent_logo_path = "fluent_logo.png"  # Add your image path here
slide.shapes.add_picture(fluent_logo_path, Inches(13), Inches(0.5), Inches(2.5), Inches(1))

# Save the presentation
presentation.save("render.pptx")