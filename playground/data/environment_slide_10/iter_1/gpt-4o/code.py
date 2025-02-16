from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Length

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set a light green background with a subtle hexagonal pattern
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(204, 255, 204)  # Light green

# Add a large white textbox in the center
left = Inches(2)
top = Inches(3)
width = Inches(12)
height = Inches(3)
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
textbox.fill.solid()
textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White

# Title Section: Add "THANK YOU"
p = text_frame.add_paragraph()
p.text = "THANK YOU"
p.font.bold = True
p.font.size = Pt(100)  # Large font size
p.font.color.rgb = RGBColor(51, 153, 51)  # Light green font
p.alignment = PP_ALIGN.CENTER  # Center the text

# Apply a stylized font and add shadow effect
p.font.name = 'Comic Sans MS'  # Example of a decorative font
p.shadow = True  # Enable shadow effect

# Save the presentation
presentation.save("render.pptx")