from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color to dark blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 128)

# Add movie icon (assuming you have an image file named 'movie_icon.png')
icon_path = 'movie_icon.png'
slide.shapes.add_picture(icon_path, Inches(0.5), Inches(0.5), Inches(1), Inches(1))

# Add title "VIEWING ACTIVITY"
title_box = slide.shapes.add_textbox(Inches(1.6), Inches(0.5), Inches(6), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "VIEWING ACTIVITY"
title_p = title_frame.paragraphs[0]
title_p.font.bold = True
title_p.font.size = Pt(28)
title_p.font.color.rgb = RGBColor(255, 255, 255)

# Add instructional text
instruction_box = slide.shapes.add_textbox(Inches(1.6), Inches(1.5), Inches(6), Inches(1))
instruction_frame = instruction_box.text_frame
instruction_frame.text = "Watch the first 2 minutes of the video."
instruction_p = instruction_frame.paragraphs[0]
instruction_p.font.size = Pt(20)
instruction_p.font.color.rgb = RGBColor(255, 255, 255)

# Add text box on the right side
text_box = slide.shapes.add_textbox(Inches(9), Inches(2.5), Inches(6), Inches(1.5))
text_frame = text_box.text_frame
text_frame.text = "Friends | Joey Doesn’t Share Food!"
text_p = text_frame.paragraphs[0]
text_p.font.size = Pt(24)
text_p.font.color.rgb = RGBColor(255, 255, 255)

# Insert the provided image of Joey from "Friends"
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(9), Inches(4), Inches(6), Inches(3))

# Insert a logo or icon at the bottom left (assuming you have an image file named 'activity_icon.png')
logo_path = 'activity_icon.png'
slide.shapes.add_picture(logo_path, Inches(0.5), Inches(7), Inches(1.5), Inches(1.5))

# Save the presentation
presentation.save("render.pptx")