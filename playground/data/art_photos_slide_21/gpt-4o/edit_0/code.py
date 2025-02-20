from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_AUTO_SHAPE_TYPE

# Create a presentation object
presentation = Presentation()

# Set the slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background color of the slide to a soft off-white shade
background_fill = slide.background.fill
background_fill.solid()
background_fill.fore_color.rgb = RGBColor(255, 255, 240)  # Soft off-white color

# Add the background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, 0, 0, presentation.slide_width, presentation.slide_height)

# Add the smaller title on the top-left corner
title_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(6), Inches(1))
title_frame = title_textbox.text_frame
title_frame.word_wrap = True

title_p = title_frame.add_paragraph()
title_p.text = "The Subject and Content of Art."
title_p.font.size = Pt(24)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(0, 0, 0)

# Set the background color of the title textbox
fill = title_textbox.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 0)

# Add the main title
main_title_textbox = slide.shapes.add_textbox(Inches(2), Inches(2), Inches(12), Inches(2))
main_title_frame = main_title_textbox.text_frame
main_title_frame.word_wrap = True

main_title_p = main_title_frame.add_paragraph()
main_title_p.text = "Three levels of meaning"
main_title_p.font.size = Pt(32)
main_title_p.font.bold = True
main_title_p.font.color.rgb = RGBColor(0, 0, 0)

# Set the background color of the main title textbox
main_title_fill = main_title_textbox.fill
main_title_fill.solid()
main_title_fill.fore_color.rgb = RGBColor(255, 255, 0)

# Add the main content
content_textbox = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(12), Inches(4))
content_frame = content_textbox.text_frame
content_frame.word_wrap = True

content_p = content_frame.add_paragraph()
content_p.text = ("3. Subjective meaning - refers to the individual meaning deliberately and "
                  "instinctively expressed by the artist using personal symbolism that stems from "
                  "his own alliance with certain objects, actions, or colors with past experiences.")
content_p.font.size = Pt(20)
content_p.font.color.rgb = RGBColor(0, 0, 0)

# Set the background color of the content textbox
content_fill = content_textbox.fill
content_fill.solid()
content_fill.fore_color.rgb = RGBColor(255, 255, 255)

# Save the presentation
presentation.save("render.pptx")