from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Define slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add title text
title_text = "NYC"
text_box_height = Inches(3)  # Height of the text box
text_box_top = (presentation.slide_height - text_box_height) / 2  # Center the text box vertically
text_box = slide.shapes.add_textbox(Inches(0), text_box_top, Inches(16), text_box_height)
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = title_text
p.font.bold = True
p.font.size = Pt(96)
p.font.color.rgb = RGBColor(255, 255, 255)  # White color
text_frame.paragraphs[0].alignment = 1  # Center alignment

# Add Airbnb logo
logo_path = image_path  # Replace with actual path if different
logo = slide.shapes.add_picture(logo_path, Inches(6.5), Inches(6), Inches(3), Inches(1.5))

# Save presentation
presentation.save("render.pptx")