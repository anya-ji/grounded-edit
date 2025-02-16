from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image, ImageEnhance

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Define slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add background image with color enhancement
image_path = "../../media/image_0.jpg"
image = Image.open(image_path)

# Enhance colors (increase yellows and oranges)
enhancer = ImageEnhance.Color(image)
image = enhancer.enhance(1.5)  # Adjust the factor as necessary for more saturation

# Save the enhanced image temporarily
enhanced_image_path = "../../media/enhanced_image.jpg"
image.save(enhanced_image_path)

slide.shapes.add_picture(enhanced_image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add title text
title_text = "NYC"  # Accent removed from "C"
text_box = slide.shapes.add_textbox(Inches(0), Inches(3), Inches(16), Inches(2))  # Adjusted vertical position
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = title_text
p.font.bold = True
p.font.size = Pt(72)  # Font size
p.font.color.rgb = RGBColor(255, 255, 255)  # White color for all letters
p.font.name = 'Arial'  # Change font to a more uniform style
text_frame.paragraphs[0].alignment = 1  # Center alignment

# Add Airbnb logo (keep original logo path)
logo_path = "../../media/image_0.jpg"  # Keep using the background image as logo
logo = slide.shapes.add_picture(logo_path, Inches(6.5), Inches(6), Inches(3), Inches(1.5))

# Save presentation
presentation.save("render.pptx")