from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide layout
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add a yellow vertical bar on the left side
left_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.5), Inches(9)
)
fill = left_bar.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 223, 0)  # Yellow color

# Add title text box on the left side
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(2))
title_frame = title_box.text_frame
title_frame.text = "Obstacles that Elon faced"
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.size = Pt(44)
title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black color
# Ensure the title is left-aligned
for paragraph in title_frame.paragraphs:
    paragraph.alignment = 0  # 0 represents left alignment

# Add bullet points on the right side
content_box = slide.shapes.add_textbox(Inches(7), Inches(1), Inches(8), Inches(7))
content_frame = content_box.text_frame

bullets = [
    "Rejection from Netscape and Ousted at ZIP2: In 1996, Elon Musk was ousted as the CEO of the company.",
    "PayPal was voted 'worst business concept of the year’: When PayPal was issued as security software for the PalmPilot and other portable devices in 1999, it was voted the worst business idea of the year.",
    "Tesla crisis: The company also had serious financial problems, and it was on the verge of going out of business.",
    "SpaceX failure: The failure of the company's first three launches generated skepticism among many investors, but the fourth one got a huge success.",
    "Running on empty: 2008 was the worst year of his life. After investing his whole money in Tesla and SpaceX, Musk was reliant on personal loans from friends."
]

for bullet in bullets:
    p = content_frame.add_paragraph()
    p.text = bullet
    p.font.size = Pt(20)
    p.space_after = Pt(14)
    p.level = 0
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black color

# Save the presentation
presentation.save("render.pptx")