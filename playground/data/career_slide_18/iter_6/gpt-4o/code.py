from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a new slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set dark background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(30, 30, 30)  # Dark background

# Add main title with adjusted spacing
title_text = "ACADEMIC      WORLD        AND        PROFESSIONAL          WORLD"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_paragraph = title_frame.add_paragraph()
title_paragraph.text = title_text

# Split title into parts to have different sizes
academic_world = title_frame.add_paragraph()
academic_world.text = "ACADEMIC WORLD"
academic_world.font.size = Pt(58)  # Larger font size for emphasis
academic_world.font.bold = True
academic_world.font.color.rgb = RGBColor(255, 255, 255)  # White color

professional_world = title_frame.add_paragraph()
professional_world.text = "PROFESSIONAL WORLD"
professional_world.font.size = Pt(58)  # Larger font size for emphasis
professional_world.font.bold = True
professional_world.font.color.rgb = RGBColor(255, 255, 255)  # White color

# Add the remaining text "AND" with original size
and_paragraph = title_frame.add_paragraph()
and_paragraph.text = "AND"
and_paragraph.font.size = Pt(44)  # Original font size
and_paragraph.font.bold = True
and_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White color

# Adjust line spacing for title
title_paragraph.space_after = Pt(10)  # Set desired spacing after the title

# Add subtitle
subtitle_text = "DEMANDS AND CHARACTERISTICS"
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(15), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.word_wrap = True
subtitle_paragraph = subtitle_frame.add_paragraph()
subtitle_paragraph.text = subtitle_text
subtitle_paragraph.font.size = Pt(32)
subtitle_paragraph.font.bold = True
subtitle_paragraph.font.color.rgb = RGBColor(255, 255, 0)  # Yellow color

# Add content text
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(15), Inches(6))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Define text paragraphs
content_text = [
    ("In summary,", Pt(28)),
    ("• Academic world focuses on education, intellectual development, and theoretical knowledge.", Pt(24)),
    ("• Assessment primarily aimed at understanding and demonstrating knowledge.", Pt(24)),
    ("• The professional world, on the other hand, emphasizes the application of knowledge and skills to achieve specific goals.", Pt(24)),
    ("• Assessment based on job performance and outcomes.", Pt(24)),
    ("• While both worlds contribute to an individual's overall development, they have distinct demands and expectations.", Pt(24)),
    ("• Transitioning from the academic to the professional world often requires adapting to a different set of priorities and challenges.", Pt(24)),
]

# Add text to content frame with consistent line spacing
for text, size in content_text:
    paragraph = content_frame.add_paragraph()
    paragraph.text = text
    paragraph.font.size = size
    paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White color
    paragraph.space_after = Pt(5)  # Set consistent spacing after each bullet point

# Save the presentation
presentation.save("render.pptx")