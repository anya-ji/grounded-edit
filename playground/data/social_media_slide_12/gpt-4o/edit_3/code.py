from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation object with desired dimensions
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background color to dark blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue

# Add title with white text
title_text = "3 Factors that have affected the process of Economic Globalization"
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5))
tf = title_box.text_frame
p = tf.add_paragraph()
p.text = title_text
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)  # White

# Add a white background rectangle for the main content area
content_box_width = Inches(14)
content_box_height = Inches(5)
content_left = Inches(1)
content_top = Inches(2.5)

content_rect = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, content_left, content_top,
    content_box_width, content_box_height
)
content_rect.fill.solid()
content_rect.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White

# Add the content text within the white background
content_text = (
    "1. Improvements in transportation and communication technology have reduced "
    "the cost of transporting goods, services and factors of production and communicating "
    "economically useful knowledge and technology.\n\n"
    "2. Tastes of individuals and societies have generally but not universally favored "
    "taking advantage of the opportunities provided by declining."
)

content_tf = content_rect.text_frame
content_tf.word_wrap = True
content_p = content_tf.add_paragraph()
content_p.text = content_text
content_p.font.size = Pt(20)
content_p.font.color.rgb = RGBColor(0, 0, 0)  # Black for bullet point text

# Save the presentation
presentation.save("render.pptx")