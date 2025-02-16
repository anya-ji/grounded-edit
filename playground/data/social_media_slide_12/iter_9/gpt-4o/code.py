from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation object with desired dimensions
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background color to dark blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue

# Add header with the title
header_text = "3 Factors that have affected the process of Economic Globalization"
header_box = slide.shapes.add_textbox(Inches(1), Inches(0), Inches(14), Inches(1))
header_tf = header_box.text_frame
header_paragraph = header_tf.add_paragraph()
header_paragraph.text = header_text
header_paragraph.font.size = Pt(36)  # Large font size
header_paragraph.font.bold = True  # Ensure header is bold
header_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White
header_paragraph.alignment = 1  # Center the text

# Set background color for header text box
header_box.fill.solid()
header_box.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue

# Add title with dark blue background and centered text
title_text = "3 Factors that have affected the process of Economic Globalization"
title_box = slide.shapes.add_textbox(Inches(1), Inches(2.0), Inches(14), Inches(2))  # Increased height and adjusted top margin
tf = title_box.text_frame
p = tf.add_paragraph()
p.text = title_text
p.font.size = Pt(44)
p.font.bold = True  # Ensure title is bold
p.font.color.rgb = RGBColor(255, 255, 255)  # White

# Set background color for title text box
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue

# Center the title text
for paragraph in tf.paragraphs:
    paragraph.space_after = Pt(20)  # Add space after title
    paragraph.alignment = 1  # Center the text

# Remove the white background rectangle for the main content area
content_box_width = Inches(14)
content_box_height = Inches(5)
content_left = Inches(1)
content_top = Inches(3)

# Instead of adding a rectangle, directly add the content text
content_text = (
    "1. Improvements in transportation and communication technology have reduced "
    "the cost of transporting goods, services and factors of production and communicating "
    "economically useful knowledge and technology.\n\n"
    "2. Tastes of individuals and societies have generally but not universally favored "
    "taking advantage of the opportunities provided by declining."
)

content_tf = slide.shapes.add_textbox(content_left, content_top, content_box_width, content_box_height).text_frame
content_tf.word_wrap = True
content_p = content_tf.add_paragraph()
content_p.text = content_text
content_p.font.size = Pt(20)
content_p.font.color.rgb = RGBColor(255, 255, 255)  # White to ensure readability against blue background

# Adjust spacing above the content box
content_tf.space_before = Pt(30)  # Increase space before content

# Center the content text and adjust spacing between bullet points for consistency
for paragraph in content_tf.paragraphs:
    paragraph.alignment = 1  # Center the content text
    paragraph.space_after = Pt(20)  # Ensure consistent space after each bullet point

# Save the presentation
presentation.save("render.pptx")