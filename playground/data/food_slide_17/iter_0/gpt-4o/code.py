from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Initialize presentation
presentation = Presentation()

# Set slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Choose a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set up title and text box
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(3), Inches(1))
tf = title.text_frame
tf.text = "PART 2"

p = tf.paragraphs[0]
p.font.size = Pt(36)
p.font.bold = True

# Add exam logo
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), Inches(2), Inches(2))

# Add instruction text
instruction_box = slide.shapes.add_textbox(Inches(3), Inches(1.75), Inches(8), Inches(1))
tf = instruction_box.text_frame
tf.text = "Read the sentences and categorize each example of the past simple/past continuous."
p = tf.paragraphs[0]
p.font.size = Pt(18)
p.font.bold = True
p.alignment = PP_ALIGN.LEFT

# Add dialogue text
dialogue = (
    "JOEY: We were out to dinner. We were getting along, having a really nice time, "
    "I was thinking she was really cool and then, out of nowhere, "
    "(she reached over and took some of my fries from my plate!)\n"
    "PHOEBE: So she took some fries, big deal!\n"
    "RACHEL: Oh yeah, Joey doesn't share food. I mean, just last week, we were having breakfast, "
    "and...and he had a couple of grapes on his plate..."
)

dialogue_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(15), Inches(2))
tf = dialogue_box.text_frame
tf.text = dialogue
tf.word_wrap = True

for p in tf.paragraphs:
    p.font.size = Pt(16)

# Add table
rows, cols = 2, 2
table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(6.5), Inches(15), Inches(2)).table

# Set table column headings
table.cell(0, 0).text = "PAST SIMPLE"
table.cell(0, 1).text = "PAST CONTINUOUS"

# Align headers to center
for col in range(cols):
    table.cell(0, col).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Insert sentence into table cell
table.cell(1, 0).text = "We were out to dinner."

# Adjust font size in the table
for cell in table.iter_cells():
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.size = Pt(14)

# Save the presentation
presentation.save("render.pptx")