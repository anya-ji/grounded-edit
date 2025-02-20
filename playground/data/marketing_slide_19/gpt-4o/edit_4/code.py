from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object and set slide dimensions
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add new slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set dark background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(18, 18, 18)  # Dark color

# Add title text "THE MARKETING MIX"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "THE MARKETING MIX"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Add Marketing Mix Diagram image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0.5), Inches(1), Inches(7), Inches(6))

# Add the word "Marketing" on the right
marketing_box = slide.shapes.add_textbox(Inches(8.5), Inches(3), Inches(7), Inches(3))
marketing_frame = marketing_box.text_frame
marketing = marketing_frame.add_paragraph()
marketing.text = "Marketing"
marketing.font.size = Pt(60)
marketing.font.bold = True
marketing.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Add text boxes for PRODUCT, PRICE, PROMOTION, and PLACE with white color
product_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(2), Inches(1))
product_frame = product_box.text_frame
product = product_frame.add_paragraph()
product.text = "PRODUCT"
product.font.size = Pt(24)
product.font.bold = True
product.font.color.rgb = RGBColor(255, 255, 255)  # White text

price_box = slide.shapes.add_textbox(Inches(4), Inches(2), Inches(2), Inches(1))
price_frame = price_box.text_frame
price = price_frame.add_paragraph()
price.text = "PRICE"
price.font.size = Pt(24)
price.font.bold = True
price.font.color.rgb = RGBColor(255, 255, 255)  # White text

promotion_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(2), Inches(1))
promotion_frame = promotion_box.text_frame
promotion = promotion_frame.add_paragraph()
promotion.text = "PROMOTION"
promotion.font.size = Pt(24)
promotion.font.bold = True
promotion.font.color.rgb = RGBColor(255, 255, 255)  # White text

place_box = slide.shapes.add_textbox(Inches(4), Inches(4), Inches(2), Inches(1))
place_frame = place_box.text_frame
place = place_frame.add_paragraph()
place.text = "PLACE"
place.font.size = Pt(24)
place.font.bold = True
place.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Save the presentation
presentation.save("render.pptx")