from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object and set slide dimensions
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add new slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set gradient background
background = slide.background
fill = background.fill
fill.gradient()
stop1 = fill.gradient_stops[0]
stop1.position = 0
stop1.color.rgb = RGBColor(0, 100, 0)  # Dark green

stop2 = fill.gradient_stops[1]
stop2.position = 1
stop2.color.rgb = RGBColor(144, 238, 144)  # Light green

# Add title text "THE MARKETING MIX"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "THE MARKETING MIX"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Add Marketing Mix Diagram image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0.5), Inches(1), Inches(7), Inches(6))

# Add the word "Marketing" on the right
marketing_box = slide.shapes.add_textbox(Inches(8.5), Inches(3), Inches(7), Inches(3))
marketing_frame = marketing_box.text_frame
marketing = marketing_frame.add_paragraph()
marketing.text = "Marketing"
marketing.font.size = Pt(60)
marketing.font.bold = True
marketing.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Save the presentation
presentation.save("render.pptx")