from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide layout
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background for contrast

# Add title
title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title.text_frame
title_para = title_frame.add_paragraph()
title_para.text = "The Technology used in Blockchain"
title_para.font.size = Pt(56)  # Increased font size
title_para.font.bold = True

# Add bulleted list with checkboxes
left_inch = Inches(1)
top_inch = Inches(2)
width_inch = Inches(7)
height_inch = Inches(5)

content = slide.shapes.add_textbox(left_inch, top_inch, width_inch, height_inch)
text_frame = content.text_frame

items = [
    "☐ Cryptographic Keys",  # Unchecked
    "☑ Network Protocol",     # Changed to checked
    "☐ Distributed Ledger Technology",  # Changed to unchecked
    "☑ Hashing"
]

for item in items:
    p = text_frame.add_paragraph()
    p.text = item
    p.font.size = Pt(28)
    p.space_after = Pt(5)

# Add image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(9), Inches(2), Inches(6), Inches(4))

# Save the presentation
presentation.save("render.pptx")