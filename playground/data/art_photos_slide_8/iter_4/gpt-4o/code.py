from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import MSO_COLOR_TYPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from PIL import Image, ImageEnhance

# Function to adjust brightness and contrast of an image
def enhance_image(image_path):
    image = Image.open(image_path)
    # Increase brightness and contrast
    enhancer = ImageEnhance.Brightness(image)
    image = enhancer.enhance(1.2)  # Increase brightness by 20%
    enhancer = ImageEnhance.Contrast(image)
    image = enhancer.enhance(1.3)  # Increase contrast by 30%
    return image

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set gradient background to a lighter and warmer shade
slide_background = slide.background
fill = slide_background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 200)  # Changed to a lighter yellow hue
fill.gradient()  # Applying gradient

# Add the main title "Still Life"
title_box = slide.shapes.add_textbox(Inches(5.5), Inches(0.5), Inches(5), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Still Life"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(48)
title_paragraph.font.bold = True
title_paragraph.alignment = 1  # Centered

# Add the smaller title "The Subject and Content of Art"
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
subtitle_box.fill.solid()
subtitle_box.fill.fore_color.rgb = RGBColor(255, 255, 0)
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "The Subject and Content of Art"
subtitle_paragraph = subtitle_frame.paragraphs[0]
subtitle_paragraph.font.size = Pt(24)
subtitle_paragraph.font.bold = True

# Insert the flower painting image
image_path_flowers = "../../media/image_1.jpg"
enhanced_flowers = enhance_image(image_path_flowers)
enhanced_flowers.save("../../media/enhanced_flowers.jpg")  # Save the enhanced image
left_image = slide.shapes.add_picture("../../media/enhanced_flowers.jpg", Inches(1), Inches(2), width=Inches(7), height=Inches(5))
left_image.shadow.inherit = False
line = left_image.line
line.color.rgb = RGBColor(255, 255, 255)
line.width = Pt(5)

# Insert the fruit painting image
image_path_fruits = "../../media/image_2.jpg"
right_image = slide.shapes.add_picture(image_path_fruits, Inches(8), Inches(2), width=Inches(7), height=Inches(5))
right_image.shadow.inherit = False
line = right_image.line
line.color.rgb = RGBColor(255, 255, 255)
line.width = Pt(5)

# Save the presentation
presentation.save("render.pptx")