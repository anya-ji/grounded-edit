from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add the main title "Still Life"
title_box = slide.shapes.add_textbox(Inches(5.5), Inches(0.5), Inches(5), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Still Life"
title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(48)
title_paragraph.font.bold = True
title_paragraph.alignment = 1  # Centered

# Add the smaller title "The Subject and Content of Art"
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
subtitle_box.fill.solid()
subtitle_box.fill.fore_color.rgb = RGBColor(255, 255, 0)
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "The Subject and Content of Art"
subtitle_paragraph = subtitle_frame.paragraphs[0]
subtitle_paragraph.font.size = Pt(24)
subtitle_paragraph.font.bold = True

# Insert the flower painting image
image_path_flowers = "../../media/image_1.jpg"
left_image = slide.shapes.add_picture(image_path_flowers, Inches(1), Inches(2), width=Inches(7), height=Inches(5))
left_image.shadow.inherit = False
line = left_image.line
line.color.rgb = RGBColor(255, 255, 255)
line.width = Pt(5)

# Insert the fruit painting image
image_path_fruits = "../../media/image_2.jpg"
right_image = slide.shapes.add_picture(image_path_fruits, Inches(8), Inches(2), width=Inches(7), height=Inches(5))
right_image.shadow.inherit = False
line = right_image.line
line.color.rgb = RGBColor(255, 255, 255)
line.width = Pt(5)

# Save the presentation
presentation.save("render.pptx")