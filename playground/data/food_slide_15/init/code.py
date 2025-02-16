from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set the background color to dark blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0x00, 0x34, 0x81)  # dark blue

# Add the title text
title_text = "VIEWING FOLLOW-UP"
text_box = slide.shapes.add_textbox(Inches(3), Inches(3.5), Inches(10), Inches(2))
text_frame = text_box.text_frame
text_frame.word_wrap = True

title_paragraph = text_frame.add_paragraph()
title_paragraph.text = title_text
title_paragraph.font.size = Pt(60)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # white

# Add the light bulb icon with a brain
# For this example, let's assume we have an image named "lightbulb_brain.png"
icon_image_path = "lightbulb_brain.png"
slide.shapes.add_picture(icon_image_path, Inches(0.5), Inches(3), Inches(2.5), Inches(2.5))

# Add the "FLUENT" logo in the top right corner
# Let's assume the logo image is named "fluent_logo.png"
logo_image_path = "fluent_logo.png"
slide.shapes.add_picture(logo_image_path, Inches(12.5), Inches(0.5), height=Cm(2))

# Save the presentation to the specified path
presentation.save("render.pptx")