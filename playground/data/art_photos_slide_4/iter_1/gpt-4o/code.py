from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide layout
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set a solid background color (light blue for example)
background_color = slide.fill
background_color.solid()
background_color.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue

# Add the smaller title at the top-left corner
title_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
title_text_box.fill.solid()
title_text_box.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow background
title_text_frame = title_text_box.text_frame
title_text_frame.text = "The Subject and Content of Art"
title_text_frame.paragraphs[0].font.size = Pt(24)
title_text_frame.paragraphs[0].font.bold = True
title_text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add the main title
main_title_text_box = slide.shapes.add_textbox(Inches(5), Inches(7), Inches(11), Inches(1))
main_title_text_frame = main_title_text_box.text_frame
main_title_text_frame.text = "Subject of Art"
main_title_text_frame.paragraphs[0].font.size = Pt(40)
main_title_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add the main content
content_text_box = slide.shapes.add_textbox(Inches(5), Inches(7.5), Inches(12), Inches(1.5))
content_text_box.fill.solid()
content_text_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
content_text_frame = content_text_box.text_frame
content_text_frame.text = "A cat catching a bird is the subject of both works above."
content_text_frame.paragraphs[0].font.size = Pt(20)
content_text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add one image
image_path = "../../media/image_1.jpg"  # Only keeping one image
left = Inches(2)
top = Inches(2)
height = Inches(4)

slide.shapes.add_picture(image_path, left, top, height=height)

# Save the presentation
presentation.save("render.pptx")