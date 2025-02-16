from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls
from pptx.enum.shapes import MSO_SHAPE

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Create a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color to black
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# Add an image of Elon Musk to the left side
image_path = "elon_musk.jpg"  # Replace with the correct path to the image
image_left = Inches(0.5)
image_top = Inches(1.5)
image_width = Inches(4.0)
slide.shapes.add_picture(image_path, image_left, image_top, width=image_width)

# Add the quote text box
quote_text = (
    "The first step is to establish that something is possible; then probability will occur."
)
text_box = slide.shapes.add_textbox(Inches(5), Inches(1), Inches(10), Inches(2.5))
text_box.fill.solid()  # Optional: Make the text box opaque
text_frame = text_box.text_frame

# Add quote text
text_paragraph = text_frame.add_paragraph()
text_paragraph.text = quote_text
text_paragraph.font.size = Pt(28)
text_paragraph.font.color.rgb = RGBColor(255, 255, 255)
text_paragraph.alignment = PP_ALIGN.LEFT

# Apply yellow color to specific words
for run in text_paragraph.runs:
    if "first step" in run.text:
        run.font.color.rgb = RGBColor(255, 255, 0)
    if "something is possible" in run.text:
        run.font.color.rgb = RGBColor(255, 255, 0)
    if "occur" in run.text:
        run.font.color.rgb = RGBColor(255, 255, 0)

# Add author text
author_text = "Elon Musk"
author_paragraph = text_frame.add_paragraph()
author_paragraph.text = author_text
author_paragraph.font.size = Pt(24)
author_paragraph.font.color.rgb = RGBColor(255, 255, 255)
author_paragraph.alignment = PP_ALIGN.RIGHT

# Add "Thank You!" text on the right side
thank_you_box = slide.shapes.add_textbox(Inches(12.5), Inches(6), Inches(3), Inches(2))
thank_you_frame = thank_you_box.text_frame
thank_you_paragraph = thank_you_frame.add_paragraph()
thank_you_paragraph.text = "Thank You!"
thank_you_paragraph.font.size = Pt(36)
thank_you_paragraph.font.bold = True
thank_you_paragraph.font.color.rgb = RGBColor(0, 0, 0)
thank_you_paragraph.alignment = PP_ALIGN.CENTER

# Add two vertical yellow lines next to "Thank You!" text
for offset in [0.1, 0.3]:
    line = slide.shapes.add_line(
        Inches(12.2 + offset), Inches(6), Inches(12.2 + offset), Inches(8)
    )
    line.line.color.rgb = RGBColor(255, 255, 0)
    line.line.width = Pt(4)

# Add a white border around the slide
left = top = Inches(0)
width = Inches(16)
height = Inches(9)
border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
border.fill.background()  # No fill
border.line.color.rgb = RGBColor(255, 255, 255)
border.line.width = Pt(3)

presentation.save("render.pptx")