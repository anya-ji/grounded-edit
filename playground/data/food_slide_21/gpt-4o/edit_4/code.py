from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation and slide layout
presentation = Presentation()
slide_layout = presentation.slide_layouts[5]  # Blank slide
slide = presentation.slides.add_slide(slide_layout)

# Slide title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(12), Inches(0.8))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Pronunciation Activity"
title.font.size = Pt(36)
title.font.bold = True
title.alignment = PP_ALIGN.LEFT

# Exam logo
exam_logo_path = "../../media/image_1.jpg"
slide.shapes.add_picture(exam_logo_path, Inches(10), Inches(0.2), Inches(1), Inches(1))

# Introductory sentence
intro_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(12), Inches(0.5))
intro_frame = intro_box.text_frame
intro_sentence = intro_frame.add_paragraph()
intro_sentence.text = "Play the audio clips below. On which words do you notice the stress/intonation?"
intro_sentence.font.size = Pt(18)

# PART 1
part1_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(12), Inches(0.5))
part1_title_frame = part1_title_box.text_frame
part1_title = part1_title_frame.add_paragraph()
part1_title.text = "PART 1"
part1_title.font.size = Pt(24)
part1_title.font.bold = True

sentences = [
    "1. “We were out to dinner, ok… We were getting along…”",
    "2. “We were having a really nice time.”",
    "3. “I was thinking she was really cool.”",
    "4. “And then, out of nowhere…”"
]

for i, sentence in enumerate(sentences, start=1):
    y_position = 2 + (i - 1) * 1.5
    slide.shapes.add_picture("../../media/image_0.jpg", Inches(0.5), Inches(y_position), Inches(0.5), Inches(0.5))
    
    sentence_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_position), Inches(9), Inches(0.5))
    sentence_frame = sentence_box.text_frame
    sentence_paragraph = sentence_frame.add_paragraph()
    sentence_paragraph.text = sentence
    sentence_paragraph.font.size = Pt(16)
    
    response_box = slide.shapes.add_textbox(Inches(1.2), Inches(y_position + 0.6), Inches(10), Inches(0.5))
    response_frame = response_box.text_frame
    response_paragraph = response_frame.add_paragraph()
    response_paragraph.text = "Stressed word(s):"
    response_paragraph.font.size = Pt(16)
    response_paragraph.font.color.rgb = RGBColor(255, 0, 0)

# PART 2
part2_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(8.5), Inches(12), Inches(0.5))
part2_title_frame = part2_title_box.text_frame
part2_title = part2_title_frame.add_paragraph()
part2_title.text = "PART 2"
part2_title.font.size = Pt(24)
part2_title.font.bold = True

instruction_box2 = slide.shapes.add_textbox(Inches(0.5), Inches(9), Inches(12), Inches(1))
instruction_frame2 = instruction_box2.text_frame
instruction2 = instruction_frame2.add_paragraph()
instruction2.text = "Repeat each sentence above out loud. Put stress/intonation on the words you marked in Part 1."
instruction2.font.size = Pt(18)

# Save the presentation
presentation.save("render.pptx")