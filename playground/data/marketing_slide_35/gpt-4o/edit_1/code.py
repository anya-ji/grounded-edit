from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Create a blank slide layout
blank_slide_layout = presentation.slide_layouts[5]  # a blank layout
slide = presentation.slides.add_slide(blank_slide_layout)

# Add title with gradient background on the left
left_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(3), Inches(9))
left_box.fill.solid()
fill = left_box.fill
fill.gradient()
# Adding gradient stops from dark teal to lighter teal
stop1, stop2 = fill.gradient_stops[0], fill.gradient_stops[1]
stop1.position, stop1.color.rgb = 0.0, RGBColor(0, 51, 51)
stop2.position, stop2.color.rgb = 1.0, RGBColor(102, 153, 153)

title_frame = left_box.text_frame
title_frame.word_wrap = True

title_p = title_frame.add_paragraph()
title_p.text = "Types of Retailers"
title_p.font.bold = True
title_p.font.size = Pt(28)
title_p.font.color.rgb = RGBColor(255, 255, 255)
title_p.alignment = PP_ALIGN.CENTER

# Add top content box for "1. Mobile retailers"
top_box = slide.shapes.add_textbox(Inches(4), Inches(2), Inches(10), Inches(2))
top_box.fill.solid()
top_box.fill.fore_color.rgb = RGBColor(255, 223, 0)  # Changed from yellow to gold color
top_frame = top_box.text_frame
top_frame.word_wrap = True
top_p = top_frame.add_paragraph()
top_p.text = "1. Mobile retailers"
top_p.font.size = Pt(24)
top_p.alignment = PP_ALIGN.CENTER

# Add arrow between the two boxes
arrow = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(9), Inches(4), Inches(1), Inches(1))

# Add bottom content box for "2. Fixed – Shop Retailers"
bottom_box = slide.shapes.add_textbox(Inches(4), Inches(6), Inches(10), Inches(2))
bottom_box.fill.solid()
bottom_box.fill.fore_color.rgb = RGBColor(255, 140, 0)  # Orange color
bottom_frame = bottom_box.text_frame
bottom_frame.word_wrap = True
bottom_p = bottom_frame.add_paragraph()
bottom_p.text = "2. Fixed – Shop Retailers"
bottom_p.font.size = Pt(24)
bottom_p.alignment = PP_ALIGN.CENTER

# Save the presentation
presentation.save("render.pptx")