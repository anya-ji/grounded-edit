from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide
slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Set the slide background to black
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# Add the title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True

p = title_frame.add_paragraph()
p.text = "Now: I, CAN"
p.font.bold = True
p.font.size = Pt(44)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Add the first bullet point with border
top_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(1.5))
top_box.fill.solid()
top_box.fill.fore_color.rgb = RGBColor(255, 165, 0)  # Orange
top_box.line.color.rgb = RGBColor(211, 211, 211)  # Light grey border
top_frame = top_box.text_frame
top_frame.word_wrap = True
top_frame.text = "Differentiate between different types of market."
top_frame.paragraphs[0].font.size = Pt(24)
top_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

# Add the second bullet point with border
bottom_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(14), Inches(1.5))
bottom_box.fill.solid()
bottom_box.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow
bottom_box.line.color.rgb = RGBColor(211, 211, 211)  # Light grey border
bottom_frame = bottom_box.text_frame
bottom_frame.word_wrap = True
bottom_frame.text = "Explain why society need different types of market."
bottom_frame.paragraphs[0].font.size = Pt(24)
bottom_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

# Add the image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(16), height=Inches(9))

# Save the presentation
presentation.save("render.pptx")