from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Define title text
title_text = "Types of graphic design"

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_frame.clear()
title_p = title_frame.paragraphs[0]
title_p.text = title_text
title_p.font.size = Pt(32)
title_p.font.bold = True

# Define text box properties
text_boxes_data = [
    ("Corporate plan", "Corporate plan has to do with the visual personality of an organization."),
    ("Marketing and advertising design", "Probably one of the most commonly known kinds of graphic design."),
    ("Website composition", "While website composition isn't really a kind of visual communication, visual communication is one component of website composition."),
    ("For more professional assistance", "Go through the site graphic design company in Dubai, seo company Dubai.")
]

# Position settings
left_margin = Inches(0.5)
top_margin = Inches(2)
box_width = Inches(15)
box_height = Inches(1.2)
vertical_spacing = Inches(1.3)

# Add text boxes with a deeper blue background
for i, (title, description) in enumerate(text_boxes_data):
    top_position = top_margin + i * vertical_spacing
    textbox = slide.shapes.add_textbox(left_margin, top_position, box_width, box_height)

    text_frame = textbox.text_frame
    text_frame.clear()

    # Title paragraph
    p_title = text_frame.add_paragraph()
    p_title.text = title
    p_title.font.size = Pt(20)
    p_title.font.bold = True
    p_title.font.color.rgb = RGBColor(255, 255, 255)  # White color

    # Description paragraph
    p_desc = text_frame.add_paragraph()
    p_desc.text = description
    p_desc.font.size = Pt(16)
    p_desc.font.color.rgb = RGBColor(255, 255, 255)  # White color

    # Set background color to a deeper blue for all sections
    textbox.fill.solid()
    textbox.fill.fore_color.rgb = RGBColor(0, 102, 204)  # Deeper blue color

# Adjust height of "Website composition" text box to match "Marketing and advertising design"
website_box = slide.shapes[-2]  # This is the "Website composition" text box
website_box.height = Inches(1.2)  # Set height to match the other boxes

# Save the presentation
presentation.save("render.pptx")