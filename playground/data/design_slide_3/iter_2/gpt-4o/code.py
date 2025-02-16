from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Insert background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.2), Inches(14), Inches(1.5))  # Adjusted positioning
title_frame = title_box.text_frame
title_paragraph = title_frame.add_paragraph()
title_paragraph.text = "Components of visual communication"
title_paragraph.font.bold = True
title_paragraph.font.size = Pt(30)  # Reduced font size
title_paragraph.font.color.rgb = RGBColor(255, 255, 0)  # bright yellow color

# Add content
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(5.5))  # Adjusted positioning
content_frame = content_box.text_frame

# Bullet point 1
bullet1 = content_frame.add_paragraph()
bullet1.text = (
    "Components of craftsmanship are the essential units of any visual plan "
    "that structure the plan's design and pass on its visual messages."
)
bullet1.font.size = Pt(20)  # Reduced font size
bullet1.font.color.rgb = RGBColor(255, 255, 0)  # bright yellow color

# Bullet point 2
bullet2 = content_frame.add_paragraph()
bullet2.text = (
    "The components of visual computerization are line, shape, size, color, "
    "Typography, Space and Texture."
)
bullet2.font.size = Pt(20)  # Reduced font size
bullet2.font.color.rgb = RGBColor(255, 255, 0)  # bright yellow color

presentation.save("render.pptx")