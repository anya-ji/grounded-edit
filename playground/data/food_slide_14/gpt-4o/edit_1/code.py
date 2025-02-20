from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()

# Set slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "PART 4"
title.font.bold = True
title.font.size = Pt(44)
title.font.color.rgb = RGBColor(0, 51, 102)  # dark blue

# Add instructions and events
content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

instruction = content_frame.add_paragraph()
instruction.text = "List the events in Joey’s second date with Sarah in order from 2 - 6. The first one has been done as an example."
instruction.font.size = Pt(24)
instruction.space_after = Pt(14)

events = [
    "1. Joey apologizes for getting angry that Sarah took his food on the first date.",
    "2. The waiter arrives with their desserts.",
    "3. Joey doesn’t like the dessert he orders and tries to order a different one.",
    "4. Sarah leaves the table for a phone call.",
    "5. Joey is unable to control himself and eats all of Sarah’s dessert.",
    "6. Joey refuses to apologize when Sarah returns."
]

# Rearranged order: 2, 3, 4, 5, 6
rearranged_events = [
    events[1],
    events[2],
    events[3],
    events[4],
    events[5]
]

for event_text in rearranged_events:
    event = content_frame.add_paragraph()
    event.text = event_text
    event.font.size = Pt(20)

# Add image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(11), Inches(6), width=Inches(4))

# Add button
button_box = slide.shapes.add_textbox(Inches(12.8), Inches(7.5), Inches(3), Inches(0.7))
button_frame = button_box.text_frame
button_frame.word_wrap = True
button = button_frame.add_paragraph()
button.text = "VIEWING ACTIVITY"
button.font.size = Pt(20)
button.font.color.rgb = RGBColor(255, 255, 255)  # white
button_box.fill.solid()
button_box.fill.fore_color.rgb = RGBColor(0, 51, 102)  # dark blue
button_box.line.color.rgb = RGBColor(0, 51, 102)

# Save presentation
presentation.save("render.pptx")