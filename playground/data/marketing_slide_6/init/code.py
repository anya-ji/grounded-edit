from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide
slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Define colors
red_color = RGBColor(255, 0, 0)
gray_color = RGBColor(169, 169, 169)
white_color = RGBColor(255, 255, 255)

# Add red background
background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, presentation.slide_width, presentation.slide_height)
background.fill.solid()
background.fill.fore_color.rgb = red_color
background.line.fill.background()

# Add gray area horizontally in the middle
gray_area = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0),
    Inches(3),
    presentation.slide_width,
    Inches(3)
)
gray_area.fill.solid()
gray_area.fill.fore_color.rgb = gray_color
gray_area.line.fill.solid()
gray_area.line.fill.fore_color.rgb = white_color
gray_area.line.width = Pt(1.5)

# Add a circular shape on the left side with the image inside
circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1), Inches(1.5), Inches(4), Inches(4))
circle.fill.user_picture("../../media/image_0.jpg")
circle.line.fill.solid()
circle.line.fill.fore_color.rgb = white_color

# Add the title text
text_box = slide.shapes.add_textbox(Inches(11), Inches(3.5), Inches(4), Inches(2))
tf = text_box.text_frame
p = tf.add_paragraph()
p.text = "WHAT IS MARKET?"
p.font.bold = True
p.font.size = Pt(44)
p.font.color.rgb = white_color

# Save the presentation
presentation.save("render.pptx")