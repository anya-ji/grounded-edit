from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(16), height=Inches(9))

# Add question text
text_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(3))
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = "Q. How often you think about social impact of your action before you act?"
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(255, 255, 255)
text_box.text_frame.word_wrap = True
text_box.text_frame.paragraphs[0].alignment = 1  # Center

# Remove the gray background by ensuring no fill for the textbox
text_box.fill.solid()   # Make sure this is present to define fill properties
text_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Set fill to white (or remove this entirely)

# Save presentation
presentation.save("render.pptx")