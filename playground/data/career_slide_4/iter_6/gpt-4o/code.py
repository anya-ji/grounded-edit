from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide
slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set dark background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(30, 30, 30)  # dark background

# Add title
title_text = "KEY REASONS WHY ACADEMIC SUCCESS IS IMPORTANT IN SOCIETY"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1.0))
text_frame = title_box.text_frame
text_frame.word_wrap = True
p = text_frame.add_paragraph()
p.text = title_text
p.font.size = Pt(24)  # Reduced font size for title to 1.5em (approximately 24 points)
p.font.bold = True
p.alignment = PP_ALIGN.CENTER
p.font.color.rgb = RGBColor(255, 255, 255)  # white color for title text

# Section 1
left_title = "Personal Growth and Development"
left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5), Inches(1.0))
left_frame = left_box.text_frame
left_frame.word_wrap = True
left_p = left_frame.add_paragraph()
left_p.text = left_title
left_p.font.size = Pt(32)  # Increased font size for section title
left_p.font.color.rgb = RGBColor(255, 255, 0)  # yellow color
left_p.alignment = PP_ALIGN.LEFT  # Adjust alignment

left_bullets = [
    "Academic success promotes personal growth through critical thinking, problem-solving, and intellectual curiosity.",
    "It leads to fulfilling careers and a sense of purpose."
]
for bullet in left_bullets:
    p = left_frame.add_paragraph()
    p.text = bullet
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)  # white color

# Add image to left section
left_image_path = "../../media/image_0.jpg"  # replace with your own image path
slide.shapes.add_picture(left_image_path, Inches(0.5), Inches(3.0), width=Inches(4))

# Section 2
middle_title = "Employability and Career Opportunities"
middle_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(5), Inches(1.0))
middle_frame = middle_box.text_frame
middle_frame.word_wrap = True
middle_p = middle_frame.add_paragraph()
middle_p.text = middle_title
middle_p.font.size = Pt(32)  # Increased font size for section title
middle_p.font.color.rgb = RGBColor(255, 255, 0)  # yellow color
middle_p.alignment = PP_ALIGN.LEFT  # Adjust alignment

middle_bullets = [
    "Academic foundation is key for many careers.",
    "Achievements can lead to better jobs, earnings, and advancement."
]
for bullet in middle_bullets:
    p = middle_frame.add_paragraph()
    p.text = bullet
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(255, 255, 255)  # white color

# Add image to middle section
middle_image_path = "../../media/image_0.jpg"  # replace with your own image path
slide.shapes.add_picture(middle_image_path, Inches(5.5), Inches(4.0), width=Inches(4))

# Section 3
right_title = "Economic Impact"
right_box = slide.shapes.add_textbox(Inches(10.5), Inches(1.5), Inches(5), Inches(1.0))
right_frame = right_box.text_frame
right_frame.word_wrap = True
right_p = right_frame.add_paragraph()
right_p.text = right_title
right_p.font.size = Pt(32)  # Increased font size for section title
right_p.font.color.rgb = RGBColor(255, 255, 0)  # yellow color
right_p.alignment = PP_ALIGN.LEFT  # Adjust alignment

right_text = (
    "Well-educated societies lead to stronger economies by creating skilled "
    "workers and reducing unemployment."
)
p = right_frame.add_paragraph()
p.text = right_text
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(255, 255, 255)  # white color

# Add image to right section
right_image_path = "../../media/image_1.jpg"  # replace with your own image path
slide.shapes.add_picture(right_image_path, Inches(10.5), Inches(4.0), width=Inches(4))

# Save the presentation
presentation.save("render.pptx")