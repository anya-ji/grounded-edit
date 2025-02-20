from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide layout
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue background color

# Title
title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title.text_frame
title_para = title_frame.add_paragraph()
title_para.text = "What do you think is the general rule for this pronunciation scheme?"
title_para.font.size = Pt(32)
title_para.font.bold = True

# Subtitle
subtitle = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(0.5))
subtitle_frame = subtitle.text_frame
subtitle_para = subtitle_frame.add_paragraph()
subtitle_para.text = "PART 3"
subtitle_para.font.size = Pt(24)
subtitle_para.font.bold = True

# Options
options = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(3))
options_frame = options.text_frame
options_frame.word_wrap = True

options_a = options_frame.add_paragraph()
options_a.text = "a) We add stress/intonation on the time or location we’re speaking about."
options_a.font.size = Pt(20)

options_b = options_frame.add_paragraph()
options_b.text = "b) We add stress/intonation on the present participle (-ing verb)."
options_b.font.size = Pt(20)

options_c = options_frame.add_paragraph()
options_c.text = "c) We add stress/intonation on the last word of the clause when retelling events."
options_c.font.size = Pt(20)

# Insert checkmark image
checkmark_path = "../../media/image_0.jpg"
slide.shapes.add_picture(checkmark_path, Inches(0.7), Inches(4.8), Inches(0.5), Inches(0.5))

# Add footer or sidebar
footer = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(8.5), Inches(16), Inches(0.5))
footer.fill.solid()
footer.fill.fore_color.rgb = RGBColor(200, 200, 250)

footer_text = footer.text_frame.add_paragraph()
footer_text.text = "PRONUNCIATION ACTIVITY"
footer_text.font.size = Pt(16)
footer_text.font.bold = True
footer_text.font.color.rgb = RGBColor(0, 0, 0)

# Save the presentation
presentation.save("render.pptx")