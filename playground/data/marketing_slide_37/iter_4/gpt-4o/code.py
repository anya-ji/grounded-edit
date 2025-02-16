from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5] 
slide = presentation.slides.add_slide(slide_layout)

# Set gradient background color to a smooth teal gradient
background = slide.background
fill = background.fill
fill.gradient()  # Change from solid to gradient
fill.gradient_stops[0].color.rgb = RGBColor(173, 255, 255)  # Light cyan at the top
fill.gradient_stops[1].color.rgb = RGBColor(0, 139, 139)  # Dark teal at the bottom

# Add blue graphic of a person with an arrow pointing right
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0.5), Inches(3), width=Inches(3), height=Inches(3), transform=None)

# Add title text on the top right
title_box = slide.shapes.add_textbox(Inches(5), Inches(0.5), Inches(10), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Think Line:"
title.font.size = Pt(36)
title.font.bold = True
title.font.color.rgb = RGBColor(255, 255, 255)  # White

# Add main text below the title
text_box = slide.shapes.add_textbox(Inches(5), Inches(1.5), Inches(10), Inches(2))
tf = text_box.text_frame
text = tf.add_paragraph()
text.text = "DO YOU THINK MARKET CREATE OPPORTUNITY. HOW?"
text.font.size = Pt(36)
text.font.color.rgb = RGBColor(173, 255, 255)  # Light cyan

# Add a small red rectangle at the top right corner for emphasis
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(14.5), Inches(0.2), Inches(1), Inches(0.5))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red

# Save the presentation
presentation.save("render.pptx")