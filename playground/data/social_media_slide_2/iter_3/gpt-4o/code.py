from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color for the title area
background = slide.shapes.add_shape(
    1, Inches(0), Inches(0), Inches(16), Inches(1.5)
)
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue

# Add title text
title_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(8), Inches(1))
title = title_box.text_frame.add_paragraph()
title.text = "DEFINITION OF TERMS"
title.font.bold = True
title.font.size = Pt(32)
title.font.color.rgb = RGBColor(255, 255, 255)  # White

# Adding blue background behind definitions
definitions_bg = slide.shapes.add_shape(
    1, Inches(0.5), Inches(1.5), Inches(15), Inches(6.5)  # Positioned according to the content
)
definitions_bg.fill.solid()
definitions_bg.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue

# Add content for Media
media_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(15), Inches(4))  # Increased height here
media_frame = media_box.text_frame
media_title = media_frame.add_paragraph()
media_title.text = "MEDIA"
media_title.font.bold = True
media_title.font.size = Pt(26)  # Increased font size for MEDIA
media_title.font.color.rgb = RGBColor(255, 255, 255)  # White text for contrast

media_text = media_frame.add_paragraph()
media_text.text = "refers to the communication channels through which we disseminate news, music, movies, education, promotional messages and other data."
media_text.font.size = Pt(20)
media_text.space_after = Pt(10)

# Add content for Culture
culture_box = slide.shapes.add_textbox(Inches(0.5), Inches(6), Inches(15), Inches(3))  # Increased height here
culture_frame = culture_box.text_frame
culture_title = culture_frame.add_paragraph()
culture_title.text = "CULTURE"
culture_title.font.bold = True
culture_title.font.size = Pt(24)  # Font size remains the same for CULTURE
culture_title.font.color.rgb = RGBColor(255, 255, 255)  # White text for contrast

culture_text = culture_frame.add_paragraph()
culture_text.text = "can be defined as all the ways of life, including arts, beliefs and institutions of a population that are passed down from generation to generation."
culture_text.font.size = Pt(20)

# Save the presentation
presentation.save("render.pptx")