from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add title
txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1))
tf = txBox.text_frame
p = tf.add_paragraph()
p.text = "Information Systems"
p.font.bold = True
p.font.size = Pt(36)
p.font.name = 'Arial'
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.LEFT

# Set the background for title
fill = txBox.fill
fill.solid()
fill.fore_color.rgb = RGBColor(169, 169, 169)  # Gray color

# Add logo image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0.5), Inches(2.5), Inches(2), Inches(2))

# Save the presentation
presentation.save("render.pptx")