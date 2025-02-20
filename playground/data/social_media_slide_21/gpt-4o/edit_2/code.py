from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object with 16x9 slide dimensions
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with blank layout
slide_layout = presentation.slide_layouts[5]  # This is a blank layout
slide = presentation.slides.add_slide(slide_layout)

# Title text
text = "THANK YOU!"

# Create a text box for the main message, centered vertically and horizontally
text_box = slide.shapes.add_textbox(Inches(4), Inches(3), Inches(8), Inches(2))  # Changed to 3 for vertical centering
text_frame = text_box.text_frame
text_frame.text = text
text_frame.paragraphs[0].font.bold = True
text_frame.paragraphs[0].font.size = Pt(80)
text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102) # Dark Blue
text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Save the slide to the specified path
presentation.save("render.pptx")