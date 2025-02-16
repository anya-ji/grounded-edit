from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation and slide
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set background color to light gray
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(211, 211, 211)  # Light gray color

# Add PART 1 title with rectangular shape and border
part1_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.2), Inches(1.5), Inches(0.5))
part1_box.fill.solid()
part1_box.fill.fore_color.rgb = RGBColor(0, 0, 255)  # Blue color
part1_box.line.color.rgb = RGBColor(0, 0, 0)  # Black border
part1_text = part1_box.text_frame.add_paragraph()
part1_text.text = "PART 1"
part1_text.font.bold = True
part1_text.font.size = Pt(24)  # Adjust font size to match target image
part1_text.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Add instruction text for PART 1
instruction_text = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(10), Inches(1))
instruction_p = instruction_text.text_frame.add_paragraph()
instruction_p.text = "Discuss/Write short answers below according to the information."
instruction_p.font.size = Pt(18)

# Add questions and time indicators
question1_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(13), Inches(1))
question1_p = question1_text.text_frame.add_paragraph()
question1_p.text = "1: What happened on Joey’s first date with Sarah?                                         0:18"
question1_p.font.size = Pt(16)

question2_text = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(13), Inches(1))
question2_p = question2_text.text_frame.add_paragraph()
question2_p.text = "2: How does Joey feel about their first date? What does he say about it?           0:48"
question2_p.font.size = Pt(16)

# Add image of Joey
joey_image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(joey_image_path, Inches(13.5), Inches(1.0), width=Inches(1.5))

# Add PART 2 title with rectangular shape and border
part2_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(3.0), Inches(1.5), Inches(0.5))
part2_box.fill.solid()
part2_box.fill.fore_color.rgb = RGBColor(0, 0, 255)  # Blue color
part2_box.line.color.rgb = RGBColor(0, 0, 0)  # Black border
part2_text = part2_box.text_frame.add_paragraph()
part2_text.text = "PART 2"
part2_text.font.bold = True
part2_text.font.size = Pt(24)  # Adjust font size to match target image
part2_text.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Add instruction for PART 2
instruction2_text = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(10), Inches(1))
instruction2_p = instruction2_text.text_frame.add_paragraph()
instruction2_p.text = "Refer back to Preview, Part 3. How similar was your date story to Joey’s?"
instruction2_p.font.size = Pt(18)

# Add dialogue subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(5.5), Inches(15), Inches(1))
subtitle_p = subtitle_box.text_frame.add_paragraph()
subtitle_p.text = "Joey: How you doin'?"
subtitle_p.font.size = Pt(18)
subtitle_p.font.color.rgb = RGBColor(0, 0, 0)  # Black text, typical for subtitles

# Add VIEWING ACTIVITY button
button = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(7.0), Inches(3), Inches(0.75))
button.fill.solid()
button.fill.fore_color.rgb = RGBColor(255, 165, 0)  # Orange color
button_text = button.text_frame.add_paragraph()
button_text.text = "VIEWING ACTIVITY"
button_text.font.bold = True
button_text.font.size = Pt(18)
button_text.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Increase the saturation of colors (by adjusting RGB values)
part1_box.fill.fore_color.rgb = RGBColor(0, 0, 255)  # Blue (unchanged)
part2_box.fill.fore_color.rgb = RGBColor(0, 0, 255)  # Blue (unchanged)
button.fill.fore_color.rgb = RGBColor(255, 140, 0)  # More vibrant orange

# Save presentation
presentation.save("render.pptx")