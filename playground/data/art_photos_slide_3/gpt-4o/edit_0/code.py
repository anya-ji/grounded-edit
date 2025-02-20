from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide
slide_layout = presentation.slide_layouts[5]  # Use a blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add a soft gradient background
fill = slide.background.fill
fill.gradient()
stop_1 = fill.gradient_stops[0]
stop_1.position = 0.0
stop_1.color.rgb = RGBColor(255, 223, 186)  # Light peach color
stop_2 = fill.gradient_stops[1]
stop_2.position = 1.0
stop_2.color.rgb = RGBColor(255, 204, 204)  # Light pink color

# Add the background image
img_path = "../../media/image_0.jpg"
slide.shapes.add_picture(img_path, 0, 0, presentation.slide_width, presentation.slide_height)

# Add a smaller title at the top-left corner
small_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(6), Inches(1))
small_title_frame = small_title_box.text_frame
small_title = small_title_frame.add_paragraph()
small_title.text = "The Subject and Content of Art."
small_title.font.size = Pt(18)
small_title.font.bold = True
small_title.font.color.rgb = RGBColor(0, 0, 0)
small_title_box.fill.solid()
small_title_box.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow background

# Add the main title
title_box = slide.shapes.add_textbox(Inches(3), Inches(1.5), Inches(10), Inches(1.5))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "The Subject of Art"
title.font.size = Pt(32)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow background

# Add bullets
content_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(12), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

bullet_1 = content_frame.add_paragraph()
bullet_1.text = "The subject is the visual focus or the image that may be extracted from examining the artworks."
bullet_1.font.size = Pt(24)
bullet_1.font.color.rgb = RGBColor(0, 0, 0)

bullet_2 = content_frame.add_paragraph()
bullet_2.text = "An artwork can be identified by its subject."
bullet_2.font.size = Pt(24)
bullet_2.font.color.rgb = RGBColor(0, 0, 0)

# Set background color for content
content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

# Save the presentation
presentation.save("render.pptx")