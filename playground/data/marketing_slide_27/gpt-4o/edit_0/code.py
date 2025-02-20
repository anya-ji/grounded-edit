from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color to dark gray
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(64, 64, 64)  # Dark gray

# Define the images
image_paths = [
    "../../media/image_1.jpg",
    "../../media/image_0.jpg",
    "../../media/image_3.jpg",
    "../../media/image_4.jpg",
    "../../media/image_2.jpg"
]

# Define rectangle positions and sizes
rectangle_positions = [
    (Inches(2), Inches(1.5)),
    (Inches(5), Inches(1.5)),
    (Inches(8), Inches(1.5)),
    (Inches(11), Inches(1.5)),
    (Inches(6.5), Inches(5))
]
rectangle_size = Inches(2.5)
corner_radius = Pt(0.5)  # Radius for rounded corners 

# Add rectangular image frames with rounded corners
for image_path, position in zip(image_paths, rectangle_positions):
    left, top = position
    pic = slide.shapes.add_picture(image_path, left, top, rectangle_size, rectangle_size)
    
    # Create a rectangular shape with rounded corners
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, rectangle_size, rectangle_size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
    shape.line.fill.background()

# Add title text
title_text = "Types of Markets"
text_box = slide.shapes.add_textbox(Inches(4.5), Inches(7.5), Inches(7), Inches(1.5))
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.add_paragraph()
p.text = title_text
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)  # White color for contrast

# Save the presentation
presentation.save("render.pptx")