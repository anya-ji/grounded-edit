from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color to dark gray
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(64, 64, 64)  # Dark gray

# Define the images
image_paths = [
    "../../media/image_1.jpg",
    "../../media/image_0.jpg",
    "../../media/image_3.jpg",
    "../../media/image_4.jpg",
    "../../media/image_2.jpg"
]

# Define circle positions and sizes
circle_positions = [
    (Inches(2), Inches(1.5)),
    (Inches(5), Inches(1.5)),
    (Inches(8), Inches(1.5)),
    (Inches(11), Inches(1.5)),
    (Inches(6.5), Inches(5))
]
circle_size = Inches(2.5)

# Add circular image frames
for image_path, position in zip(image_paths, circle_positions):
    left, top = position
    # Add the picture and apply the circular mask
    pic = slide.shapes.add_picture(image_path, left, top, circle_size, circle_size)
    
    # Create a transparent circle mask
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, circle_size, circle_size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Set to white instead of black
    shape.fill.background()  # Make the shape transparent
    shape.line.fill.background()

# Add title text
title_text = "Types of Markets"
text_box = slide.shapes.add_textbox(Inches(4.5), Inches(7.5), Inches(7), Inches(1.5))
text_frame = text_box.text_frame
text_frame.word_wrap = True

p = text_frame.add_paragraph()
p.text = title_text
p.font.size = Pt(24)  # Reduced font size for better visual balance
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)  # White color for contrast

# Save the presentation
presentation.save("render.pptx")