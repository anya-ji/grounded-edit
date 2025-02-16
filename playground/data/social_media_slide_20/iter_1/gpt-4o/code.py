from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a PowerPoint presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]  # choosing a blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set dark blue background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 32, 96)  # Adjusted to a new shade of blue

# Add a title text box
title_text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1))
title_frame = title_text_box.text_frame
title_text = title_frame.add_paragraph()
title_text.text = "Three perspectives on global cultural flows:"
title_text.font.size = Pt(36)
title_text.font.color.rgb = RGBColor(255, 255, 255)  # White
title_text.font.bold = True
title_frame.word_wrap = True

# Add a subtitle text box for "3. CULTURAL CONVERGENCE"
subtitle_text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(1))
subtitle_frame = subtitle_text_box.text_frame
subtitle_text = subtitle_frame.add_paragraph()
subtitle_text.text = "3. CULTURAL CONVERGENCE"
subtitle_text.font.size = Pt(40)
subtitle_text.font.color.rgb = RGBColor(255, 255, 255)  # White
subtitle_text.font.bold = True
subtitle_frame.word_wrap = True

# Add content text box for bullet points
content_text_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.5), Inches(13), Inches(4))  # Adjusted from 4.5 to 3.5 for spacing
content_frame = content_text_box.text_frame
content_frame.word_wrap = True

# Bullet point 1
bullet_1 = content_frame.add_paragraph()
bullet_1.text = "approach stresses homogeneity introduced by globalization."
bullet_1.font.size = Pt(24)
bullet_1.font.color.rgb = RGBColor(255, 255, 255)  # White

# Bullet point 2
bullet_2 = content_frame.add_paragraph()
bullet_2.text = ("Cultures are deemed to be radically altered by strong flows, "
                 "while cultural imperialism happens when one culture imposes itself on "
                 "and tends to destroy at least parts of another culture.")
bullet_2.font.size = Pt(24)
bullet_2.font.color.rgb = RGBColor(255, 255, 255)  # White

# Save the presentation
presentation.save("render.pptx")