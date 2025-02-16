from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), width=Inches(16), height=Inches(9))

# Title text
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title = title_frame.add_paragraph()
title.text = "Q2. HOTS"
title.font.size = Pt(40)
title.font.bold = True
title.font.color.rgb = RGBColor(255, 255, 255)  # White color

# Main text
main_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(4))
main_frame = main_box.text_frame
main_frame.word_wrap = True
main_text = main_frame.add_paragraph()
main_text.text = "Consider yourself as a producer, list the various things you must be taking care of till your product reaches the consumers. (3 MARKS)"
main_text.font.size = Pt(24)
main_text.font.bold = True
main_text.font.color.rgb = RGBColor(100, 100, 100)  # Darker color for better contrast

presentation.save("render.pptx")