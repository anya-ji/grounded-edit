from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Background gradient
background = slide.shapes.add_shape(
    1,  # Use the rectangle shape
    0, 0, Inches(16), Inches(9)
)
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 128)  # Dark Blue
fill.gradient()
stop = fill.gradient_stops
stop[0].position = 0.0
stop[0].color.rgb = RGBColor(0, 0, 128)  # Dark Blue
stop[1].position = 1.0
stop[1].color.rgb = RGBColor(0, 0, 255)  # Lighter Blue

# Title text
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title = title_frame.add_paragraph()
title.text = "Q2. HOTS"
title.font.size = Pt(40)
title.font.bold = True
title.font.color.rgb = RGBColor(255, 255, 255)

# Main text
main_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(4))
main_frame = main_box.text_frame
main_frame.word_wrap = True
main_text = main_frame.add_paragraph()
main_text.text = "Consider yourself as a producer, list the various things you must be taking care of till your product reaches the consumers. (3 MARKS)"
main_text.font.size = Pt(24)
main_text.font.bold = True
main_text.font.color.rgb = RGBColor(255, 255, 255)

presentation.save("render.pptx")