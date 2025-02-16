from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set dark blue background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue background

# Add a light bulb icon on the left side
left_icon = slide.shapes.add_shape(MSO_SHAPE.LIGHTBULB, Inches(0.5), Inches(2.5), Inches(1.5), Inches(2.5))
light_bulb_fill = left_icon.fill
light_bulb_fill.solid()
light_bulb_fill.fore_color.rgb = RGBColor(255, 255, 255)  # White color for icon

# Add centered title text
title_box = slide.shapes.add_textbox(Inches(3), Inches(3.5), Inches(10), Inches(1.5))
title_text_frame = title_box.text_frame
title_text_frame.text = "‘DATE GONE WRONG’ STORY"
title_paragraph = title_text_frame.paragraphs[0]
title_paragraph.font.bold = True
title_paragraph.font.size = Pt(44)
title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White font
title_text_frame.text_anchor = MSO_SHAPE.CENTER

# Add "FLUENT" logo in the top right corner
# Assuming 'fluent_logo.png' is available in the directory
fluent_logo_path = "fluent_logo.png"
slide.shapes.add_picture(fluent_logo_path, Inches(12), Inches(0.5), Inches(3), Inches(1))

# Save the presentation
presentation.save("render.pptx")