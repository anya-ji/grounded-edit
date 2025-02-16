from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation object
presentation = Presentation()

# Set dimensions for widescreen
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "PART 3"
title.font.size = Pt(40)
title.font.bold = True

# Instruction text
instruction_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(1))
instruction_frame = instruction_box.text_frame
instruction = instruction_frame.add_paragraph()
instruction.text = "Fill in the correct past simple or past continuous forms in the short story below."
instruction.font.size = Pt(20)

# Text with blank spaces
text_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10), Inches(4))
text_frame = text_box.text_frame
text = text_frame.add_paragraph()
text.text = ("Last week, I went (I / go) on a date with a nice boy (I / meet) recently. "
             "(We / skate) in the park, (we / laugh), and (he / make) some really funny jokes, "
             "and then, out of nowhere (I / fall) down (I / break) my leg!")
text.font.size = Pt(18)

# Adding text boxes for answers
verbs = ["I / go", "I / meet", "We / skate", "we / laugh", "he / make", "I / fall", "I / break"]
positions = [
    (Inches(3.5), Inches(3)), 
    (Inches(7), Inches(3)), 
    (Inches(2.6), Inches(3.5)), 
    (Inches(6.8), Inches(3.5)), 
    (Inches(9.3), Inches(3.5)), 
    (Inches(3.6), Inches(4)), 
    (Inches(5.6), Inches(4))
]

# Add a text box in front of each verb in parentheses
for pos in positions:
    answer_box = slide.shapes.add_textbox(pos[0], pos[1], Inches(1.5), Inches(0.5))
    frame = answer_box.text_frame
    answer = frame.add_paragraph()
    frame.word_wrap = True
    answer_box.fill.solid()
    answer_box.fill.fore_color.rgb = RGBColor(230, 230, 230)  # Light background color
    answer_box.fill.transparency = 1.0  # Set transparency to 100% (0% opacity)

# Insert exam icon
icon_path = "../../media/image_1.jpg"
slide.shapes.add_picture(icon_path, Inches(12), Inches(0.5), Inches(2), Inches(2))

# Insert image of people ice skating
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(11), Inches(5), Inches(4), Inches(3))

# Save presentation
presentation.save("render.pptx")