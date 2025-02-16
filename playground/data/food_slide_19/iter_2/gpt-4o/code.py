from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation object
presentation = Presentation()

# Set dimensions for widescreen
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add background image of snow and trees
background_image_path = "../../media/background_snow_trees.jpg"
slide.shapes.add_picture(background_image_path, 0, 0, width=presentation.slide_width, height=presentation.slide_height)

# Title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "PART 3"
title.font.size = Pt(40)
title.font.bold = True

# Instruction text
instruction_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(14), Inches(1))
instruction_frame = instruction_box.text_frame
instruction = instruction_frame.add_paragraph()
instruction.text = "Fill in the correct past simple or past continuous forms in the short story below."
instruction.font.size = Pt(20)

# Text with blank spaces
text_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(10), Inches(4))
text_frame = text_box.text_frame
text = text_frame.add_paragraph()
text.text = ("Last week, I went (I / go) on a date with a nice boy (I / meet) recently. "
             "(We / skate) in the park, (we / laugh), and (he / make) some really funny jokes, "
             "and then, out of nowhere (I / fall) down (I / break) my leg!")
text.font.size = Pt(18)

# Insert exam icon
icon_path = "../../media/image_1.jpg"
slide.shapes.add_picture(icon_path, Inches(12), Inches(0.5), Inches(2), Inches(2))

# Insert image of people ice skating
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(11), Inches(5), Inches(4), Inches(3))

# Save presentation
presentation.save("render.pptx")