from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color to a darker shade
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(30, 30, 30)  # Dark background color

# Add the conclusion text
content_text = (
    "In today’s world top grades are valuable but they are not the only path to success. "
    "Balance is the Key. Pursue your interests and passions with a holistic approach to "
    "personal and professional development. Focus on your unique strengths and define your "
    "own paths to success."
)
text_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(3))
text_frame = text_box.text_frame

p = text_frame.add_paragraph()
p.text = content_text
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Add the title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(1))
title_frame = title_box.text_frame

title_p = title_frame.add_paragraph()
title_p.text = "Conclusion"
title_p.font.size = Pt(48)
title_p.font.color.rgb = RGBColor(255, 255, 255)
title_p.alignment = PP_ALIGN.CENTER

# Insert the new trophy image
trophy_path = "../../media/trophy_new_design.jpg"  # Updated path for the new trophy image
trophy = slide.shapes.add_picture(trophy_path, Inches(8), Inches(3.5), Inches(2), Inches(2))  # Adjust size and position as needed

# Insert the yellow banner with text
banner_path = "../../media/image_1.jpg"
banner = slide.shapes.add_picture(banner_path, Inches(10), Inches(0.1), Inches(5), Inches(1))

text_box_banner = slide.shapes.add_textbox(Inches(11), Inches(0.2), Inches(4), Inches(0.5))
text_frame_banner = text_box_banner.text_frame

p_banner = text_frame_banner.add_paragraph()
p_banner.text = "To Download Visit www.topicsforseminar.com"
p_banner.font.size = Pt(12)
p_banner.font.color.rgb = RGBColor(255, 255, 255)
p_banner.alignment = PP_ALIGN.RIGHT

# Insert the background illustration
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(3.5), Inches(3), Inches(9), Inches(5))

# Save the presentation
presentation.save("render.pptx")