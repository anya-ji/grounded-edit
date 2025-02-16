from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5] 
slide = presentation.slides.add_slide(slide_layout)

# Add the background images
image_left_path = "../../media/image_0.jpg"
image_right_path = "../../media/blurred_image.jpg"
slide.shapes.add_picture(image_left_path, Inches(0), Inches(0), Inches(8), Inches(9))
slide.shapes.add_picture(image_right_path, Inches(8), Inches(0), Inches(8), Inches(9))

# Add the title text on the right side
title_box = slide.shapes.add_textbox(Inches(9), Inches(1), Inches(6), Inches(1.5))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "MARKET AROUND US CHAPTER-6"
title.font.size = Pt(44)
title.font.bold = True
title.font.color.rgb = RGBColor(255, 255, 255)

# Add attribution text at the bottom
attribution_box = slide.shapes.add_textbox(Inches(1), Inches(8), Inches(14), Inches(1))
attribution_frame = attribution_box.text_frame
attribution = attribution_frame.add_paragraph()
attribution.text = "This Photo by Unknown author is licensed under CC BY."
attribution.font.size = Pt(12)
attribution.font.color.rgb = RGBColor(255, 255, 255)

# Save the presentation
presentation.save("render.pptx")