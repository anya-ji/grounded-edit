from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

# Opens and resizes the image to fit in half of a 16x9 slide.
def resize_image(image_path, target_height):
    img = Image.open(image_path)
    aspect_ratio = img.width / img.height
    new_width = int(aspect_ratio * target_height / 2)
    return img.resize((new_width, target_height))

# Load and resize the image
image_path = "../../media/image_0.jpg"
image = resize_image(image_path, 900)

# Save the resized image
resized_image_path = "resized_image.jpg"
image.save(resized_image_path)

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5] # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add resized image to the left half
left = Inches(0)
top = Inches(0)
height = presentation.slide_height
slide.shapes.add_picture(resized_image_path, left, top, height=height)

# Add teal rectangle over right half
right = presentation.slide_width / 2
width = presentation.slide_width / 2
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, right, top, width, height)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0, 128, 128)  # Teal color

# Add text box for the question
text = "Q. Do you think 'vocal for local' is not just a slogan but it is the need of the time."
text_box = slide.shapes.add_textbox(right, top, width, height)
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = text
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(255, 255, 255)  # White color
text_box.text_frame.vertical_anchor = MSO_SHAPE.MIDDLE  # Center vertically

# Add small red rectangle in top right corner
red_rect_width = Inches(1)
red_rect_height = Inches(0.5)
red_rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, presentation.slide_width - red_rect_width, 0, red_rect_width, red_rect_height)
red_rect.fill.solid()
red_rect.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red color

# Save presentation
presentation.save("render.pptx")