from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_slide():
    # Initialize presentation
    presentation = Presentation()
    presentation.slide_width = Inches(16)
    presentation.slide_height = Inches(9)
    
    # Use a blank slide layout
    slide_layout = presentation.slide_layouts[5]
    slide = presentation.slides.add_slide(slide_layout)
    
    # Set background to a light green with hexagonal pattern
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(210, 255, 210)  # Modified light green
    
    # Add a centered large white text box for title and content
    left = Inches(1)
    top = Inches(1)
    width = Inches(14)
    height = Inches(7)
    textbox = slide.shapes.add_textbox(left, top, width, height)
    textbox.fill.solid()
    textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
    
    # Title: "CONCLUSION"
    title = textbox.text_frame.add_paragraph()
    title.text = "CONCLUSION"
    title.font.bold = True
    title.font.size = Pt(48)
    title.font.color.rgb = RGBColor(0, 128, 0)  # Green font
    title.alignment = PP_ALIGN.LEFT
    
    # Content
    content_text = (
        "Since the dawn of time, human beings have solemnly dependent on mother "
        "nature to fulfil their needs from inhaling air to sustenance to food requirements. "
        "But we are unable to handle the greatness of nature, we have overused, exploited, wasted "
        "natural resources which will affect badly on future generations. So it’s our duty to utilize "
        "the resources carefully. If we do not then that day is not so far where there will be no water, "
        "no food, no oxygen and life will be impossible."
    )
    content = textbox.text_frame.add_paragraph()
    content.text = content_text
    content.font.size = Pt(20)
    content.font.color.rgb = RGBColor(0, 0, 0)  # Black font
    content.alignment = PP_ALIGN.LEFT
    
    # Save presentation
    presentation.save("render.pptx")

# Execute the function to create the slide
create_slide()