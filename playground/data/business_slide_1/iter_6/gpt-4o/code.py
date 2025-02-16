from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color
background_color = RGBColor(252, 110, 100)  # Coral
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = background_color

# Draw a rectangle to create a bold background for the logo
logo_background = slide.shapes.add_shape(
    1,  # Shape type: 1 = Rectangle
    Inches(3), Inches(2), Inches(10), Inches(5)
)
logo_background.fill.solid()
logo_background.fill.fore_color.rgb = RGBColor(252, 110, 100)  # Same as background

# Adjust position and size to center the logo, ensuring full opacity
logo_icon_path = "../../media/image_1_solid_white.jpg"  # Changed to solid white logo
logo_shape = slide.shapes.add_picture(logo_icon_path, Inches(3), Inches(2), Inches(10), Inches(5))

# Add Airbnb text logo on the right, centered vertically
logo_text_path = "../../media/image_0.jpg"
slide.shapes.add_picture(logo_text_path, Inches(8.5), Inches(2.5), Inches(5), Inches(2.5))

# Add title text "Business Case"
title_tx_box = slide.shapes.add_textbox(Inches(8.5), Inches(5), Inches(7), Inches(2))
title_frame = title_tx_box.text_frame
p = title_frame.add_paragraph()
p.text = "Business Case"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)  # White

# Add names below the title
names = ["Daniel Consuegra", "Alejandra Del Chiaro", "Maria Camila Echeverri"]
y_position = 6.5
for name in names:
    name_box = slide.shapes.add_textbox(Inches(8.5), Inches(y_position), Inches(7), Inches(0.5))
    tf = name_box.text_frame
    p = tf.add_paragraph()
    p.text = name
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)  # White
    y_position += 0.5

# Save the presentation
presentation.save("render.pptx")