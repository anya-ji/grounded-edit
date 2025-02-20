from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Define image path
image_path = "../../media/image_0.jpg"

# Increase the overall size of the image
image_width = Inches(5)  # Increased size
image_height = Inches(5)  # Increased size

left = (presentation.slide_width - image_width) / 2
top = (presentation.slide_height - image_height) / 2

slide.shapes.add_picture(image_path, left, top, width=image_width, height=image_height)

# Save the presentation
presentation.save("render.pptx")