from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Define image path
image_path = "../../media/image_0.jpg"

# Add the image in the center of the slide
# Centering calculations
image_width = Inches(3)
image_height = Inches(3)

left = (presentation.slide_width - image_width) / 2
top = (presentation.slide_height - image_height) / 2

slide.shapes.add_picture(image_path, left, top, width=image_width, height=image_height)

# Add a text box with drop shadow effect
text_box = slide.shapes.add_textbox(left, top + image_height + Inches(0.2), width=Inches(6), height=Inches(1))
text_frame = text_box.text_frame
text_frame.text = "Your Text Here"  # Replace with your desired text

# Apply drop shadow effect
for paragraph in text_frame.paragraphs:
    paragraph.font.shadow = True

# Save the presentation
presentation.save("render.pptx")