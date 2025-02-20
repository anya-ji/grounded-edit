from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Define image path
image_path = "../../media/image_0.jpg"

# Add the image in the center of the slide
# Centering calculations
image_width = Inches(3)
image_height = Inches(3)

left = (presentation.slide_width - image_width) / 2
top = (presentation.slide_height - image_height) / 2

slide.shapes.add_picture(image_path, left, top, width=image_width, height=image_height)

# Add text with increased font size
text_box = slide.shapes.add_textbox(left, top + image_height + Inches(0.5), width=Inches(5), height=Inches(1))
text_frame = text_box.text_frame
text = text_frame.add_paragraph()
text.text = "global.media.cultures"
text.font.size = Pt(36)  # Adjusting the font size to 36 points for increased visibility

# Save the presentation
presentation.save("render.pptx")