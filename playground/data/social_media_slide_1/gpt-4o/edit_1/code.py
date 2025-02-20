from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Define image path
image_path = "../../media/image_0.jpg"

# Add the image in the center of the slide
image_width = Inches(3)
image_height = Inches(3)

left = (presentation.slide_width - image_width) / 2
top = (presentation.slide_height - image_height) / 2

slide.shapes.add_picture(image_path, left, top, width=image_width, height=image_height)

# Add text with color change
text_box_left = Inches(4)
text_box_top = Inches(4)
text_box_width = Inches(8)
text_box_height = Inches(2)

text_box = slide.shapes.add_textbox(text_box_left, text_box_top, text_box_width, text_box_height)
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = "global.media cultures"
p.font.size = Pt(40)

# Change the color of the word "cultures" to dark blue
p.font.color.rgb = RGBColor(0, 0, 139)  # Dark blue color for "cultures"

# Save the presentation
presentation.save("render.pptx")