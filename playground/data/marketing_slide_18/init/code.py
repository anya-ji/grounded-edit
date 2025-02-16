from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add title
title_box = slide.shapes.add_textbox(Inches(0), Inches(0.3), Inches(16), Inches(1))
title_frame = title_box.text_frame
title_frame.text = 'More About the Market'
title_frame.paragraphs[0].font.size = Pt(44)
title_frame.paragraphs[0].alignment = 1  # Center

# Colors
background_color = RGBColor(173, 216, 230)  # Light blue
box_color = RGBColor(240, 248, 255)  # Alice blue
text_color = RGBColor(255, 255, 255)  # White

# Set background color
background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    Inches(0), Inches(0), presentation.slide_width, presentation.slide_height
)
background.fill.solid()
background.fill.fore_color.rgb = background_color
background.line.color.rgb = background_color  # No border

# Define box positions
left_x = Inches(0.5)
middle_x = Inches(6)
right_x = Inches(11.5)
y_positions = [Inches(1.5), Inches(3.5), Inches(5.5)]
box_width, box_height = Inches(4), Inches(1.5)

# Define content for each box
content = [
    ["Step 1", "Marketing", "Step 2"],
    ["Storage", "Step 3", "Sales"],
    ["Step 4", "After sales services", "Ensure Customers"]
]

# Add boxes and arrows
for row, y in enumerate(y_positions):
    for col, text in enumerate(content[row]):
        x = [left_x, middle_x, right_x][col]
        box = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, x, y, box_width, box_height
        )
        box.fill.solid()
        box.fill.fore_color.rgb = box_color
        box.line.color.rgb = box_color  # No border

        # Add text
        txBox = slide.shapes.add_textbox(x, y, box_width, box_height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(24)
        p.font.color.rgb = text_color
        p.alignment = 1  # Center

        # Add arrows if not the last column
        if col < 2:
            arrow = slide.shapes.add_connector(
                MSO_SHAPE.LINE_ARROW, x + box_width, y + box_height/2, [middle_x, right_x][col], y + box_height/2
            )
            arrow.line.color.rgb = RGBColor(255, 255, 255)  # White arrows

# Save presentation
presentation.save("render.pptx")