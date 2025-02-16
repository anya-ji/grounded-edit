from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.oxml import parse_xml
from pptx.oxml.ns import nsdecls

# Initialize presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add background image
image_path = "../../media/image_0.jpg"
left = 0
top = 0
width = Inches(9)
height = Inches(9)

pic = slide.shapes.add_picture(image_path, left, top, width, height)
pic.crop_right = 0.5  # Crop the right side to make a half-circled shape

# Add semi-transparent overlay on the right side
overlay = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(9), 0, Inches(7), Inches(9)
)
fill = overlay.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)
fill.fore_color.alpha = 153  # 60% transparency

# Title
title_box = slide.shapes.add_textbox(Inches(9.5), Inches(1), Inches(6), Inches(1))
title = title_box.text_frame.add_paragraph()
title.text = "ELO's"
title.font.size = Pt(40)
title.font.bold = True
title.font.color.rgb = RGBColor(255, 255, 255)  # Light color for contrast

# Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(9.5), Inches(2.5), Inches(6), Inches(1))
subtitle = subtitle_box.text_frame.add_paragraph()
subtitle.text = "1. From production to consumption"
subtitle.font.size = Pt(28)
subtitle.font.color.rgb = RGBColor(255, 255, 255)  # Light color for contrast

# Attribution
attrib_box = slide.shapes.add_textbox(Inches(9.5), Inches(8), Inches(6), Inches(1))
attrib = attrib_box.text_frame.add_paragraph()
attrib.text = "This Photo by Unknown author is licensed under CC BY-SA."
attrib.font.size = Pt(12)
attrib.font.color.rgb = RGBColor(255, 255, 255)  # Light color for contrast

# Save presentation
presentation.save("render.pptx")