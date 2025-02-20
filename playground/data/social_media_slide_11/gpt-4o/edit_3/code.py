from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5] 
slide = presentation.slides.add_slide(slide_layout)

# Set the slide background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Add the Main Title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
title_tf = title_box.text_frame
title_tf.word_wrap = True
p = title_tf.add_paragraph()
p.text = "Various Media drive Various forms of Global Integration"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue background

# Add the Subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.75), Inches(15), Inches(0.8))
subtitle_tf = subtitle_box.text_frame
subtitle_tf.word_wrap = True
p = subtitle_tf.add_paragraph()
p.text = "GLOBAL INTEGRATION"
p.font.size = Pt(36)  # Changed font size to 36
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER
subtitle_box.fill.solid()
subtitle_box.fill.fore_color.rgb = RGBColor(0, 51, 102)

# Add the Content
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(15), Inches(5))
content_tf = content_box.text_frame
content_tf.word_wrap = True

bullet_points = [
    "the process of increasing the degree of economic and political integration among countries around the world.",
    "Global integration can involve the processes of product standardization and technology development centralization."
]

for point in bullet_points:
    p = content_tf.add_paragraph()
    p.text = point
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0)
    p.space_after = Pt(14)  # Spacing between bullet points

content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

# Save the presentation
presentation.save("render.pptx")