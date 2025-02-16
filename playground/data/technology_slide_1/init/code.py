from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set background color
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

# Insert first image on the left
left_image_path = "../../media/image_1.jpg"
left_img = slide.shapes.add_picture(left_image_path, Inches(0), Inches(1), Inches(7.5), Inches(6.5))

# Insert title on the right
title_text_box = slide.shapes.add_textbox(Inches(8), Inches(1), Inches(7.5), Inches(1))
title_text_frame = title_text_box.text_frame
title_text_frame.word_wrap = True

title_p = title_text_frame.add_paragraph()
title_p.text = "Blockchain Technology"
title_p.font.size = Pt(44)
title_p.font.bold = True
title_p.alignment = PP_ALIGN.LEFT

# Insert second image on the right
right_image_path = "../../media/image_0.jpg"
right_img = slide.shapes.add_picture(right_image_path, Inches(8), Inches(2.5), Inches(7.5), Inches(5))

# Save presentation
presentation.save("render.pptx")