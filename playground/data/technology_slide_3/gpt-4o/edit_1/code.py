from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object and set slide dimensions
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add the slide title
title_text = "Origin of Blockchain Technology"
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True

# Configure the title paragraph
title_p = title_frame.add_paragraph()
title_p.text = title_text
title_p.font.bold = True
title_p.font.size = Pt(44)
title_p.alignment = PP_ALIGN.LEFT

# Add bullet points
bullet_text = (
    "- The blockchain was created by a person (or group of people) using the name (or pseudonym) Satoshi Nakamoto in 2008 to serve as the public distributed ledger for bitcoin cryptocurrency transactions.\n"
    "- The implementation of the blockchain within bitcoin made it the first digital currency to solve the double-spending problem without the need of a trusted authority or central server.\n"
    "- Private blockchains have been proposed for business use."
)

content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Add each bullet point as a paragraph with adjusted line spacing
for bullet in bullet_text.split("\n"):
    bullet_p = content_frame.add_paragraph()
    bullet_p.text = bullet
    bullet_p.font.size = Pt(24)
    bullet_p.alignment = PP_ALIGN.LEFT
    bullet_p.space_after = Pt(12)  # Adjust line spacing to 1.5 (approximately)

# Make sure the background is simple to enhance readability
slide_background = slide.background
fill = slide_background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

# Save the presentation
presentation.save("render.pptx")