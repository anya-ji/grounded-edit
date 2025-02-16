from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5] 
slide = presentation.slides.add_slide(slide_layout)

# Set background color to a softer shade of green
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(184, 255, 184) # softer green background

# Add title and content text box
text_box = slide.shapes.add_textbox(Inches(2), Inches(1), Inches(12), Inches(5.5))
text_frame = text_box.text_frame

# Add Title
title = text_frame.add_paragraph()
title.text = "NEGATIVE IMPACT"
title.font.bold = True
title.font.size = Pt(44)
title.font.color.rgb = RGBColor(0, 153, 0) # bold green font

# Add Content
content = [
    "Humans are continuously cutting down trees and littering have a negative impact on environment.",
    "Pollution is everywhere, from the trash thrown out on the free way to the millions of metric tons of pollution pumped into the atmosphere every year.",
    "Burning fossil fuels is the prime culprit for climate change. When oil gas and coal are burned they release carbon dioxide and other harmful gases which trap heat in atmosphere and cause it to warm up.",
    "The increasing use of fertilizers and pesticides to protect crops have damaged the fertility of land.",
    "Food wastage is a huge issue with an estimated one third of all food produced globally is being wasted."
]

for item in content:
    p = text_frame.add_paragraph()
    p.text = item
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0, 0, 0) # black font

text_frame.word_wrap = True

# Insert Deforestation Image
image_path_1 = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path_1, Inches(0.5), Inches(7), Inches(7.5), Inches(2))

# Insert Pollution Image
image_path_2 = "../../media/image_1.jpg"
slide.shapes.add_picture(image_path_2, Inches(8), Inches(7), Inches(7.5), Inches(2))

# Save presentation
presentation.save("render.pptx")