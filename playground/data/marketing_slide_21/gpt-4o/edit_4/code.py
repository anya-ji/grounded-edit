from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set a dark gradient background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(54, 57, 63)

# Function to add a text box
def add_textbox(slide, text, left, top, width, height, font_size=24, bold=False, color=RGBColor(0, 0, 0), align=PP_ALIGN.LEFT):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    p = text_frame.add_paragraph()
    p.text = text
    p.font.bold = bold
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.alignment = align
    return textbox

# Define positions for the flowchart elements
left_margin = Inches(1)
top_margin = Inches(1)

# Define sizes for the rectangles
width = Inches(2)
height = Inches(1)

# Add flowchart elements
elements = [
    ("Direct", left_margin, top_margin),
    ("Retail", left_margin, top_margin + Inches(1.25)),
    ("Wholesale", left_margin, top_margin + Inches(2.5)),
    ("Retailer", left_margin + Inches(5), top_margin + Inches(1.25)),  # Moved this down
    ("Producer or manufacturer", left_margin + Inches(3), top_margin + Inches(1.25)),  # Adjusted position
    ("Wholesaler", left_margin + Inches(5), top_margin + Inches(2.5)),
    ("Consumer", left_margin + Inches(7.5), top_margin + Inches(1.25))
]

# Add rectangles with text
for text, left, top in elements:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(255, 140, 0)  # Orange color
    shape.line.color.rgb = RGBColor(255, 255, 255)  # White border
    add_textbox(slide, text, left, top, width, height, bold=True, align=PP_ALIGN.CENTER)

# Add title at the bottom
add_textbox(slide, "Sales and Distribution", Inches(4), Inches(7.5), Inches(8), Inches(1), font_size=32, bold=True, color=RGBColor(0, 0, 0), align=PP_ALIGN.CENTER)

# Add attribution note
add_textbox(slide, "This Photo by Unknown author is licensed under CC BY-SA-NC.", Inches(9), Inches(8), Inches(7), Inches(0.5), font_size=10, color=RGBColor(255, 255, 255), align=PP_ALIGN.RIGHT)

# Save the presentation
presentation.save("render.pptx")