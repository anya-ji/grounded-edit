from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
presentation = Presentation()

# Set slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(30, 30, 30)

# Function to add a textbox with specified properties
def add_textbox(slide, x, y, w, h, text, font_size, bold=False, color=RGBColor(255, 255, 255)):
    textbox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    text_frame = textbox.text_frame
    text_frame.clear()
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    text_frame.word_wrap = True
    return textbox

# Add slide title
add_textbox(slide, 4, 0.5, 8, 1, "ACADEMIC WORLD AND PROFESSIONAL WORLD: DEMANDS AND CHARACTERISTICS", 32, True, RGBColor(255, 255, 255)).text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add Academic World section
add_textbox(slide, 0.5, 1.5, 7, 1, "ACADEMIC WORLD", 28, True, RGBColor(0, 102, 204))
add_textbox(slide, 0.5, 2.5, 7, 0.5, "Evaluation and Feedback", 24, True, RGBColor(255, 204, 0))
add_textbox(slide, 0.5, 3.2, 7, 1, "Assessment:\n• Students are assessed through exams, assignments, and projects, with a focus on understanding and demonstrating knowledge.", 20)
add_textbox(slide, 0.5, 4.5, 7, 1, "Feedback:\n• Feedback is provided by educators to help students improve their understanding and skills.", 20)
add_textbox(slide, 0.5, 6, 7, 0.5, "Timeframe", 24, True, RGBColor(255, 204, 0))
add_textbox(slide, 0.5, 6.7, 7, 1, "Structured:\n• Academic schedules are typically structured into semesters or quarters, with predefined courses and timelines.", 20)

# Add Professional World section
add_textbox(slide, 8.5, 1.5, 7, 1, "PROFESSIONAL WORLD", 28, True, RGBColor(0, 153, 76))
add_textbox(slide, 8.5, 2.5, 7, 0.5, "Evaluation and Feedback", 24, True, RGBColor(255, 204, 0))
add_textbox(slide, 8.5, 3.2, 7, 1, "Assessment:\n• Professionals are evaluated based on their performance, results, and contributions to their organizations or clients.", 20)
add_textbox(slide, 8.5, 4.5, 7, 1, "Feedback:\n• Feedback is often tied to specific work outcomes and is focused on improvement.", 20)
add_textbox(slide, 8.5, 6, 7, 0.5, "Timeframe", 24, True, RGBColor(255, 204, 0))
add_textbox(slide, 8.5, 6.7, 7, 1, "Varied:\n• Professional work often follows project-based or ongoing timelines, with varying levels of structure and flexibility.", 20)

# Save the presentation
presentation.save("render.pptx")