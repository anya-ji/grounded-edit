from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()

# Define slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background image
image_path = "../../media/image_0.jpg"
left = top = Inches(0)
pic = slide.shapes.add_picture(image_path, left, top, width=Inches(16), height=Inches(9))

# Add the smaller title in the top-left corner
small_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(5), Inches(1))
small_title = small_title_box.text_frame.add_paragraph()
small_title.text = "The Subject and Content of Art"
small_title.font.size = Pt(18)
small_title.font.bold = True
small_title.font.color.rgb = RGBColor(0, 0, 0)
small_title_box.fill.solid()
small_title_box.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow background

# Add the main title
main_title_box = slide.shapes.add_textbox(Inches(6), Inches(1), Inches(8), Inches(1.5))
main_title = main_title_box.text_frame.add_paragraph()
main_title.text = "The Subject of Art"
main_title.font.size = Pt(32)
main_title.font.bold = True
main_title.font.color.rgb = RGBColor(0, 0, 0)
main_title_box.fill.solid()
main_title_box.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow background

# Add bullet points
content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(4.5))  # Adjusted height
text_frame = content_box.text_frame
text_frame.word_wrap = True

# Set the background of the content box to a vibrant color
content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(255, 204, 0)  # Changed to a vibrant orange background

# Add first bullet point
p1 = text_frame.add_paragraph()
p1.text = "The subject of art is usually anything that is represented in the artwork."
p1.font.size = Pt(24)
p1.font.color.rgb = RGBColor(0, 0, 0)  # Changed to black
p1.font.bold = True  # Make the text bold

# Add second bullet point
p2 = text_frame.add_paragraph()
p2.text = "It may be a person, object, scene or event."
p2.font.size = Pt(24)
p2.font.color.rgb = RGBColor(0, 0, 0)  # Changed to black
p2.font.bold = True  # Make the text bold

# Add a darker outline to the content box
content_box.line.color.rgb = RGBColor(0, 0, 0)  # Black outline
content_box.line.width = Pt(2)  # Set line width

# Save the presentation
presentation.save("render.pptx")