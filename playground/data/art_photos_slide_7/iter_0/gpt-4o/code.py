from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add main title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.1), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_frame.text = "The Subject and Content of Art"
title_frame.paragraphs[0].font.size = Pt(24)
title_frame.paragraphs[0].font.bold = True
title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(255, 255, 0)

# Add main content title
main_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(15), Inches(1))
main_title_frame = main_title_box.text_frame
main_title_frame.word_wrap = True
main_title_frame.text = "Representational or Objective Art"
main_title_frame.paragraphs[0].font.size = Pt(32)
main_title_frame.paragraphs[0].font.bold = True
main_title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)
main_title_box.fill.solid()
main_title_box.fill.fore_color.rgb = RGBColor(255, 255, 0)

# Add main points
content_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(14), Inches(6))
content_frame = content_box.text_frame
content_frame.word_wrap = True

p = content_frame.add_paragraph()
p.text = "Examples:"
p.font.size = Pt(28)

p = content_frame.add_paragraph()
p.text = "A. *Still life* is a work of art that does not have life."
p.font.size = Pt(24)
p.space_after = Pt(6)
p.level = 1  # Sub-bullet

p = content_frame.add_paragraph()
p.text = "examples: (food, flower, plants, rocks or shells) or man-made (drinking glasses, books, vases, jewelry, coins and so on) in an artificial setting."
p.font.size = Pt(24)
p.level = 1  # Sub-bullet

# Set content box background to light gray for readability
content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(211, 211, 211)  # Light Gray

presentation.save("render.pptx")