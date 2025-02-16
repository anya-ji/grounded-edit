from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Create a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add a light gradient background
background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), presentation.slide_width, presentation.slide_height  # Adjusted width
)
fill = background.fill
fill.gradient()  # Change to gradient fill
gradient_stops = fill.gradient_stops
gradient_stops[0].color.rgb = RGBColor(173, 216, 230)  # Light blue color
gradient_stops[1].color.rgb = RGBColor(144, 238, 144)  # Light green color for gradient depth

# Remove the text box
# (The text box code has been omitted)

# Add image of character holding money without surrounding rectangle
image_path_0 = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path_0, Inches(5.5), Inches(2.5), height=Inches(4))  # Adjusted position to center

# Add image of various coins and bills
image_path_1 = "../../media/image_1.jpg"
slide.shapes.add_picture(image_path_1, Inches(9), Inches(2.5), height=Inches(4))  # Adjusted position to center

# Save the presentation
presentation.save("render.pptx")