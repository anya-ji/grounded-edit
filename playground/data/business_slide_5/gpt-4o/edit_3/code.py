from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]  # 5 is a blank slide
slide = presentation.slides.add_slide(slide_layout)

# Add title with teal background
title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, presentation.slide_width, Inches(1))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = RGBColor(0, 128, 128)  # Teal color
title_shape.line.color.rgb = RGBColor(0, 128, 128)  # No border

# Add title text
title_text = title_shape.text_frame
p = title_text.add_paragraph()
p.text = "Values"
p.font.bold = True
p.font.size = Pt(40)
p.font.color.rgb = RGBColor(0, 0, 0)  # Black text
title_text.word_wrap = True

# Add content text box
content_left = 0.5
content_top = 1.2
content_width = 8.5
content_height = 6.5

content_textbox = slide.shapes.add_textbox(Inches(content_left), Inches(content_top), Inches(content_width), Inches(content_height))
text_frame = content_textbox.text_frame

# Add content
values = [
    ("Champion the Mission", "We’re united with our community to create a world where anyone can belong anywhere."),
    ("Be a Host", "We're caring, open, and encouraging to everyone we work with."),
    ("Simplify", ""),
    ("Every Frame Matters", ""),
    ("Be a Cereal Entrepreneur", "We’re determined and creative in transforming our bold ambitions into reality."),
    ("Embrace the Adventure", "We’re driven by curiosity, optimism, and the belief that every person can grow."),
]

# Loop through values and add to the slide
for title, desc in values:
    p = text_frame.add_paragraph()
    p.text = title
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black
    if desc:
        p = text_frame.add_paragraph()
        p.text = desc
        p.font.size = Pt(20)
        p.font.color.rgb = RGBColor(0, 0, 0)  # Black

text_frame.word_wrap = True

# Add image to the slide
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(9.5), Inches(1.2), height=Inches(6.5))

# Save the presentation
presentation.save("render.pptx")