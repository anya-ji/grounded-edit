from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]  # 5 is a blank slide
slide = presentation.slides.add_slide(slide_layout)

# Add title with a more vivid teal background
title_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, presentation.slide_width, Inches(1))
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = RGBColor(0, 150, 150)  # Brighter teal color
title_shape.line.color.rgb = RGBColor(0, 150, 150)  # No border

# Add title text
title_text = title_shape.text_frame
p = title_text.add_paragraph()
p.text = "Values"
p.font.bold = True
p.font.size = Pt(40)
p.font.color.rgb = RGBColor(255, 255, 255)  # White text for higher contrast
title_text.word_wrap = True

# Add content with city posters and logos
cities = [
    ("NEW YORK", "path/to/ny_image.jpg"),
    ("TOKYO", "path/to/tokyo_image.jpg"),
    ("RIO DE JANEIRO", "path/to/rio_image.jpg"),
    ("CITY FOUR", "path/to/city4_image.jpg"),
    ("CITY FIVE", "path/to/city5_image.jpg"),
    ("CITY SIX", "path/to/city6_image.jpg"),
]

grid_left = 0.5
grid_top = 1.5
grid_width = 12
grid_height = 6

# Loop to add city posters with the Airbnb logo
for index, (city_name, image_path) in enumerate(cities):
    left = grid_left + (index % 3) * (grid_width / 3)
    top = grid_top + (index // 3) * (grid_height / 2)
    poster_shape = slide.shapes.add_picture(image_path, Inches(left), Inches(top), width=Inches(4))
    
    # Add city name text
    city_frame = slide.shapes.add_textbox(Inches(left), Inches(top + poster_shape.height), Inches(4), Inches(0.5))
    text_frame = city_frame.text_frame
    p = text_frame.add_paragraph()
    p.text = city_name
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(255, 255, 255)  # White color for city names
    text_frame.word_wrap = True

# Save the presentation
presentation.save("render.pptx")