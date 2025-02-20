from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue

# Title text
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
title_frame = title.text_frame
title_frame.clear()
p = title_frame.paragraphs[0]
p.text = "GLOBAL MEDIA CULTURES"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)  # White
p.alignment = PP_ALIGN.CENTER

# Bullet points text
left_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(7), Inches(6))
tf = left_textbox.text_frame
tf.clear()

bullet_points = [
    "Globalization entails the spread of various cultures",
    "Globalization also involves the spread of ideas",
    "Globalization relies on media as its main conduit for the spread of global culture and ideas"
]

for point in bullet_points:
    p = tf.add_paragraph()
    p.text = f"✔ {point}"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)  # White
    p.space_before = Pt(14)

# Set justification for each bullet point
for paragraph in tf.paragraphs:
    paragraph.alignment = PP_ALIGN.JUSTIFY

# Insert image on the right side
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(8), Inches(2), Inches(7.5), Inches(6))

# Save the presentation
presentation.save("render.pptx")