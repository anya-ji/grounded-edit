from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image
import numpy as np

# Function to increase image saturation
def increase_saturation(image_path, increase_by=1.15):
    img = Image.open(image_path)
    img = img.convert("RGBA")
    
    # Convert to numpy array
    arr = np.array(img)
    
    # Convert RGB to HSV
    rgb = arr[..., :3]
    hsv = (arr[..., :3] / 255.0).dot([[0.299], [0.587], [0.114]])  # Calculate brightness
    s = np.linalg.norm(rgb - hsv, axis=-1)  # Calculate saturation
    
    # Increase saturation
    s = np.clip(s * increase_by, 0, 1)
    
    # Example rendering, you can handle RGB values accordingly
    arr[..., :3] = (arr[..., :3] * s[..., np.newaxis]).astype(np.uint8)
    
    # Save the modified image
    new_image_path = "enhanced_image.png"
    Image.fromarray(arr).save(new_image_path)
    
    return new_image_path

# Create a presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue

# Title text
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
title_frame = title.text_frame
title_frame.clear()
p = title_frame.paragraphs[0]
p.text = "GLOBAL MEDIA CULTURES"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)  # White
p.alignment = PP_ALIGN.CENTER

# Bullet points text
left_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(7), Inches(6))
tf = left_textbox.text_frame
tf.clear()

bullet_points = [
    "Globalization entails the spread of various cultures",
    "Globalization also involves the spread of ideas",
    "Globalization relies on media as its main conduit for the spread of global culture and ideas"
]

for point in bullet_points:
    p = tf.add_paragraph()
    p.text = f"✔ {point}"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(255, 255, 255)  # White
    p.space_before = Pt(14)

# Process image and insert it on the right side
image_path = "../../media/image_0.jpg"
enhanced_image_path = increase_saturation(image_path)
slide.shapes.add_picture(enhanced_image_path, Inches(8), Inches(2), Inches(7.5), Inches(6))

# Save the presentation
presentation.save("render.pptx")