from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(10, 31, 69)  # dark blue

# Change the title and make it bold and larger
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
title_frame = title_box.text_frame
p = title_frame.add_paragraph()
p.text = "Friends | Joey Doesn't Share Food!"
p.font.size = Pt(28)  # Increased font size
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)  # white

# Add image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(9), Inches(2), height=Inches(4))

# Add text above image (removed duplicate title)
text_box = slide.shapes.add_textbox(Inches(9), Inches(1.3), Inches(6), Inches(0.6))
text_frame = text_box.text_frame
text_frame.text = "Friends | Joey Doesn't Share Food!"
text_frame.paragraphs[0].font.size = Pt(20)
text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # white

# Add text on the left side
left_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(4), Inches(1))
left_text_frame = left_text_box.text_frame
left_paragraph = left_text_frame.add_paragraph()
left_paragraph.text = "Finish watching the video."
left_paragraph.font.size = Pt(20)
left_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # white

# Add "FLUENT" logo or text
fluent_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(3), Inches(0.5))
fluent_frame = fluent_box.text_frame
fluent_paragraph = fluent_frame.add_paragraph()
fluent_paragraph.text = "FLUENT"
fluent_paragraph.font.size = Pt(24)
fluent_paragraph.font.bold = True
fluent_paragraph.font.color.rgb = RGBColor(255, 223, 77)  # contrasting color

# Save presentation
presentation.save("render.pptx")