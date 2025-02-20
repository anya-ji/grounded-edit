from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Increase top margin by adjusting the Y position
top_bar = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0.5), Inches(16), Inches(1)
)
fill = top_bar.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 105, 180) # Pink color

# Adjust title textbox position due to increased margin
title_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.7), Inches(5), Inches(0.8))
title_frame = title_box.text_frame
title_frame.clear()

title = title_frame.add_paragraph()
title.text = "Vision"
title.font.bold = True
title.font.size = Pt(32)
title.font.color.rgb = RGBColor(255, 255, 255)

# Insert the airbnb logo at the bottom left corner
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0.5), Inches(7), Inches(2), Inches(2))

# Save the presentation
presentation.save("render.pptx")