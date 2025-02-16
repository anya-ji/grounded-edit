from pptx import Presentation
from pptx.util import Inches
from PIL import Image, ImageEnhance

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide layout
slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout for flexibility
slide = presentation.slides.add_slide(slide_layout)

# Insert the provided image with adjusted positioning and size
image_path = "../../media/image_0.jpg"

# Open the image, adjust brightness and contrast, and save it temporarily
with Image.open(image_path) as img:
    enhancer_brightness = ImageEnhance.Brightness(img)
    img = enhancer_brightness.enhance(1.2)  # Increase brightness
    enhancer_contrast = ImageEnhance.Contrast(img)
    img = enhancer_contrast.enhance(1.2)  # Increase contrast

    # Create a light blue background image
    background_color = (173, 216, 230)  # Light blue RGB
    background = Image.new('RGB', img.size, background_color)

    # Paste the modified image onto the background
    background.paste(img, (0, 0), img)

    temp_image_path = "../../media/temp_image.jpg"
    background.save(temp_image_path)

# Add the modified image to the slide
slide.shapes.add_picture(temp_image_path, Inches(0.5), Inches(1), Inches(15), Inches(7))  # Adjusted size and position

# Save the presentation
presentation.save("render.pptx")