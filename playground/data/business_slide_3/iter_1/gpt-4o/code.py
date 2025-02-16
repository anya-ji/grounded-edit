from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide layout
slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout for flexibility
slide = presentation.slides.add_slide(slide_layout)

# Insert the provided image with adjusted positioning and size
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0.5), Inches(1), Inches(15), Inches(7))  # Adjusted size and position

# Save the presentation
presentation.save("render.pptx")