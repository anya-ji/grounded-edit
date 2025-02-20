from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide layout
slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout for flexibility
slide = presentation.slides.add_slide(slide_layout)

# Insert the provided image prominently
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(2), Inches(2), Inches(12), Inches(6))

# Add a text box with white text
left = Inches(2)
top = Inches(1)
width = Inches(12)
height = Inches(1.5)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = "Your white text here"  # Replace with your desired text
p.font.color.rgb = RGBColor(255, 255, 255)  # Set the text color to white

# Save the presentation
presentation.save("render.pptx")