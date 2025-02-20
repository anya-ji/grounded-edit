from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide layout
slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout for flexibility
slide = presentation.slides.add_slide(slide_layout)

# Insert the provided image prominently
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(2), Inches(2), Inches(12), Inches(6))

# Add text box with padding
from pptx.util import Pt
from pptx.dml.color import RGBColor

left = Inches(1)
top = Inches(1)
width = Inches(14)
height = Inches(2)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
text_frame.text = "We imagine a world where you can belong anywhere."
for paragraph in text_frame.paragraphs:
    paragraph.font.size = Pt(32)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
    paragraph.space_after = Pt(14)

# Save the presentation
presentation.save("render.pptx")