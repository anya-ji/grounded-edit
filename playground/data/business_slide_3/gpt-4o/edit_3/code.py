from pptx import Presentation
from pptx.util import Inches

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide layout
slide_layout = presentation.slide_layouts[5]  # Use a blank slide layout for flexibility
slide = presentation.slides.add_slide(slide_layout)

# Insert the provided image prominently
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(2), Inches(2), Inches(12), Inches(6))

# Add a centered text box
left = Inches(2)
top = Inches(8)
width = Inches(12)
height = Inches(1)
textbox = slide.shapes.add_textbox(left, top, width, height)
text_frame = textbox.text_frame
p = text_frame.add_paragraph()
p.text = "Your centered text here"
p.alignment = 1  # Set alignment to center (1 = center alignment)

# Save the presentation
presentation.save("render.pptx")