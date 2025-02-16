from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add abstract background image
abstract_image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(abstract_image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add a semi-transparent white shape for text background
shape_left = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(2))
shape_left.fill.solid()
shape_left.fill.fore_color.rgb = RGBColor(255, 255, 255)
shape_left.fill.transparency = 0.2
shape_left.line.fill.background()

# Add the main title
title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(16), Inches(1))
title_frame = title_box.text_frame
title_frame.text = "Landscapes, seascapes, moonscapes, and cityscape"

title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(40)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(0, 0, 0)
title_paragraph.alignment = PP_ALIGN.CENTER

# Add smaller title
small_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(8), Inches(0.5))
small_title_frame = small_title_box.text_frame
small_title_frame.text = "The Subject and Content of Art"

small_title_paragraph = small_title_frame.paragraphs[0]
small_title_paragraph.font.size = Pt(24)
small_title_paragraph.font.bold = True
small_title_paragraph.font.color.rgb = RGBColor(0, 0, 0)

# Add a semi-transparent white layer for images background
shape_bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(2), Inches(16), Inches(7))
shape_bottom.fill.solid()
shape_bottom.fill.fore_color.rgb = RGBColor(255, 255, 255)
shape_bottom.fill.transparency = 0.2
shape_bottom.line.fill.background()

# Add the city view image
city_image_path = "../../media/image_1.jpg"
slide.shapes.add_picture(city_image_path, Inches(1), Inches(2.5), Inches(7), Inches(5))

# Add the countryside image
countryside_image_path = "../../media/image_2.jpg"
slide.shapes.add_picture(countryside_image_path, Inches(8), Inches(2.5), Inches(7), Inches(5))

# Save the presentation
presentation.save("render.pptx")