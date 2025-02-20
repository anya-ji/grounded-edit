from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

# Add slide
slide_layout = prs.slide_layouts[5]  # blank layout
slide = prs.slides.add_slide(slide_layout)

# Set slide background color to black
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black background

# Add title
title_text = "KEY REASONS WHY ACADEMIC SUCCESS IS IMPORTANT IN SOCIETY"
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_p = title_frame.add_paragraph()
title_p.text = title_text
title_p.font.bold = True
title_p.font.size = Pt(44)
title_p.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Section 1: Global Competitiveness
sec1_img_path = "../../media/image_0.jpg"
slide.shapes.add_picture(sec1_img_path, Inches(1), Inches(2.5), Inches(4), Inches(3))

sec1_title_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(4), Inches(1))
sec1_title_frame = sec1_title_box.text_frame
sec1_title_p = sec1_title_frame.add_paragraph()
sec1_title_p.text = "Global Competitiveness"
sec1_title_p.font.bold = True
sec1_title_p.font.size = Pt(30)
sec1_title_p.font.color.rgb = RGBColor(255, 255, 0)  # Yellow text

sec1_text_box = slide.shapes.add_textbox(Inches(1), Inches(6.2), Inches(4.5), Inches(1.8))
sec1_text_frame = sec1_text_box.text_frame
sec1_text_frame.word_wrap = True
sec1_text_p = sec1_text_frame.add_paragraph()
sec1_text_p.text = "Education is key for a country's competitiveness in the global market."
sec1_text_p.font.size = Pt(20)
sec1_text_p.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Section 2: Scientific Advancement and Research
sec2_img_path = "../../media/image_1.jpg"
slide.shapes.add_picture(sec2_img_path, Inches(6), Inches(2.5), Inches(4), Inches(3))

sec2_title_box = slide.shapes.add_textbox(Inches(6), Inches(5.5), Inches(4), Inches(1))
sec2_title_frame = sec2_title_box.text_frame
sec2_title_p = sec2_title_frame.add_paragraph()
sec2_title_p.text = "Scientific Advancement and Research"
sec2_title_p.font.bold = True
sec2_title_p.font.size = Pt(30)
sec2_title_p.font.color.rgb = RGBColor(255, 255, 0)  # Yellow text

sec2_text_box = slide.shapes.add_textbox(Inches(6), Inches(6.2), Inches(4.5), Inches(1.8))
sec2_text_frame = sec2_text_box.text_frame
sec2_text_frame.word_wrap = True
sec2_text_p = sec2_text_frame.add_paragraph()
sec2_text_p.text = "Academic success drives research, leading to advancements that improve society."
sec2_text_p.font.size = Pt(20)
sec2_text_p.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Section 3: Cultural Enrichment and Social Cohesion
sec3_img_path = "../../media/image_2.jpg"
slide.shapes.add_picture(sec3_img_path, Inches(11), Inches(2.5), Inches(4), Inches(3))

sec3_title_box = slide.shapes.add_textbox(Inches(11), Inches(5.5), Inches(4), Inches(1))
sec3_title_frame = sec3_title_box.text_frame
sec3_title_p = sec3_title_frame.add_paragraph()
sec3_title_p.text = "Cultural Enrichment and Social Cohesion"
sec3_title_p.font.bold = True
sec3_title_p.font.size = Pt(30)
sec3_title_p.font.color.rgb = RGBColor(255, 255, 0)  # Yellow text

sec3_text_box = slide.shapes.add_textbox(Inches(11), Inches(6.2), Inches(4.5), Inches(1.8))
sec3_text_frame = sec3_text_box.text_frame
sec3_text_frame.word_wrap = True
sec3_text_p = sec3_text_frame.add_paragraph()
sec3_text_p.text = "Education fosters cultural understanding and social cohesion through exposure to diverse perspectives and experiences, building a more inclusive society."
sec3_text_p.font.size = Pt(20)
sec3_text_p.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Save presentation
prs.save("examples/career/slide_6/gpt_4o.pptx")