from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add title
title = "Blockchain in Bitcoin"
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
title_paragraph = title_frame.add_paragraph()
title_paragraph.text = title
title_paragraph.font.size = Pt(36)
title_paragraph.font.bold = True

# Add body text
body_text = """CRYPTOCURRENCY bitcoin is one of the crypto currency
Bitcoin is one of the most popular and successful implementations of blockchain technology. It is an open source cryptocurrency that uses distributed peer-to-peer computing. There is no need of a central authority to manage bitcoin network."""
body_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(4))
body_frame = body_box.text_frame
body_frame.word_wrap = True

for line in body_text.split('\n'):
    p = body_frame.add_paragraph()
    p.text = line
    p.font.size = Pt(18)

# Add image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(12), Inches(2), Inches(3.5), Inches(3.5))

# Save presentation
presentation.save("render.pptx")