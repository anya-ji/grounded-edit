from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Get the blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color to a light color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 204)  # Light yellow background

# Add background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, 0, 0, presentation.slide_width, presentation.slide_height)

# Add main title
title_box = slide.shapes.add_textbox(Inches(3), Inches(1), Inches(10), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True
p = title_frame.add_paragraph()
p.text = "KEEPING ART"
p.font.bold = True
p.font.size = Pt(40)  # Changed to a smaller font size
p.font.color.rgb = RGBColor(0, 0, 0)  # Changed to black for better readability

# Style title background to solid yellow
fill = title_box.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 0)  # Changed to a solid yellow

# Add smaller title in top-left corner
small_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
small_title = small_title_box.text_frame
small_title.word_wrap = True
p = small_title.add_paragraph()
p.text = "The Subject and Content of Art"
p.font.bold = True
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(0, 0, 0)

# Style small title box background
small_fill = small_title_box.fill
small_fill.solid()
small_fill.fore_color.rgb = RGBColor(255, 255, 0)

# Add content box
content_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(12), Inches(5))
content_frame = content_box.text_frame
content_frame.word_wrap = True
content_fill = content_box.fill
content_fill.solid()
content_fill.fore_color.rgb = RGBColor(255, 255, 255)  # Changed background to white
content_box.fill.transparency = 0  # Ensure box is fully opaque

# Add content
content = [
    "1. National Pride and Glory",
    "✓ nations keep art because it is good for the economy and for the business."
]

for idx, text in enumerate(content):
    p = content_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(28)
    if idx == 0:
        p.font.bold = True

# Save the presentation
presentation.save("render.pptx")