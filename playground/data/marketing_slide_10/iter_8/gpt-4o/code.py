from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Create gradient background (darker gradient)
background = slide.shapes.add_shape(
    1,  # msoShapeRectangle
    0, 0, Inches(16), Inches(9)
)
fill = background.fill
fill.gradient()
fill.gradient_stops[0].color.rgb = RGBColor(0, 80, 80)  # Darker Cyan
fill.gradient_stops[1].color.rgb = RGBColor(0, 0, 139)  # Dark Blue
fill.gradient_angle = 90  # Adjust gradient angle

# Add gradient overlay (dark blue-green)
overlay = slide.shapes.add_shape(
    1,  # msoShapeRectangle
    0, 0, Inches(16), Inches(9)
)
overlay_fill = overlay.fill
overlay_fill.gradient()
overlay_fill.gradient_stops[0].color.rgb = RGBColor(0, 139, 139)  # Dark Cyan
overlay_fill.gradient_stops[1].color.rgb = RGBColor(0, 0, 139)  # Dark Blue
overlay_fill.gradient_angle = 90  # Adjust gradient angle

# Add soft focus effect by adding a blurred rectangle
blurred_background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    0, 0, Inches(16), Inches(9)
)
blurred_fill = blurred_background.fill
blurred_fill.solid()
blurred_fill.fore_color.rgb = RGBColor(0, 139, 139)  # Dark Cyan
# Simulate a blur effect by adjusting transparency
blurred_fill.transparency = 0.3

# Add abstract shapes with gradients for visual interest
for i in range(5):  # Adding 5 subtle shapes
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(i * 2), Inches(i), Inches(1.5), Inches(1.5)
    )
    shape_fill = shape.fill
    shape_fill.gradient()
    shape_fill.gradient_stops[0].color.rgb = RGBColor(255, 255, 255)  # White
    shape_fill.gradient_stops[1].color.rgb = RGBColor(0, 255, 255)  # Cyan
    shape_fill.gradient_angle = 45  # Gradient angle for variety
    shape.fill.transparency = 0.5  # Slight transparency for subtlety

# Add title text box
title_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
title_frame = title_text_box.text_frame
title_p = title_frame.add_paragraph()
title_p.text = "HOMEWORK"
title_p.font.bold = False  # Adjusted to non-bold
title_p.font.size = Pt(80)  # Adjusted font size
title_p.font.color.rgb = RGBColor(255, 255, 255)

# Add question text box
question_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(15), Inches(3))
question_frame = question_text_box.text_frame
question_p = question_frame.add_paragraph()
question_p.text = "Q1. Explain how a chain of market is formed? What purpose does it serve. (3+2=5 MARKS)"
question_p.font.size = Pt(32)
question_p.font.color.rgb = RGBColor(255, 255, 255)

# Save the presentation
presentation.save("render.pptx")