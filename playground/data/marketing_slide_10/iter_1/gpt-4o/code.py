from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set image path
image_path = "../../media/image_0.jpg"

# Insert background image
slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add title text box
title_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
title_frame = title_text_box.text_frame
title_p = title_frame.add_paragraph()
title_p.text = "HOMEWORK"
title_p.font.bold = False  # Adjusted to non-bold
title_p.font.size = Pt(80)  # Adjusted font size
title_p.font.color.rgb = RGBColor(255, 255, 255)

# Add question text box
question_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(15), Inches(3))
question_frame = question_text_box.text_frame
question_p = question_frame.add_paragraph()
question_p.text = "Q1. Explain how a chain of market is formed? What purpose does it serve. (3+2=5 MARKS)"
question_p.font.size = Pt(32)
question_p.font.color.rgb = RGBColor(255, 255, 255)

# Add a semi-transparent white rectangle over the central area of the image
from pptx.util import Pt

left_inch = Inches(7)  # Adjusted for the central figure's position
top_inch = Inches(3)   # Adjusted for height
width_inch = Inches(2)  # Width of the overlay
height_inch = Inches(2)  # Height of the overlay
overlay = slide.shapes.add_shape(
    1,  # 1 corresponds to the shape type for rectangle
    left_inch, top_inch, width_inch, height_inch
)
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(255, 255, 255)
overlay.fill.transparency = 0.5  # Adjust the transparency here

# Save the presentation
presentation.save("render.pptx")