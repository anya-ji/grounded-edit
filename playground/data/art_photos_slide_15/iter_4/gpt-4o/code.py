from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageEnhance

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background color to a light color (e.g., light gray)
background_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, presentation.slide_width, presentation.slide_height)
background_shape.fill.solid()
background_shape.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

# Add the main title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.3), Inches(14), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Non-representational or Non-objective Art"
title.font.size = Pt(48)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)  # Changed to black for contrast
title.alignment = PP_ALIGN.CENTER

# Add subtitle in the top-left corner
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(6), Inches(0.5))
subtitle_frame = subtitle_box.text_frame
subtitle = subtitle_frame.add_paragraph()
subtitle.text = "The Subject and Content of Art"
subtitle.font.size = Pt(28)
subtitle.font.color.rgb = RGBColor(0, 0, 0)  # Changed to black for contrast

# Enhance both the saturation and brightness of the left painting image
image1_path = "../../media/image_2.jpg"
enhanced_image_path = "../../media/enhanced_image_2.jpg"

# Open the image
image = Image.open(image1_path)
# Enhance saturation
saturation_enhancer = ImageEnhance.Color(image)
image_enhanced = saturation_enhancer.enhance(2.0)  # Increased saturation factor from 1.8 to 2.0

# Enhance brightness
brightness_enhancer = ImageEnhance.Brightness(image_enhanced)
image_final = brightness_enhancer.enhance(1.2)  # Adjust the brightness factor
image_final.save(enhanced_image_path)

# Add images side by side
image2_path = "../../media/image_3.jpg"
left = Inches(1)
top = Inches(2.5)
image_width = Inches(6)
image_height = Inches(4.5)

pic1 = slide.shapes.add_picture(enhanced_image_path, left, top, width=image_width, height=image_height)
pic2 = slide.shapes.add_picture(image2_path, left + image_width + Inches(1), top, width=image_width, height=image_height)

# Add white borders around images
for pic in [pic1, pic2]:
    line = pic.line
    line.color.rgb = RGBColor(255, 255, 255)
    line.width = Pt(4)

# Save the presentation
presentation.save("render.pptx")