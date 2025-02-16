from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the slide background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 87, 87)  # Coral color

# Add the logos image, centered horizontally
image_path = "../../media/image_0.jpg"
left = (presentation.slide_width - Inches(8)) / 2
top = (presentation.slide_height - Inches(4)) / 2
logo = slide.shapes.add_picture(image_path, left, top, width=Inches(8))

# Add text below each logo
text_box1 = slide.shapes.add_textbox(left + Inches(1.2), top + Inches(3.5), Inches(2), Inches(1))
tf1 = text_box1.text_frame
p1 = tf1.add_paragraph()
p1.text = "Belo"
p1.font.size = Pt(24)
p1.font.bold = True
p1.font.color.rgb = RGBColor(255, 255, 255)

text_box2 = slide.shapes.add_textbox(left + Inches(5.2), top + Inches(3.5), Inches(2.5), Inches(1))
tf2 = text_box2.text_frame
p2 = tf2.add_paragraph()
p2.text = "Community"
p2.font.size = Pt(24)
p2.font.bold = True
p2.font.color.rgb = RGBColor(255, 255, 255)

# Save the presentation
presentation.save("render.pptx")