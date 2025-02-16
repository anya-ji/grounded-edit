from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Select a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Create a gradient background
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 204, 0)  # Dominant yellow color
fill.gradient()  # Start gradient fill
stops = fill.gradient_stops
stops[0].position = 0.0
stops[0].color.rgb = RGBColor(255, 204, 0)  # Yellow
stops[1].position = 0.5
stops[1].color.rgb = RGBColor(255, 153, 51)  # Orange
stops[2].position = 1.0
stops[2].color.rgb = RGBColor(255, 255, 255)  # White

# Title: Three levels of meaning
title_box = slide.shapes.add_textbox(Inches(0), Inches(1), Inches(16), Inches(1.5))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Three levels of meaning"
title.font.size = Pt(40)
title.font.bold = True
title.font.color.rgb = RGBColor(255, 255, 255)
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(255, 204, 0)  # Brighter yellow with opacity
title_box.fill.fore_color.alpha = 180  # Adjust opacity for better contrast

# Smaller title: The Subject and Content of Art
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(5), Inches(0.8))  # Adjusted vertical position
subtitle_frame = subtitle_box.text_frame
subtitle = subtitle_frame.add_paragraph()
subtitle.text = "The Subject and Content of Art"
subtitle.font.size = Pt(24)
subtitle.font.color.rgb = RGBColor(255, 255, 255)
subtitle_box.fill.solid()
subtitle_box.fill.fore_color.rgb = RGBColor(255, 204, 0)  # Brighter yellow with opacity
subtitle_box.fill.fore_color.alpha = 180  # Adjust opacity for better contrast

# Content box
content_box = slide.shapes.add_textbox(Inches(1), Inches(4), Inches(14), Inches(5))  # Adjusted vertical position
content_frame = content_box.text_frame

# Add content
content = [
    "1. Factual meaning - the literal statement or narrative content in the work "
    "that can be directly apprehended because the objects presented are easily recognized.",
    "2. Conventional meaning - refers to the special meaning that the certain object "
    "or color has for a particular culture or group of people when it is shown in an artwork."
]

for point in content:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(20)
    p.space_after = Pt(10)

content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

# Save the presentation to a file
presentation.save("render.pptx")