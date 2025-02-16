from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Select a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add the background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), presentation.slide_width, presentation.slide_height)

# Title: Three levels of meaning
title_box = slide.shapes.add_textbox(Inches(0), Inches(1), Inches(16), Inches(1.5))  # Adjusted width here
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Three levels of meaning"
title.font.size = Pt(40)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)  # Black text
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow background

# Smaller title: The Subject and Content of Art
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(5), Inches(0.8))
subtitle_frame = subtitle_box.text_frame
subtitle = subtitle_frame.add_paragraph()
subtitle.text = "The Subject and Content of Art"
subtitle.font.size = Pt(24)
subtitle.font.color.rgb = RGBColor(0, 0, 0)  # Black text
subtitle_box.fill.solid()
subtitle_box.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow background

# Content box
content_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(5))
content_frame = content_box.text_frame

# Add content
content = [
    "1. Factual meaning - the literal statement or narrative content in the work " 
    "that can be directly apprehended because the objects presented are easily recognized.",
    "2. Conventional meaning - refers to the special meaning that the certain object " 
    "or color has for a particular culture or group of people when it is shown in an artwork."
]

for point in content:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(20)
    p.space_after = Pt(10)

content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

# Save the presentation to a file
presentation.save("render.pptx")