from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background image
img_path = "../../media/image_0.jpg"
slide.shapes.add_picture(img_path, Inches(0), Inches(0), width=presentation.slide_width, height=presentation.slide_height)

# Add main title
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(2))
title_frame = title_box.text_frame
title_frame.text = "Non-representational or Non-objective Art"
title_frame.paragraphs[0].font.size = Pt(40)
title_frame.paragraphs[0].bold = True
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(255, 255, 0)
title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

# Add subtitle
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = "Sources of the Subject"
subtitle_frame.paragraphs[0].font.size = Pt(32)
subtitle_frame.paragraphs[0].bold = True
subtitle_box.fill.solid()
subtitle_box.fill.fore_color.rgb = RGBColor(255, 255, 0)
subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Changed font color to black

# Add bullet points
content_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(14), Inches(4))
content_frame = content_box.text_frame
content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(255, 255, 255)

bullet_points = [
    "Nature",
    "History",
    "Greek and Roman Mythology",
    "Religion",
    "Sacred Oriental Text",
]

for point in bullet_points:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(28)

# Add smaller title
small_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
small_title_frame = small_title_box.text_frame
small_title_frame.text = "The Subject and Content of Art"
small_title_frame.paragraphs[0].font.size = Pt(20)
small_title_frame.paragraphs[0].bold = True
small_title_box.fill.solid()
small_title_box.fill.fore_color.rgb = RGBColor(255, 255, 0)
small_title_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)

# Save the presentation
presentation.save("render.pptx")