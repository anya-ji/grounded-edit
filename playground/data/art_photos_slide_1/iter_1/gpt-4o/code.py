from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Initialize presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5]  # Use a blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set background image
background_path = "../../media/image_0.jpg"
slide.shapes.add_picture(background_path, 0, 0, presentation.slide_width, presentation.slide_height)

# Function to add a text box
def add_text_box(slide, text, top, height, font_size, color=RGBColor(0, 0, 0)):
    left = Inches(0)
    width = presentation.slide_width
    text_box = slide.shapes.add_textbox(left, top, width, height)

    # Set the text
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = color  # Updated to allow custom color
    p.alignment = PP_ALIGN.CENTER

# Add white rectangular bar behind the title
def add_background_bar(slide, top, height, width):
    left = Inches(0)
    slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = RGBColor(255, 255, 255)  # White color

# Add white bar for text contrast
add_background_bar(slide, Inches(2), Inches(2), presentation.slide_width)

# Add title
add_text_box(slide, "GE 7:", Inches(2), Inches(1), 60)

# Add subtitle with white color
add_text_box(slide, "ART APPRECIATION", Inches(3), Inches(1), 48, RGBColor(255, 255, 255))  # White color

# Save presentation
presentation.save("render.pptx")