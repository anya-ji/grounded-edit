from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.shapes.shapetree import SlideShapes

# Initialize presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5]  # Use a blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set background image
background_path = "../../media/image_0.jpg"
slide.shapes.add_picture(background_path, 0, 0, presentation.slide_width, presentation.slide_height)

# Function to add a text box
def add_text_box(slide, text, top, height, font_size):
    left = Inches(0)
    width = presentation.slide_width
    text_box = slide.shapes.add_textbox(left, top, width, height)

    # Set the text
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black
    p.alignment = PP_ALIGN.CENTER

# Add title with updated font size
add_text_box(slide, "GE 7:", Inches(2), Inches(1), 48)  # Changed font size to 48

# Add subtitle
add_text_box(slide, "ART APPRECIATION", Inches(3), Inches(1), 48)

# Save presentation
presentation.save("render.pptx")