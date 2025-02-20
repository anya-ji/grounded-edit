from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide
slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add title
title_text = "Art and Everyday Life"
text_box = slide.shapes.add_textbox(Inches(0), Inches(0.2), Inches(16), Inches(1))
text_frame = text_box.text_frame
text_frame.text = title_text
title_paragraph = text_frame.paragraphs[0]
title_paragraph.font.size = Pt(48)  # Font size change for better visibility
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(50, 50, 50)  # Neutral color

# Paths to images
image_paths = [
    "../../media/image_2.jpg",  # Sculpting
    "../../media/image_0.jpg",  # Painting
    "../../media/image_4.jpg",  # Field
    "../../media/image_1.jpg",  # Market
    "../../media/image_3.jpg"   # Shopping mall
]

# Coordinates for image positions
positions = [
    (Inches(0.5), Inches(1.5)),  # Top left
    (Inches(11.5), Inches(1.5)), # Top right
    (Inches(6), Inches(3.5)),    # Center
    (Inches(0.5), Inches(6)),    # Bottom left
    (Inches(11.5), Inches(6))    # Bottom right
]

# Add images in circular frames
for img_path, pos in zip(image_paths, positions):
    left, top = pos
    picture = slide.shapes.add_picture(img_path, left, top, width=Inches(4), height=Inches(4))
    
    # Making the images circular
    picture.crop_top = 0.0
    picture.crop_bottom = 0.0
    picture.crop_left = 0.0
    picture.crop_right = 0.0
    
# Save presentation
presentation.save("render.pptx")