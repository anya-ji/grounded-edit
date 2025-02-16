from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

# Create presentation and slide
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(240, 248, 255)  # Light background color

# Add decorative elements
left_top_shape = slide.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(0), Inches(0), Inches(2), Inches(2)
)
left_top_shape.fill.solid()
left_top_shape.fill.fore_color.rgb = RGBColor(255, 182, 193)  # Light pink

right_bottom_shape = slide.shapes.add_shape(
    MSO_SHAPE.OVAL, Inches(14), Inches(7), Inches(2), Inches(2)
)
right_bottom_shape.fill.solid()
right_bottom_shape.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue

# Add logo in top left corner
logo_path = "../../media/fluentize_logo.jpg"  # replace with actual path
slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.5), Inches(2), Inches(1))

# Add title text
title_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8), Inches(2))
title_text_frame = title_text_box.text_frame
title_text = title_text_frame.add_paragraph()
title_text.text = "Friends | Joey Doesn't Share Food! (B1)"
title_text.font.bold = True
title_text.font.size = Pt(44)
title_text.font.color.rgb = RGBColor(0, 0, 139)  # Dark blue

# Add image in a circular shape
image_path = "../../media/image_0.jpg"
image_left = Inches(11)
image_top = Inches(2)
image_size = Inches(3.5)
mask_circle = slide.shapes.add_shape(
    MSO_AUTO_SHAPE_TYPE.OVAL, image_left, image_top, image_size, image_size
)
pic = slide.shapes.add_picture(image_path, image_left, image_top, image_size, image_size)
mask_circle.fill.solid()
pic.left, pic.top, pic.width, pic.height = image_left, image_top, image_size, image_size  # Align image precisely

# Add footer text
footer_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(8), Inches(10), Inches(1))
footer_text_frame = footer_text_box.text_frame
footer_text = footer_text_frame.add_paragraph()
footer_text.text = "Fluentize, LLC. Copyright 2023. For use only with license."
footer_text.font.size = Pt(12)

# Save presentation
presentation.save("render.pptx")