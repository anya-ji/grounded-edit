from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add title text
text_box = slide.shapes.add_textbox(Inches(0), Inches(0.2), Inches(16), Inches(1))  # Adjusted position
text_frame = text_box.text_frame
text_frame.clear()  # Clear any existing content

p = text_frame.add_paragraph()
p.text = "Innovations (Future Products or Services)"
p.font.bold = True  # Changed to bold font
p.font.size = Pt(24)  # Adjusted font size to 24
p.font.color.rgb = RGBColor(255, 165, 0)  # Kept font color as bold orange
p.alignment = 1  # Center aligned

# Save the presentation
presentation.save("render.pptx")