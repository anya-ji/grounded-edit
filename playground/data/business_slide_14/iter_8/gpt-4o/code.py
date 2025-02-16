from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Removed the solid orange rectangle for contrast

# Add title text
text_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(15), Inches(1))
text_frame = text_box.text_frame
text_frame.clear()  # Clear any existing content

p = text_frame.add_paragraph()
p.text = "Innovations (Future Products or Services)"
p.font.bold = False  # Changed to non-bold font for less intrusion
p.font.size = Pt(24)  # Adjusted font size to 24 for a less intrusive appearance
p.font.color.rgb = RGBColor(255, 165, 0)  # Kept font color as bold orange

# Save the presentation
presentation.save("render.pptx")