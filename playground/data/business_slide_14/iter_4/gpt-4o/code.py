from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add a more vivid orange rectangle for contrast
left_inch = 0.5
top_inch = 0.2
width_inch = 15
height_inch = 1.5
shape = slide.shapes.add_shape(
    5, Inches(left_inch), Inches(top_inch), Inches(width_inch), Inches(height_inch)
)
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(255, 165, 0)  # Changed to a vivid orange color

# Add title text
text_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.4), Inches(15), Inches(1))
text_frame = text_box.text_frame
text_frame.clear()  # Clear any existing content

p = text_frame.add_paragraph()
p.text = "Innovations (Future Products or Services)"
p.font.bold = False  # Change to less bold font
p.font.size = Pt(36)  # Adjust size to be lighter
p.font.color.rgb = RGBColor(255, 255, 255)  # Changed font color to white

# Save the presentation
presentation.save("render.pptx")