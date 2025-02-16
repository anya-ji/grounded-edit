from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue color

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(3), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "PART 1"
title.font.bold = True
title.font.size = Pt(32)

# Add instructions
instructions_box = slide.shapes.add_textbox(Inches(0.5), Inches(1), Inches(10), Inches(1.5))
instructions_frame = instructions_box.text_frame
instructions = instructions_frame.add_paragraph()
instructions.text = ("Describe each date you see below. "
                     "Which do you think is the best for a first date? Rank them in order from 1 (best) - 6 (worst). "
                     "Explain your choices.")
instructions.font.size = Pt(16)

# Define image coordinates and size
x_offset = 0.5
y_offset = 2.5
image_width = 2
image_height = 1.5
images = ["../../media/image_0.jpg"] * 6

# Add images A to F
labels = ["A", "B", "C", "D", "E", "F"]
coords = [(x_offset, y_offset), (x_offset + 2.5, y_offset), (x_offset + 5, y_offset),
          (x_offset, y_offset + 2), (x_offset + 2.5, y_offset + 2), (x_offset + 5, y_offset + 2)]

for i, coord in enumerate(coords):
    img_path = images[i]
    slide.shapes.add_picture(img_path, Inches(coord[0]), Inches(coord[1]), Inches(image_width), Inches(image_height))
    label_box = slide.shapes.add_textbox(Inches(coord[0] + 0.8), Inches(coord[1] - 0.3), Inches(0.2), Inches(0.3))
    label_frame = label_box.text_frame
    label = label_frame.add_paragraph()
    label.text = labels[i]
    label.font.bold = True
    label.font.size = Pt(18)

# Add ranking section
ranking_box = slide.shapes.add_textbox(Inches(8.5), Inches(2.5), Inches(7), Inches(3))
ranking_frame = ranking_box.text_frame
for i in range(1, 7):
    line = ranking_frame.add_paragraph()
    line.text = f"{i}."

ranking_frame.word_wrap = True
ranking_frame.auto_size = True

# Add note at bottom
note_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(14), Inches(0.5))
note_frame = note_box.text_frame
note = note_frame.add_paragraph()
note.text = "(1 = the best first date idea / 6 = the worst first date idea)"
note.font.size = Pt(12)

# Add "PREVIEW ACTIVITY" button
button = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(12.5), Inches(7.8), Inches(3), Inches(0.8))
button.fill.solid()
button.fill.fore_color.rgb = RGBColor(0, 102, 204)  # Dark blue color
button_line = button.line
button_line.color.rgb = RGBColor(255, 255, 255)
button_text_frame = button.text_frame
button_text = button_text_frame.add_paragraph()
button_text.text = "PREVIEW ACTIVITY"
button_text.font.size = Pt(16)
button_text.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Save the presentation
presentation.save("render.pptx")