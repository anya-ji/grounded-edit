from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Define colors
title_color = RGBColor(0, 102, 204)  # Dark blue for title
text_color = RGBColor(0, 0, 0)  # Black for general text
button_color = RGBColor(255, 140, 0)  # Orange for buttons
option_bg_color = RGBColor(240, 240, 240)  # Light gray for options

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(5), Inches(1))
title_frame = title_box.text_frame
title_txt = title_frame.add_paragraph()
title_txt.text = "PART 3"
title_txt.font.size = Pt(36)
title_txt.font.bold = True
title_txt.font.color.rgb = title_color

# Add question with icon
icon_path = "examples/food/slide_22/media/image_0.jpg"  # Use provided checkmark icon
slide.shapes.add_picture(icon_path, Inches(0.5), Inches(1.5), Inches(0.5), Inches(0.5))

question_box = slide.shapes.add_textbox(Inches(1.1), Inches(1.5), Inches(14), Inches(1))
question_frame = question_box.text_frame
question_txt = question_frame.add_paragraph()
question_txt.text = "What do you think is the general rule for this pronunciation scheme?"
question_txt.font.size = Pt(24)
question_txt.font.color.rgb = text_color

# Add options with checkboxes
options = [
    "a) We add stress/intonation on the time or location we’re speaking about.",
    "b) We add stress/intonation on the present participle (-ing verb).",
    "c) We add stress/intonation on the last word of the clause when retelling events."
]

top = Inches(2.5)

for option in options:
    # Add checkbox
    checkbox = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(0.4), Inches(0.4))
    checkbox.fill.solid()
    checkbox.fill.fore_color.rgb = option_bg_color  # Set checkbox background color

    # Add option text
    option_box = slide.shapes.add_textbox(Inches(1.1), top, Inches(14), Inches(0.5))
    option_frame = option_box.text_frame
    option_txt = option_frame.add_paragraph()
    option_txt.text = option
    option_txt.font.size = Pt(20)
    option_txt.font.color.rgb = text_color

    top += Inches(0.7)

# Add buttons
button1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11), Inches(7.5), Inches(3), Inches(0.75))
button1.fill.solid()
button1.fill.fore_color.rgb = button_color
button1.text = "Check your answers"
button1.text_frame.paragraphs[0].font.size = Pt(20)
button1.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

button2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(11), Inches(8.3), Inches(3), Inches(0.75))
button2.fill.solid()
button2.fill.fore_color.rgb = RGBColor(0, 102, 204)  # Set a different color for contrast
button2.text = "PRONUNCIATION ACTIVITY"
button2.text_frame.paragraphs[0].font.size = Pt(20)
button2.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

# Save the presentation
presentation.save("render.pptx")