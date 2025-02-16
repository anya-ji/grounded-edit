from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide to the presentation
slide_layout = presentation.slide_layouts[5]  # Blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Set background to white color
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

# Add centered text
text = "Thank You!"
text_box = slide.shapes.add_textbox(Inches(4), Inches(3.5), Inches(8), Inches(2))
text_frame = text_box.text_frame
text_frame.text = text

# Apply text formatting
p = text_frame.paragraphs[0]
p.font.size = Pt(60)
p.font.bold = True
p.font.color.rgb = RGBColor(0, 0, 0)
p.font.name = 'Serif'
# Removed italic font style
# p.font.italic = True  
text_box.text_frame.text_anchor = 'middle'

# Removed the code to add a thin horizontal line

# Save the presentation
presentation.save("render.pptx")