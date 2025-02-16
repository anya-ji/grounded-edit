from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide to the presentation
slide_layout = presentation.slide_layouts[5]  # Blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Set background to white color
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

# Add centered text
text = "Thank You!"
text_box = slide.shapes.add_textbox(Inches(4), Inches(3.5), Inches(8), Inches(2))
text_frame = text_box.text_frame
text_frame.text = text

# Apply text formatting
p = text_frame.paragraphs[0]
p.font.size = Pt(60)
# p.font.bold = True  # Removed font bold setting
p.font.color.rgb = RGBColor(0, 0, 0)
p.font.name = 'Times New Roman'  # Changed font to Times New Roman
text_box.rotation = 0  # Removed tilt for straight alignment

# Save the presentation
presentation.save("render.pptx")