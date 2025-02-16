from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Initialize presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5]  # blank slide
slide = presentation.slides.add_slide(slide_layout)

# Add title
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Lessons that we can learn from Elon Musk's life."
title.font.bold = True
title.font.size = Pt(32)

# Add bullet points
left_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(7), Inches(5))
left_frame = left_box.text_frame
left_frame.word_wrap = True
left_frame.margin_bottom = 0

points = [
    "Work hardly",
    "Divergent thinking",
    "Our first goal should not be to make money.",
    "Never fear to failure",
    "Take a risk.",
    "Follow own passions"
]

for point in points:
    bullet = left_frame.add_paragraph()
    bullet.text = point
    bullet.font.size = Pt(24)
    bullet.level = 0

# Add image with quote
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(9), Inches(1.5), Inches(6), Inches(6.8))

# Set background colors
left_background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(8), Inches(9)
)
left_background.fill.solid()
left_background.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White

right_background = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(8), Inches(0), Inches(8), Inches(9)
)
right_background.fill.solid()
right_background.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow

# Adjust visibility of layers
slide.shapes._spTree.remove(left_background._element)
slide.shapes._spTree.insert(2, left_background._element)

slide.shapes._spTree.remove(right_background._element)
slide.shapes._spTree.insert(1, right_background._element)

# Save presentation
presentation.save("render.pptx")