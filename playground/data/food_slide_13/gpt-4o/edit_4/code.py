from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Select a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add slide title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_frame.clear()
title_paragraph = title_frame.add_paragraph()
title_paragraph.text = "PART 4"
title_paragraph.font.size = Pt(48)
title_paragraph.font.bold = True

# Add instructions
instruction_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(15), Inches(1.5))
instruction_frame = instruction_box.text_frame
instruction_frame.clear()
instruction_paragraph = instruction_frame.add_paragraph()
instruction_paragraph.text = "List the events in Joey’s second date with Sarah in order from 2 - 6. The first one has been done as an example."
instruction_paragraph.font.size = Pt(20)

# Add checklist
checklist_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(4))
checklist_frame = checklist_box.text_frame

events = [
    "1. Sarah leaves the table for a phone call.",
    "2. Joey doesn’t like the dessert he orders and tries to order a different one.",
    "3. Joey refuses to apologize when Sarah returns.",
    "4. Joey apologizes for getting angry that Sarah took his food on the first date.",
    "5. Joey is unable to control himself and eats all of Sarah’s dessert.",
    "6. The waiter arrives with their desserts."
]

for event in events:
    paragraph = checklist_frame.add_paragraph()
    paragraph.text = event
    paragraph.space_after = Pt(20)  # Increased spacing
    paragraph.font.size = Pt(18)

# Add image of Joey
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(12), Inches(6.5), width=Inches(3))

# Add buttons
# Check your answers button
button_1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10), Inches(5), Inches(2), Inches(0.6))
button_1.fill.solid()
button_1.fill.fore_color.rgb = RGBColor(255, 255, 255)  # Transparent background
button_1.line.color.rgb = RGBColor(255, 192, 203)  # Pink border
button_1.text = "Check your answers"
button_1.text_frame.paragraphs[0].font.size = Pt(14)
button_1.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 105, 180)  # Pink text

# VIEWING ACTIVITY button
button_2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(10), Inches(6), Inches(2), Inches(0.6))
button_2.fill.solid()
button_2.fill.fore_color.rgb = RGBColor(0, 0, 139)  # Dark blue background
button_2.line.color.rgb = RGBColor(0, 0, 139)  # Dark blue border
button_2.text = "VIEWING ACTIVITY"
button_2.text_frame.paragraphs[0].font.size = Pt(14)
button_2.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)  # White text

# Save the presentation
presentation.save("render.pptx")