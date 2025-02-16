from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Choose a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set a darker background color for the slide
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)  # Keeping the background black

# Add slide title
title = slide.shapes.title
title.text = "IMPORTANCE OF SKILLS BEYOND ACADEMICS"
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.size = Pt(44)
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(208, 208, 208)  # Changed to Light Gray

# Creating a textbox for bullet points
left_margin = Inches(1)
top_margin = Inches(1.5)
width = Inches(14)
height = Inches(6)
textbox = slide.shapes.add_textbox(left_margin, top_margin, width, height)
text_frame = textbox.text_frame

# Define bullet points
bullet_points = [
    ("Financial Literacy", [
        "Skills include budgeting, saving, investing, and making informed financial decisions."
    ]),
    ("Leadership and Teamwork", [
        "Leadership skills involve motivating and guiding others toward common goals.",
        "Teamwork skills enable collaboration and the ability to work effectively in diverse groups."
    ]),
    ("Time Management", [
        "Effective time management and organizational skills improve productivity and reduce stress."
    ]),
    ("Creativity and Innovation", [
        "These skills lead to the development of new ideas, products, and solutions."
    ]),
    ("Cultural Competence and Diversity", [
        "In an increasingly globalized world, understanding and respecting cultural differences is essential.",
        "Cultural competence promotes inclusivity and effective interactions in multicultural environments."
    ])
]

# Define colors for bullet points
highlight_color = RGBColor(208, 208, 208)  # Changed to Light Gray
text_color = RGBColor(208, 208, 208)  # Default color for subpoints

# Add bullet points to slide
for heading, subpoints in bullet_points:
    # Add heading
    p = text_frame.add_paragraph()
    p.text = heading
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = highlight_color
    p.space_before = Pt(10)
    
    # Add subpoints
    for subpoint in subpoints:
        p = text_frame.add_paragraph()
        p.text = subpoint
        p.level = 1
        p.font.size = Pt(20)
        # Change color of subpoints for Leadership and Teamwork to gray
        if heading == "Leadership and Teamwork":
            p.font.color.rgb = text_color
        else:
            p.font.color.rgb = text_color  # Default color

# Save the presentation
presentation.save("render.pptx")