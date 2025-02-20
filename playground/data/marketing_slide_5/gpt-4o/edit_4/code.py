from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add image to the left side as background
image_path = "../../media/image_0.jpg"
left = top = 0
pic = slide.shapes.add_picture(image_path, left, top, Inches(8), presentation.slide_height)

# Add a semi-transparent overlay rectangle on the image
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, left, top, pic.width, presentation.slide_height
)
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)
fill.fore_color.opacity = 0.4  # 60% transparency
shape.line.color.rgb = RGBColor(0, 0, 0)  # No border

# Create a text box for the title and points on the right side
text_box = slide.shapes.add_textbox(Inches(8.2), Inches(1), Inches(7.6), Inches(6))
text_frame = text_box.text_frame
text_frame.word_wrap = True

# Title
p = text_frame.add_paragraph()
p.text = "ELO's"
p.font.bold = True
p.font.size = Pt(48)
p.font.color.rgb = RGBColor(255, 255, 255)  # White color

# First point
p = text_frame.add_paragraph()
p.text = "1. What is market?"
p.font.size = Pt(32)
p.space_after = Pt(20)  # Add some space after paragraph
p.font.color.rgb = RGBColor(255, 255, 255)  # White color

# Second point
p = text_frame.add_paragraph()
p.text = "2. How product reach to market?"
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(255, 255, 255)  # White color

# Add footer with light gray background
footer_left = Inches(0)
footer_top = Inches(7.8)  # Position it near the bottom
footer_width = Inches(16)
footer_height = Inches(1.2)
footer_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, footer_left, footer_top, footer_width, footer_height
)
footer_fill = footer_shape.fill
footer_fill.solid()
footer_fill.fore_color.rgb = RGBColor(211, 211, 211)  # Light gray

# Add text to the footer
footer_text_box = slide.shapes.add_textbox(footer_left, footer_top, footer_width, footer_height)
footer_text_frame = footer_text_box.text_frame
footer_p = footer_text_frame.add_paragraph()
footer_p.text = "This presentation uses images under CC BY-SA license"
footer_p.font.size = Pt(12)

# Save the presentation
presentation.save("render.pptx")