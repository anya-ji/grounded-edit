from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add the background color as a subtle gradient
gradient_fill = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(9)
)
gradient_fill.fill.gradient()
gradient_fill.fill.gradient_stops[0].color.rgb = RGBColor(0, 102, 204)  # Light blue
gradient_fill.fill.gradient_stops[1].color.rgb = RGBColor(255, 255, 255)  # White
gradient_fill.fill.gradient_angle = 45

# Add the background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add a dark overlay
overlay = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(9)
)
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
overlay.fill.transparency = 0.5
overlay.line.color.rgb = RGBColor(0, 0, 0)

# Add the text box with red background
left = Inches(4)
top = Inches(3)
width = Inches(8)
height = Inches(3)
text_box = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, left, top, width, height
)
text_box.fill.solid()
text_box.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red background
text_box.text = "Which type of market is beneficial for the customers? Why?"
text_box.text_frame.text = "Which type of market is beneficial for the customers? Why?"

# Format the text
text_frame = text_box.text_frame
for paragraph in text_frame.paragraphs:
    paragraph.font.size = Pt(28)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Save the presentation
presentation.save("render.pptx")