from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add the grocery store background image
background_image_path = "../../media/grocery_store.jpg"  # Changed to grocery store image
slide.shapes.add_picture(background_image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add a dark overlay if needed (optional, depending on visual preference)
overlay = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(9)
)
overlay.fill.solid()
overlay.fill.fore_color.rgb = RGBColor(0, 0, 0)
overlay.fill.transparency = 0.3  # Adjust transparency for better visibility of the background
overlay.line.color.rgb = RGBColor(0, 0, 0)

# Add the text box with red background and rounded corners
left = Inches(4)
top = Inches(3)
width = Inches(8)
height = Inches(3)
text_box = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
)
text_box.fill.solid()
text_box.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red background
text_box.text = "Which type of market is beneficial for the customers? Why?"

# Format the text
text_frame = text_box.text_frame
for paragraph in text_frame.paragraphs:
    paragraph.font.size = Pt(28)
    paragraph.font.bold = True
    paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Add shadow effect to the text box
text_box.shadow.visible = True
text_box.shadow.blur_radius = Pt(5)
text_box.shadow.distance = Pt(3)
text_box.shadow.angle = 45

# Save the presentation
presentation.save("render.pptx")