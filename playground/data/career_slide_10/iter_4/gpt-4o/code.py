from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color to dark
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(20, 20, 20)

# Add title with proper centering
title_box = slide.shapes.add_textbox(Inches(0), Inches(0.5), Inches(16), Inches(1))
title = title_box.text_frame.add_paragraph()
title.text = "YOU DON'T NEED A DEGREE TO BE SUCCESSFUL"
title.font.bold = True
title.font.size = Pt(44)
title.font.color.rgb = RGBColor(255, 255, 255)
title_box.text_frame.text = title.text  # Clear previous text and add

# Center align the title
title_box.text_frame.paragraphs[0].alignment = 1  # 1 = center alignment

# List of image paths and descriptions
people_info = [
    ("../../media/image_0.jpg", "Richard Branson", 
     "Business magnate and commercial astronaut. Founder of the Virgin Group. Left school at 16."),
    ("../../media/image_2.jpg", "Mark Twain", 
     "Greatest American writer, humorist, entrepreneur, publisher, and lecturer. Did not have formal education beyond elementary school."),
    ("../../media/image_1.jpg", "Steve Jobs", 
     "Co-founder of Apple Inc., pioneer of the microcomputer revolution. Attended Reed College in Portland, Oregon, for one semester before dropping out."),
    ("../../media/image_3.jpg", "Oprah Winfrey", 
     "Famous media mogul, talk show host, actress, and philanthropist. Attended Tennessee State University but did not complete her college education.")
]

# Add images and descriptions
spacing = 0.5  # Set consistent spacing
total_width = 2 * len(people_info) + spacing * (len(people_info) - 1)
left_inch = (presentation.slide_width.inches - total_width) / 2  # Center the images horizontally
top_inch = 2

for image_path, name, description in people_info:
    # Add circular pictures
    pic = slide.shapes.add_picture(image_path, Inches(left_inch), Inches(top_inch), Inches(2), Inches(2.5))
    
    # Update the top position for description to be below the picture
    desc_top_inch = top_inch + 2.5 + Inches(0.2)  # Space below the picture

    left_inch += 2 + spacing  # Add picture width and spacing
    
    # Add name and description
    text_box = slide.shapes.add_textbox(pic.left, desc_top_inch, Inches(2), Inches(1))  # Changed width to 2 for uniformity
    text_frame = text_box.text_frame

    name_paragraph = text_frame.add_paragraph()
    name_paragraph.text = name
    name_paragraph.font.size = Pt(20)
    name_paragraph.font.bold = True
    name_paragraph.font.color.rgb = RGBColor(255, 255, 0)

    desc_paragraph = text_frame.add_paragraph()
    desc_paragraph.text = description
    desc_paragraph.font.size = Pt(14)
    desc_paragraph.font.color.rgb = RGBColor(255, 255, 255)

# Save presentation
presentation.save("render.pptx")