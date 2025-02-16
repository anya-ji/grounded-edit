from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background color to coral
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 102, 102)

# Add the title text
text_box = slide.shapes.add_textbox(Inches(0), Inches(3.5), width=Inches(16), height=Inches(2))
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = "This is Bélo."
p.font.bold = True
p.font.size = Pt(60)
p.font.color.rgb = RGBColor(255, 255, 255)

# Set a rounder, softer font style
p.font.name = 'Arial Rounded MT Bold'  # Example of a rounder font
p.alignment = PP_ALIGN.CENTER

# Save the presentation
presentation.save("render.pptx")