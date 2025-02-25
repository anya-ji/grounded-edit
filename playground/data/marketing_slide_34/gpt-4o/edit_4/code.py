from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation and set slide dimensions
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the slide background to a gradient dark green
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 51, 25)  # Dark green color

# Add title text to the slide
title_text = "Let's Recall:"
title_box = slide.shapes.add_textbox(Inches(3), Inches(1), Inches(10), Inches(2))
title_frame = title_box.text_frame
title_frame.word_wrap = True
title_p = title_frame.add_paragraph()
title_p.text = title_text
title_p.font.size = Pt(54)
title_p.font.color.rgb = RGBColor(255, 255, 255)  # White color
title_p.font.bold = True
title_box.text_frame.paragraphs[0].alignment = 1  # Center alignment

# Add subtitle text to the slide
subtitle_text = "Differentiate between Wholesaler and Retailers"
# Adjusted position to add margin at bottom
subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(3.5), Inches(12), Inches(2))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.word_wrap = True
subtitle_p = subtitle_frame.add_paragraph()
subtitle_p.text = subtitle_text
subtitle_p.font.size = Pt(32)
subtitle_p.font.color.rgb = RGBColor(255, 255, 255)  # White color
subtitle_p.font.bold = True
subtitle_box.text_frame.paragraphs[0].alignment = 1  # Center alignment

# Add a small red rectangle in the top right corner for visual interest
red_rect = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(14), Inches(0.5), Inches(1), Inches(0.5))
red_rect.fill.solid()
red_rect.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red color

# Save the presentation
presentation.save("render.pptx")