from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()

# Set dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide
slide_layout = presentation.slide_layouts[5]  # Using a blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set background to a gradient of dark teal to lighter teal
background = slide.background
fill = background.fill
fill.gradient()  # Set to gradient fill
fill.gradient_stops[0].color.rgb = RGBColor(0, 50, 50)  # Dark teal
fill.gradient_stops[1].color.rgb = RGBColor(0, 255, 255)  # Lighter teal

# Add title text
title_text = "Research Work:"
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5))
title_tf = title_box.text_frame
title_p = title_tf.add_paragraph()
title_p.text = title_text
title_p.font.bold = True
title_p.font.size = Pt(44)
title_p.font.color.rgb = RGBColor(255, 255, 0)  # Yellow text

# Add content text
content_text = "Find out how people buy and sell products over the internet. Are the goods sold at cheaper or higher rates than other market?"
content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(2))
content_tf = content_box.text_frame
content_p = content_tf.add_paragraph()
content_p.text = content_text
content_p.font.size = Pt(24)
content_p.font.color.rgb = RGBColor(255, 255, 0)  # Yellow text
content_p.alignment = None  # Left alignment

# Save the presentation
presentation.save("render.pptx")