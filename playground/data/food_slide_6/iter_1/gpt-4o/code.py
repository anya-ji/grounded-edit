from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

presentation = Presentation()
slide_layout = presentation.slide_layouts[5]  # Blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Add title button
title_button = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.2), Inches(2.5), Inches(0.5)
)
title_button.fill.solid()
title_button.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue color
title_button.text_frame.text = "PREVIEW DISCUSSION"
title_text = title_button.text_frame.paragraphs[0]
title_text.font.size = Pt(14)
title_text.font.bold = True

# Add main question
main_question = slide.shapes.add_textbox(
    Inches(3), Inches(2.5), Inches(10), Inches(2)
)
main_question_text = main_question.text_frame.add_paragraph()
main_question_text.text = "What do you think are some good ways to make a good first impression on a date?"
main_question_text.font.size = Pt(24)
main_question_text.font.bold = True
main_question_text.font.color.rgb = RGBColor(0, 0, 139)  # Dark blue

# Change thought bubble graphic to a speech bubble on the left with vibrant colors
speech_bubble = slide.shapes.add_shape(
    MSO_SHAPE.CLOUD_CALLOUT, Inches(1), Inches(2.5), Inches(1.5), Inches(1)
)
speech_bubble.fill.solid()
speech_bubble.fill.fore_color.rgb = RGBColor(255, 230, 230)  # Light pink color
speech_bubble.text_frame.text = "?"
speech_bubble_text = speech_bubble.text_frame.paragraphs[0]
speech_bubble_text.font.size = Pt(20)
speech_bubble_text.font.bold = True
speech_bubble_text.font.color.rgb = RGBColor(255, 105, 180)  # Hot pink

# Add illustration of two people
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(10), Inches(1.5), Inches(4), Inches(4))

# Add buttons at the bottom
# Play button icon (using a rounded rectangle as a placeholder)
play_button = slide.shapes.add_shape(
    MSO_SHAPE.PENTAGON, Inches(0.5), Inches(6), Inches(1), Inches(0.5)
)
play_button.fill.solid()
play_button.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue
play_button.text_frame.text = "▶"

# Viewing activity button
view_activity_button = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(11.5), Inches(6), Inches(3), Inches(0.5)
)
view_activity_button.fill.solid()
view_activity_button.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue
view_activity_button.text_frame.text = "VIEWING ACTIVITY"
activity_text = view_activity_button.text_frame.paragraphs[0]
activity_text.font.size = Pt(14)

# Save the presentation
presentation.save("render.pptx")