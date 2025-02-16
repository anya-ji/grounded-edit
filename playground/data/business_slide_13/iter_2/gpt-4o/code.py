from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image, ImageFilter

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Load and blur the background image
background_image_path = "../../media/image_0.jpg"
background_image = Image.open(background_image_path)
blurred_image = background_image.filter(ImageFilter.GaussianBlur(radius=5))  # Adjust the radius as needed
blurred_image_path = "blurred_image.jpg"
blurred_image.save(blurred_image_path)

# Set background image
slide.shapes.add_picture(blurred_image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Save the presentation
presentation.save("render.pptx")