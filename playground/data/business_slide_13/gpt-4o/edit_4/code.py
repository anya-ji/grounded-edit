from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color to light green
background_color = RGBColor(209, 232, 185)  # #D1E8B9
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = background_color

# Add a title
title = "Innovations (Future Products or Services)"
left = Inches(0.5)
top = Inches(0.5)
width = Inches(15)
height = Inches(1.5)

text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = title
p.font.bold = True
p.font.size = Pt(44)
p.alignment = PP_ALIGN.LEFT

# Set font color and outline for contrast
p.font.color.rgb = RGBColor(255, 255, 255)  # White font
p.shadow = True  # Add shadow for visual contrast

# Save the presentation
presentation.save("render.pptx")