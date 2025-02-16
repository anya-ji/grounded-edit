from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background image
background_image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(background_image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Save the presentation
presentation.save("render.pptx")