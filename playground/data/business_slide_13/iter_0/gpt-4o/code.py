from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background image
background_image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(background_image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Add a title
title = "Innovations (Future Products or Services)"
left = Inches(0.5)
top = Inches(0.5)
width = Inches(15)
height = Inches(1.5)

text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = title
p.font.bold = False  # Change font weight to lighter
p.font.size = Pt(44)
p.alignment = PP_ALIGN.LEFT

# Set font color and outline for contrast
p.font.color.rgb = RGBColor(255, 255, 255)  # White font
p.shadow = True  # Add shadow for visual contrast

# Save the presentation
presentation.save("render.pptx")