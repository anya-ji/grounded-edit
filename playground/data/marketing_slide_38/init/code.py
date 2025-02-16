from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.shapes.auto_shape import Shape

# Helper function to convert RGB to hex for gradients
def rgb_to_hex(rgb_color):
    return '{:02x}{:02x}{:02x}'.format(rgb_color[0], rgb_color[1], rgb_color[2])

# Initialize presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Create slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add title with gradient background
title_placeholder = slide.shapes.title
title_placeholder.text = "Opportunity Created by Market"
title_placeholder.left = Inches(2.5)
title_placeholder.top = Inches(0.5)
title_placeholder.width = Inches(11)
title_placeholder.height = Inches(1.5)
title_placeholder.text_frame.fit_text(maximum_size=48)
title_text_frame = title_placeholder.text_frame

# Set gradient fill for title placeholder
title_shape: Shape = title_placeholder
fill = title_shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 128, 128)  # Dark teal background
title_bg_frame = title_text_frame.auto_shape
grad_fill = title_bg_frame.fill
grad_fill.gradient()
stop1 = grad_fill.gradient_stops[0]
stop2 = grad_fill.gradient_stops[1]
stop1.position = 0.0
stop1.color.rgb = RGBColor(0, 128, 128)  # Dark teal
stop2.position = 1.0
stop2.color.rgb = RGBColor(175, 238, 238)  # Light teal

# Add icons with corresponding labels
icons_info = [
    ("Factory", "Demand of Raw Material"),
    ("Graph", "Demand for Labour"),
    ("Shop", "More Shops in Market"),
    ("Bus", "Transportations"),
    ("Buildings", "Infrastructure"),
    ("Person", "Employment")
]

# Define position variables
icon_size = Inches(1)
spacing = Inches(0.5)
start_x = (presentation.slide_width - (icon_size * len(icons_info) + spacing * (len(icons_info) - 1))) / 2
start_y = Inches(3)

for i, (icon_name, label) in enumerate(icons_info):
    # Add colored square with rounded corners
    left = start_x + i * (icon_size + spacing)
    top = start_y
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, icon_size, icon_size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(92, 172, 238)  # Light Blue Color
    shape.line.color.rgb = RGBColor(255, 255, 255)
    
    # Add label below the icon
    text_box = slide.shapes.add_textbox(left, top + icon_size + Inches(0.2), icon_size, Inches(0.5))
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = label
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 0, 0)  # Black color for text

# Save presentation
presentation.save("render.pptx")