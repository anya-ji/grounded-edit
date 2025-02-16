from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Select a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add bright blue gradient background rectangle
gradient_background = slide.shapes.add_shape(
    1,  # msoShapeRectangle
    0,  # left
    0,  # top
    Inches(16),  # width
    Inches(9),  # height
)
gradient_background.fill.solid()
gradient_background.fill.fore_color.rgb = RGBColor(100, 149, 237)  # Bright blue background

# Add title with speech bubble graphic
title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
title.fill.solid()  # Add this line to enable filling the background color
title.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Change to dark blue background
title.text_frame.text = "LESSON REFLECTION:"
p = title.text_frame.paragraphs[0]
p.font.bold = True
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(255, 255, 255)  # Change font color to white for contrast

# Change the font style for the title
p.font.name = 'Arial'  # Set to a cleaner, more modern font style

# Add speech bubble image
speech_bubble_image_path = "../../media/image_1.jpg"
slide.shapes.add_picture(speech_bubble_image_path, Inches(0.5), Inches(1.5), height=Inches(1))

# Add the same title at the top of the slide
top_title = slide.shapes.add_textbox(Inches(1), Inches(0), Inches(10), Inches(1))
top_title.fill.solid()  # Add this line to enable filling the background color
top_title.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Change to dark blue background
top_title.text_frame.text = "LESSON REFLECTION:"
top_p = top_title.text_frame.paragraphs[0]
top_p.font.bold = True
top_p.font.size = Pt(36)
top_p.font.color.rgb = RGBColor(255, 255, 255)  # Change font color to white for contrast
top_p.font.name = 'Arial'  # Set to a cleaner, more modern font style

# Add questions within speech bubbles
questions = [
    "1. Did you learn or use any new vocabulary, grammar, phrases, etc. in this lesson?",
    "2. What did you enjoy about the lesson?",
    "3. Is there anything you found difficult about the lesson?"
]

# Specify position for questions
for i, question in enumerate(questions):
    left = Inches(0.5) 
    top = Inches(1 + i * 1.5)  # Adjust vertical spacing
    question_box = slide.shapes.add_textbox(left, top, Inches(6), Inches(1.2))
    question_box.fill.solid()
    question_box.fill.fore_color.rgb = RGBColor(255, 229, 204)  # Light contrasting color for bubble
    question_box.shadow.visible = True  # Add shadow for better definition
    question_box.rounding = 0.2  # Add rounded corners effect
    question_p = question_box.text_frame.add_paragraph()
    question_p.text = question
    question_p.font.size = Pt(20)
    question_p.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue font
    question_p.alignment = PP_ALIGN.LEFT

# Add blackboard image with label
blackboard_image_path = "../../media/image_0.jpg"
pic = slide.shapes.add_picture(blackboard_image_path, Inches(11.5), Inches(5.5), height=Inches(3))

# Add label below the blackboard image
label_box = slide.shapes.add_textbox(Inches(11.5), Inches(8.5), Inches(3), Inches(0.5))
label_p = label_box.text_frame.add_paragraph()
label_p.text = "QUIZ & REVIEW ACTIVITY"
label_p.font.bold = True
label_p.font.size = Pt(14)
label_p.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue font

# Add a thin dark border around the main content area
border_left = slide.shapes.add_shape(
    1,  # msoShapeRectangle
    Inches(0.5), Inches(0.5), Inches(15), Inches(8), 
)

border_left.line.color.rgb = RGBColor(0, 0, 0)  # Black border
border_left.line.width = Pt(1)  # Thin border

# Save the presentation
presentation.save("render.pptx")