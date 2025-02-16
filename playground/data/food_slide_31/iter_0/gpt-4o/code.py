from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Select a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 204, 204)  # Softer vibrant background

# Add title with speech bubble graphic
title = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(10), Inches(1))
title.text_frame.text = "LESSON REFLECTION:"
p = title.text_frame.paragraphs[0]
p.font.bold = True
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue font

# Add speech bubble image
speech_bubble_image_path = "../../media/image_1.jpg"
slide.shapes.add_picture(speech_bubble_image_path, Inches(0.5), Inches(1.5), height=Inches(1))

# Add questions
questions = [
    "1. Did you learn or use any new vocabulary, grammar, phrases, etc. in this lesson? If so, what?",
    "2. What did you enjoy about the lesson?",
    "3. Is there anything you found difficult about the lesson?"
]

left = Inches(1)
top = Inches(2.5)
width = Inches(14)
height = Inches(4)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_frame = text_box.text_frame

for question in questions:
    p = text_frame.add_paragraph()
    p.text = question
    p.space_after = Pt(14)
    p.font.size = Pt(20)
    p.alignment = PP_ALIGN.LEFT

# Add blackboard image with label
blackboard_image_path = "../../media/image_0.jpg"
pic = slide.shapes.add_picture(blackboard_image_path, Inches(11.5), Inches(5.5), height=Inches(3))

# Add label below the blackboard image
label_box = slide.shapes.add_textbox(Inches(11.5), Inches(8.5), Inches(3), Inches(0.5))
label_p = label_box.text_frame.add_paragraph()
label_p.text = "QUIZ & REVIEW ACTIVITY"
label_p.font.bold = True
label_p.font.size = Pt(14)
label_p.font.color.rgb = RGBColor(0, 51, 102)  # Dark blue font

# Save the presentation
presentation.save("render.pptx")