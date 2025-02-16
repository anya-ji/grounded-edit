from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from PIL import Image, ImageEnhance

# Create a presentation object
presentation = Presentation()

# Set slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Choose a slide layout and add slide
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Change the background to a solid color (e.g., white)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)  # White color

# Add text to the slide
text = "As an individual, how can we ensure equality in society?"
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(8), Inches(3))
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = text
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 0)  # Yellow
p.alignment = PP_ALIGN.CENTER

# Function to enhance image
def enhance_image(image_path):
    with Image.open(image_path) as img:
        # Increase saturation and contrast
        converter = ImageEnhance.Color(img)
        img = converter.enhance(1.5)  # Increase saturation
        converter = ImageEnhance.Contrast(img)
        img = converter.enhance(1.5)  # Increase contrast
        output_path = image_path.replace(".jpg", "_enhanced.jpg")
        img.save(output_path)  # Save the enhanced image
        return output_path

# Add images to the slide
image_path1 = "../../media/image_1.jpg"
image_path2 = "../../media/image_0.jpg"
enhanced_image1 = enhance_image(image_path1)
enhanced_image2 = enhance_image(image_path2)

top_image = slide.shapes.add_picture(enhanced_image1, Inches(10), Inches(2), width=Inches(5))
bottom_image = slide.shapes.add_picture(enhanced_image2, Inches(10), Inches(5.5), width=Inches(5))

# Add attribution text
attribution_box = slide.shapes.add_textbox(Inches(10), Inches(8.5), Inches(5.5), Inches(0.5))
attribution_frame = attribution_box.text_frame
attribution_p = attribution_frame.add_paragraph()
attribution_p.text = "This Photo by Unknown author is licensed under CC BY-NC."
attribution_p.font.size = Pt(10)
attribution_p.alignment = PP_ALIGN.RIGHT

# Save the presentation
presentation.save("render.pptx")