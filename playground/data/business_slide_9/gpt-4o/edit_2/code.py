from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the slide background color to white
slide.shapes.add_shape(
    1,  # This is the shape type for a rectangle
    0, 0,  # Position (x, y) for the shape
    presentation.slide_width,  # Width of the shape
    presentation.slide_height  # Height of the shape
).fill.solid()  # Solid fill
slide.shapes[-1].fill.fore_color.rgb = RGBColor(255, 255, 255)  # Set color to white

# Add main title text
title = slide.shapes.title
title.text = "Airbnb is all about people and not about the places at all."
title.text_frame.paragraphs[0].font.size = Pt(40)
title.text_frame.paragraphs[0].font.bold = True
title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Change text color to black for visibility
title.top = Inches(0.5)
title.left = Inches(2)

# Add subtitle text
textbox = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(12), Inches(2))
text_frame = textbox.text_frame
p = text_frame.add_paragraph()
p.text = ("The passion of its founders and the love and generosity of people around the world "
          "was changing the culture of travel with places and personal experiences unlike any other. "
          "Belonging the world over was Airbnb’s truth and its differentiator.")
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(0, 0, 0)  # Change text color to black for visibility

# Add the main image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(1), Inches(3), Inches(14), Inches(4))

# Add the quote
quote_textbox = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(12), Inches(1))
quote_text_frame = quote_textbox.text_frame
quote = quote_text_frame.add_paragraph()
quote.text = "IT'S ABOUT THE PEOPLE, NOT THE PLACES"
quote.font.size = Pt(36)
quote.font.bold = True
quote.font.color.rgb = RGBColor(0, 0, 0)  # Change text color to black for visibility

# Save the presentation
presentation.save("render.pptx")