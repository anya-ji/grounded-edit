from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create presentation and set slide dimensions
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide
slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Introductory sentence configuration
intro_text = "The media have a very important impact on cultural globalization in two mutually interdependent ways:"
intro_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1.5))
intro_frame = intro_box.text_frame
intro_p = intro_frame.add_paragraph()
intro_p.text = intro_text
intro_p.font.size = Pt(28)
intro_p.font.bold = True
intro_p.font.color.rgb = RGBColor(0, 51, 102)  # dark blue color
intro_p.alignment = PP_ALIGN.CENTER

# Title configuration
title_text = "The media have a very important impact on cultural globalization in two mutually interdependent ways:"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(15), Inches(1.5))
title_frame = title_box.text_frame
title_p = title_frame.add_paragraph()
title_p.text = title_text
title_p.font.size = Pt(28)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(0, 51, 102)  # dark blue color
title_p.alignment = PP_ALIGN.LEFT  # Change alignment to left

# Main content as a numbered list
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(15), Inches(5))
content_frame = content_box.text_frame
content_frame.word_wrap = True
content_p1 = content_frame.add_paragraph()
content_p1.text = "1. The media provide an extensive transnational transmission of cultural products and;"
content_p1.space_before = Pt(10)  # Add space before this point
content_p1.space_after = Pt(10)  # Add space after this point
content_p1.level = 0
content_p1.font.size = Pt(20)
content_p1.font.color.rgb = RGBColor(0, 0, 0)  # black color
content_p1.alignment = PP_ALIGN.LEFT  # Change alignment to left

content_p2 = content_frame.add_paragraph()
content_p2.text = "2. It contributes to the formation of communicative networks and social structures."
content_p2.space_before = Pt(10)  # Add space before this point
content_p2.space_after = Pt(10)  # Add space after this point
content_p2.level = 0
content_p2.font.size = Pt(20)
content_p2.font.color.rgb = RGBColor(0, 0, 0)  # black color
content_p2.alignment = PP_ALIGN.LEFT  # Change alignment to left

# Additional text
additional_text = "Global media cultures create a continuous cultural exchange."
additional_box = slide.shapes.add_textbox(Inches(0.5), Inches(8), Inches(15), Inches(1))
additional_frame = additional_box.text_frame
additional_p = additional_frame.add_paragraph()
additional_p.text = additional_text
additional_p.font.size = Pt(18)
additional_p.font.color.rgb = RGBColor(0, 0, 0)  # black color
additional_p.alignment = PP_ALIGN.CENTER

# Save presentation
presentation.save("render.pptx")