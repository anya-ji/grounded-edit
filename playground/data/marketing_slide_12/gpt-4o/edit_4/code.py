from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Set the background gradient color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(64, 64, 64)  # Dark gray gradient

# Add title text
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame
p = title_frame.add_paragraph()
p.text = "Now: I, CAN"
p.font.bold = True
p.font.size = Pt(40)
p.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Add a white box for bullet points with thin orange border
bullet_box = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(1), Inches(2), Inches(8), Inches(3)
)
bullet_box.fill.solid()
bullet_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White fill
line = bullet_box.line
line.color.rgb = RGBColor(255, 165, 0)  # Orange border

# Add bullet points
text_frame = bullet_box.text_frame
text_frame.word_wrap = True

bullet1 = text_frame.add_paragraph()
bullet1.text = "Define the term Marketing"
bullet1.font.size = Pt(20)

# Add downward arrow
arrow1 = slide.shapes.add_shape(
    MSO_SHAPE.DOWN_ARROW, Inches(4.5), Inches(3), Inches(0.5), Inches(0.5)
)
arrow1.fill.solid()
arrow1.fill.fore_color.rgb = RGBColor(255, 165, 0)  # Orange arrow

bullet2 = text_frame.add_paragraph()
bullet2.text = "Explain various things involved how products reach to market"
bullet2.font.size = Pt(20)

# Insert the new bullet point
bullet3 = text_frame.add_paragraph()
bullet3.text = "• Explain various things involved how products reach to market"
bullet3.font.size = Pt(20)

# Add another downward arrow
arrow2 = slide.shapes.add_shape(
    MSO_SHAPE.DOWN_ARROW, Inches(4.5), Inches(4), Inches(0.5), Inches(0.5)
)
arrow2.fill.solid()
arrow2.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow arrow

# Insert the textured image on the right
img_path = "../../media/image_0.jpg"
slide.shapes.add_picture(img_path, Inches(9), Inches(1), Inches(6), Inches(7))

# Save the presentation
presentation.save("render.pptx")