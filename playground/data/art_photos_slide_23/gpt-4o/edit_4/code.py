from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, 0, 0, presentation.slide_width, presentation.slide_height)

# Add title textbox
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1.5))
title_tf = title_box.text_frame
title_tf.word_wrap = True
title_p = title_tf.add_paragraph()
title_p.text = "KEEPING ART"
title_p.font.size = Pt(44)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0x00)  # Yellow background

# Add subtitle textbox
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.75), Inches(15), Inches(1))
subtitle_tf = subtitle_box.text_frame
subtitle_tf.word_wrap = True
subtitle_p = subtitle_tf.add_paragraph()
subtitle_p.text = "2. Museums and Private Collections"
subtitle_p.font.size = Pt(32)
subtitle_p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
subtitle_box.fill.solid()
subtitle_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0x00)  # Yellow background

# Add main content textbox
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(3), Inches(15), Inches(5))
content_tf = content_box.text_frame
content_tf.word_wrap = True
content_p = content_tf.add_paragraph()
content_p.text = ("• Museums are the repositories of much of the art in most "
                  "countries and make them available for public viewing through "
                  "either permanent or temporary exhibition.")
content_p.font.size = Pt(24)  # Updated font size from 20 to 24
content_p.font.color.rgb = RGBColor(0x00, 0x00, 0x00)
content_box.fill.solid()
content_box.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)  # White background

# Save presentation
presentation.save("render.pptx")