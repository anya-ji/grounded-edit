from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
slide_layout = presentation.slide_layouts[5]  # Use a blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set slide dimensions to 16:9 ratio
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Set slide background color to white
slide_bg = slide.background
fill = slide_bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)  # Change to white

# Insert the funnel image
funnel_image_path = "../../media/image_0.jpg"
funnel_img = slide.shapes.add_picture(funnel_image_path, Inches(0.5), Inches(2), Inches(7), Inches(5))

# Insert the word cloud image
word_cloud_image_path = "../../media/image_1.jpg"
word_cloud_img = slide.shapes.add_picture(word_cloud_image_path, Inches(8.5), Inches(2), Inches(7), Inches(5))

# Add a caption below the funnel image
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(7), Inches(1))
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = "Credit: Steve Parizi"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(0, 0, 0)  # Change to black font for visibility

# Add "Marketing" with adjusted positioning
marketing_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(2), Inches(1))  # Adjusted position for better alignment
marketing_frame = marketing_box.text_frame
p_marketing = marketing_frame.add_paragraph()
p_marketing.text = "Marketing"
p_marketing.font.size = Pt(24)  # Adjust size as needed
p_marketing.font.color.rgb = RGBColor(255, 0, 0)  # Red color, adjust as needed

# Change position of "Sales" to align with funnel
sales_box_left = slide.shapes.add_textbox(Inches(5.0), Inches(3.0), Inches(2), Inches(1))  # Adjusted position for better alignment
sales_frame_left = sales_box_left.text_frame
p_sales_left = sales_frame_left.add_paragraph()
p_sales_left.text = "Sales"
p_sales_left.font.size = Pt(24)  # Adjust size as needed
p_sales_left.font.color.rgb = RGBColor(0, 0, 255)  # Change color to blue

sales_box_right = slide.shapes.add_textbox(Inches(6.0), Inches(3.0), Inches(2), Inches(1))  # Adjusted position for better alignment
sales_frame_right = sales_box_right.text_frame
p_sales_right = sales_frame_right.add_paragraph()
p_sales_right.text = "Sales"
p_sales_right.font.size = Pt(24)  # Same size as left
p_sales_right.font.color.rgb = RGBColor(0, 0, 255)  # Change color to blue

# Save the presentation
presentation.save("render.pptx")