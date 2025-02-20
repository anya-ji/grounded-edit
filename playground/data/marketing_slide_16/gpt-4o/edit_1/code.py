from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
slide_layout = presentation.slide_layouts[5]  # Use a blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set slide dimensions to 16:9 ratio
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Set slide background color to black
slide_bg = slide.background
fill = slide_bg.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)

# Insert the funnel image
funnel_image_path = "../../media/image_0.jpg"
funnel_img = slide.shapes.add_picture(funnel_image_path, Inches(0.5), Inches(2), Inches(7), Inches(5))

# Insert the word cloud image
word_cloud_image_path = "../../media/image_1.jpg"
word_cloud_img = slide.shapes.add_picture(word_cloud_image_path, Inches(8.5), Inches(2), Inches(7), Inches(5))

# Add a caption below the funnel image
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(7), Inches(1))
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = "Credit: Steve Parizi"
p.font.size = Pt(14)
p.font.color.rgb = RGBColor(255, 255, 255)  # White font color

# Add text for "Sales" with navy blue color
sales_text_box = slide.shapes.add_textbox(Inches(0.0), Inches(3.5), Inches(2), Inches(1))
sales_text_frame = sales_text_box.text_frame
sales_p = sales_text_frame.add_paragraph()
sales_p.text = "Sales"
sales_p.font.size = Pt(24)
sales_p.font.color.rgb = RGBColor(0, 0, 128)  # Navy blue font color

# Save the presentation
presentation.save("render.pptx")