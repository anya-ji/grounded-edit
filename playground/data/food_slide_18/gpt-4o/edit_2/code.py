from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light background color

# Title "PART 2"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
title_tf = title_box.text_frame
title_tf.text = "PART 2"
title_p = title_tf.paragraphs[0]
title_p.font.size = Pt(36)
title_p.font.bold = True

# Add chat icon
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.5), Inches(2), Inches(2))

# Instruction text
instruction_box = slide.shapes.add_textbox(Inches(3), Inches(1.5), Inches(10), Inches(1))
instruction_tf = instruction_box.text_frame
instruction_tf.text = "Discuss the questions."
instruction_p = instruction_tf.paragraphs[0]
instruction_p.font.size = Pt(32)
instruction_p.font.bold = True

# Dialogue text
dialogue_text = (
    "\"JOEY: We were out to dinner. We were getting along, having a really nice time, "
    "I was thinking she was really cool and then, out of nowhere, "
    "(she reached over and took some of my fries from my plate!)\""
    "\n\n\"PHOEBE: So she took some fries, big deal!\""
    "\n\n\"RACHEL: Oh yeah, Joey doesn't share food. I mean, just last week, we were having breakfast, "
    "and...and he had a couple of grapes on his plate...\""
)
dialogue_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(15), Inches(3))
dialogue_tf = dialogue_box.text_frame
dialogue_p = dialogue_tf.add_paragraph()
dialogue_p.text = dialogue_text
dialogue_p.font.size = Pt(20)
dialogue_p.alignment = PP_ALIGN.RIGHT  # Right-align the dialogue text

# Question 1
question1_box = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(15), Inches(1))
question1_tf = question1_box.text_frame
question1_tf.word_wrap = True
question1_p = question1_tf.add_paragraph()
question1_p.text = "1. When is the past continuous used in comparison with the past simple?"
question1_p.font.size = Pt(18)
question1_p.font.bold = True

# Question 2
question2_box = slide.shapes.add_textbox(Inches(0.5), Inches(8), Inches(15), Inches(1))
question2_tf = question2_box.text_frame
question2_tf.word_wrap = True
question2_p = question2_tf.add_paragraph()
question2_p.text = "2. How is the structure for the past simple different than the past continuous?"
question2_p.font.size = Pt(18)
question2_p.font.bold = True

# Button for "VIEWING FOLLOW-UP"
button_box = slide.shapes.add_textbox(Inches(12), Inches(8), Inches(3.5), Inches(0.8))
button_tf = button_box.text_frame
button_green_color = RGBColor(0, 123, 255)  # Button color
button_p = button_tf.add_paragraph()
button_p.text = "VIEWING FOLLOW-UP"
button_p.font.size = Pt(16)
button_p.font.bold = True
button_p.alignment = PP_ALIGN.CENTER

# Save the presentation
presentation.save("render.pptx")