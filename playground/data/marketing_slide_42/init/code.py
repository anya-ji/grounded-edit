from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

presentation = Presentation()
slide_layout = presentation.slide_layouts[5]  # Blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Now: I, Can"
title.font.size = Pt(36)
title.font.bold = True

# Add bullet points
left = Inches(0.5)
top = Inches(1.5)
width = Inches(8)
height = Inches(3)
text_box = slide.shapes.add_textbox(left, top, width, height)
tf = text_box.text_frame
tf.word_wrap = True

bullet_points = [
    "Explain how market create opportunities.",
    "Identify and list the people who earn more profit in Market."
]

for point in bullet_points:
    p = tf.add_paragraph()
    p.text = point
    p.font.size = Pt(24)

# Add gradient background
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, 0, 0, presentation.slide_width // 2, presentation.slide_height
)
shape.fill.solid()
fill = shape.fill
fill.gradient()

grad_stop1 = fill.gradient().stops[0]
grad_stop1.position = 0.0
grad_stop1.color.rgb = RGBColor(0, 100, 0)  # Dark green

grad_stop2 = fill.gradient().stops[1]
grad_stop2.position = 1.0
grad_stop2.color.rgb = RGBColor(144, 238, 144)  # Light green

# Add red rectangle for emphasis
red_rect = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(13.5), Inches(0.5), Inches(2), Inches(1)
)
red_rect.fill.solid()
red_rect.fill.fore_color.rgb = RGBColor(255, 0, 0)

# Add image of two men in a field
img_path = "../../media/image_1.jpg"
slide.shapes.add_picture(img_path, Inches(0.5), presentation.slide_height - Inches(3), width=None, height=Inches(3))

presentation.save("render.pptx")