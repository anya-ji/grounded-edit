from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw

# Create a presentation object and set the dimensions
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]  # Typically a blank slide
slide = presentation.slides.add_slide(slide_layout)

# Load and insert the main background image
bg_image = "../../media/image_0.jpg"
slide.shapes.add_picture(bg_image, 0, 0, Inches(16), Inches(9))

# Open shopping mall image and make it circular
mall_image_path = "../../media/image_1.jpg"
mall_img = Image.open(mall_image_path)
mall_img = mall_img.crop((0, 0, mall_img.size[0], mall_img.size[0]))
mask = Image.new('L', mall_img.size, 0)
draw = ImageDraw.Draw(mask)
draw.ellipse((0, 0) + mall_img.size, fill=255)

# Add the circular mall image
mall_img_path = "../slidesbench/examples/marketing/slide_7/media/circular_mall.png"
mall_img.putalpha(mask)
mall_img.save(mall_img_path)
slide.shapes.add_picture(mall_img_path, Inches(0.5), Inches(1), Inches(4), Inches(4))

# Dim the right background using a semi-transparent rectangle
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(8), 0, Inches(8), Inches(9)
)
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 0)
fill.transparency = 0.3
shape.line.fill.background()

# Open food market image and make it circular
food_img = Image.open(bg_image)
food_img = food_img.crop((0, 0, food_img.size[0], food_img.size[0]))
mask = Image.new('L', food_img.size, 0)
draw = ImageDraw.Draw(mask)
draw.ellipse((0, 0) + food_img.size, fill=255)

# Add the circular food market image
food_img_path = "../slidesbench/examples/marketing/slide_7/media/circular_food.png"
food_img.putalpha(mask)
food_img.save(food_img_path)
slide.shapes.add_picture(food_img_path, Inches(11.5), Inches(5), Inches(3), Inches(3))

# Insert text
text_box = slide.shapes.add_textbox(Inches(9), Inches(2), Inches(6), Inches(1))
p = text_box.text_frame.add_paragraph()
p.text = "Market is a place where products or services are bought or sold."
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(0, 0, 0)  # Changed to black for better contrast

# Add attribution text
attr_text_box = slide.shapes.add_textbox(Inches(0.5), Inches(8), Inches(15), Inches(0.5))
attr_p = attr_text_box.text_frame.add_paragraph()
attr_p.text = "This Photo by Unknown author is licensed under CC BY-SA."
attr_p.font.size = Pt(12)
attr_p.font.color.rgb = RGBColor(255, 255, 255)

# Save the presentation
presentation.save("render.pptx")