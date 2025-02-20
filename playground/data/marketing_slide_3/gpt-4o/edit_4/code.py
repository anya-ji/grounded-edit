from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()

# Set slide width and height
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Create a blank slide layout
slide_layout = presentation.slide_layouts[5]  # blank slide
slide = presentation.slides.add_slide(slide_layout)

# Set the slide background to solid red
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)

# Add a text box with white text centered on the slide
text = "Where do you buy your products from?"
text_box = slide.shapes.add_textbox(Inches(4), Inches(3.5), Inches(8), Inches(2))
text_box.text_frame.word_wrap = True

p = text_box.text_frame.add_paragraph()
p.text = text
p.font.size = Pt(44)
p.font.color.rgb = RGBColor(255, 255, 255)
p.font.name = 'Calibri'  # Clean, modern font

# Center the text box on the slide with a 2 cm bottom margin
bottom_margin = 2 * 28.35  # 2 cm in points (1 cm = 28.35 points)
text_box.top = int((presentation.slide_height - text_box.height - bottom_margin) / 2)

# Save the presentation
presentation.save("render.pptx")