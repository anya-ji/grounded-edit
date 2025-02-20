from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()

# Set slide width and height
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Create a blank slide layout
slide_layout = presentation.slide_layouts[5]  # blank slide
slide = presentation.slides.add_slide(slide_layout)

# Set the slide background to a muted dark red color (maroon)
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(128, 0, 0)  # Maroon color

# Add a text box with white text centered on the slide
text = "Where do you buy your products from?"
text_box = slide.shapes.add_textbox(Inches(4), Inches(3.5), Inches(8), Inches(2))
text_box.text_frame.word_wrap = True

p = text_box.text_frame.add_paragraph()
p.text = text
p.font.size = Pt(44)
p.font.color.rgb = RGBColor(255, 255, 255)
p.font.name = 'Calibri'  # Clean, modern font

# Center the text box on the slide
text_box.left = int((presentation.slide_width - text_box.width) / 2)
text_box.top = int((presentation.slide_height - text_box.height) / 2)

# Save the presentation
presentation.save("render.pptx")