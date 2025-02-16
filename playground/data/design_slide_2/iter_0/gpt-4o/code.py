from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add the title
title_box = slide.shapes.add_textbox(Inches(8.5), Inches(0.5), Inches(7), Inches(1))
title_frame = title_box.text_frame
p = title_frame.add_paragraph()
p.text = "What is Graphic Design?"
p.font.bold = True
p.font.size = Pt(32)

# Add bullet points
content_box = slide.shapes.add_textbox(Inches(8.5), Inches(2), Inches(7), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

bullet_points = [
    "Graphic design is the art of arranging and making visual substance to impart thoughts and messages.",
    "They rejuvenate visual ideas, most generally through visual communication programming.",
    "Visit the best SEO agency in Dubai; it provides you a good guideline for graphic design and other services."
]

for point in bullet_points:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.space_before = Pt(10)

# Insert the word cloud image
img_path = '../slidesbench/examples/design/slide_2/media/image_0.jpg'
left = Inches(1)
top = Inches(1.75)
height = Inches(6)
pic = slide.shapes.add_picture(img_path, left, top, height=height)

# Add geometric shapes
shapes_colors = [RGBColor(255, 69, 0), RGBColor(60, 179, 113), RGBColor(30, 144, 255)]

for _ in range(3):
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(1.5), Inches(2), Inches(0.5), Inches(0.5)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = shapes_colors.pop(0)

triangle = slide.shapes.add_shape(
    MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(2), Inches(6), Inches(1), Inches(1)
)
triangle.fill.solid()
triangle.fill.fore_color.rgb = RGBColor(255, 0, 255)

# Save the presentation
presentation.save("render.pptx")