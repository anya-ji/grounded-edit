from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Define a slide layout and add a slide
slide_layout = presentation.slide_layouts[5]  # blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Set background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(64, 64, 64)  # dark gray background

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "ZENEROM UAE"
title.font.size = Pt(40)
title.font.color.rgb = RGBColor(255, 255, 255)  # White color

# Add image on the left side
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0.5), Inches(1.6), width=Inches(8))

# Adjust text box size and position on the right side
text_box = slide.shapes.add_textbox(Inches(8.6), Inches(1.6), Inches(6), Inches(5))  # Moved right
text_frame = text_box.text_frame

# Set lighter semi-transparent background for the text box
fill = text_box.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)  # Changed to white background
fill.transparency = 0.2  # Set transparency to 20% for 80% opacity

# Add paragraphs
p1 = text_frame.add_paragraph()
p1.text = ("ZENEROM is a leading ")
run = p1.add_run()
run.text = "digital marketing company"
run.font.color.rgb = RGBColor(0, 102, 204)  # Highlight in blue
p1.add_run().text = (" situated in Dubai. SEO company Dubai provides top SEO "
                     "Google search organic services to clients all over the world. "
                     "We can help your website indexed on top search engines and rate "
                     "higher in order for it to be found.")
p1.font.size = Pt(16)  # Reduced font size
p1.font.color.rgb = RGBColor(0, 0, 0)  # Change to black text

p2 = text_frame.add_paragraph()
p2.text = "Contact us for more details;"
p2.font.size = Pt(16)  # Reduced font size
p2.font.color.rgb = RGBColor(0, 0, 0)  # Change to black text

p3 = text_frame.add_paragraph()
p3.text = "+916282126012"
p3.font.size = Pt(16)  # Reduced font size
p3.font.color.rgb = RGBColor(0, 0, 0)  # Change to black text

# Save the presentation
presentation.save("render.pptx")