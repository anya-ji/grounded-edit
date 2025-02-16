from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_text = title_frame.add_paragraph()
title_text.text = "5 Stages of Development of Media"
title_text.font.bold = True
title_text.font.size = Pt(36)
title_text.font.color.rgb = RGBColor(255, 255, 255)  # White text
title_box.fill.solid()
title_box.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark Blue background

# Add subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(8), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle = subtitle_frame.add_paragraph()
subtitle.text = "3. PRINTING PRESS"
subtitle.font.bold = True
subtitle.font.size = Pt(36)  # Increased font size for better readability
subtitle.font.color.rgb = RGBColor(0, 51, 102)  # Dark Blue text

# Add main content with bullet points
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(8), Inches(5))
content_frame = content_box.text_frame

# Create bullet points with modified formatting
points = [
    "• It allowed the continuous production, reproduction, and",
    "• circulation of print materials."
]

for point in points:
    content = content_frame.add_paragraph()
    content.text = point
    content.font.size = Pt(28)  # Increased font size for better readability
    content.font.color.rgb = RGBColor(0, 0, 0)  # Black text
    content.space_after = Pt(10)  # Set uniform spacing between bullet points

# Add an illustration of a printing press (update the image path)
image_path = "../../media/image_dynamic_pose.jpg"  # Change to the new image path
slide.shapes.add_picture(image_path, Inches(10), Inches(2), Inches(5), Inches(5))

# Save the presentation
presentation.save("render.pptx")