from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Create a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Define slide background as dark blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 51, 102)

# Add title text
title_text = "VIEWING ACTIVITY"
title_box = slide.shapes.add_textbox(Inches(2), Inches(1.5), Inches(12), Inches(2))
title_frame = title_box.text_frame
title_frame.text = title_text
title = title_frame.paragraphs[0]
title.font.bold = True
title.font.size = Pt(60)
title.font.color.rgb = RGBColor(255, 255, 255)

# Add subtitle text
subtitle_text = "Friends | Joey Doesn't Share Food!"
subtitle_box = slide.shapes.add_textbox(Inches(3.5), Inches(3), Inches(9), Inches(1))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.text = subtitle_text
subtitle = subtitle_frame.paragraphs[0]
subtitle.font.size = Pt(32)
subtitle.font.color.rgb = RGBColor(255, 255, 255)

# Insert film camera icon
# You may need to use a local path or ensure the file is in your working directory
# Example icon file name: 'film_camera_icon.png' (ensure it's within the working directory)
film_camera_icon_path = 'film_camera_icon.png'
slide.shapes.add_picture(film_camera_icon_path, Inches(1), Inches(1.5), Inches(1.5), Inches(1.5))

# Insert Fluent logo
# You may need to use a local path or ensure the file is in your working directory
# Example logo file name: 'fluent_logo.png' (ensure it's within the working directory)
fluent_logo_path = 'fluent_logo.png'
slide.shapes.add_picture(fluent_logo_path, Inches(13), Inches(0.5), Inches(2), Inches(1))

# Save the presentation
presentation.save("render.pptx")