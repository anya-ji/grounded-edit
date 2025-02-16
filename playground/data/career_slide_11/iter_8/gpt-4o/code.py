from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide
slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(34, 34, 34)  # Dark background color

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_frame.clear()
title_p = title_frame.paragraphs[0]
title_p.text = "IMPORTANCE OF SKILLS BEYOND ACADEMICS"
title_p.font.size = Pt(40)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(255, 255, 255)  # White color for the title

# Add content
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(15), Inches(7))
content_frame = content_box.text_frame

# Add "Effective Communication" header
p_effective_communication = content_frame.add_paragraph()
p_effective_communication.text = "Effective Communication"
p_effective_communication.font.color.rgb = RGBColor(255, 255, 255)  # White color for the header
p_effective_communication.font.bold = True  # Make the header bold
p_effective_communication.font.size = Pt(28)
p_effective_communication.space_after = Pt(10)  # Space after header

# Add bullet point for "Effective Communication"
p_sub_communication = content_frame.add_paragraph()
p_sub_communication.text = "Effective communication is key in collaborating with others and sharing ideas effectively."
p_sub_communication.level = 1
p_sub_communication.space_before = Pt(5)  # Adjusted space before sub-point for better alignment
p_sub_communication.space_after = Pt(10)  # Increased space after sub-point
p_sub_communication.font.size = Pt(20)
p_sub_communication.font.color.rgb = RGBColor(255, 255, 255)  # White color for sub-points
p_sub_communication.space_after = Pt(15)  # More space after sub-points for clarity
p_sub_communication.font.bold = True  # Make sub-point bold for consistency

# Update other bullet points with the additional sections
bullet_points = [
    ("Problem Solving and Critical Thinking", [
        "Critical thinking skills enable individuals to analyze information, evaluate options, and make informed decisions.",
        "Problem-solving skills are essential for addressing challenges and finding innovative solutions in a wide range of situations."
    ]),
    ("Emotional Intelligence", [
        "Emotional intelligence enhances empathy, interpersonal relationships, and the ability to work effectively in teams."
    ]),
    ("Adaptability and Resilience", [
        "These skills help individuals navigate unexpected setbacks, cope with stress, and embrace change as an opportunity for growth."
    ])
]

for main_point, sub_points in bullet_points:
    p_main = content_frame.add_paragraph()
    p_main.text = main_point
    
    # Change to white color for consistency and make bold
    p_main.font.color.rgb = RGBColor(255, 255, 255)  # White color for the main point
    p_main.font.bold = True  # Make the main point bold
    p_main.font.size = Pt(28)
    p_main.space_after = Pt(10)  # Increased space after main point

    for sp in sub_points:
        p_sub = content_frame.add_paragraph()
        p_sub.text = sp
        p_sub.level = 1
        p_sub.space_before = Pt(5)  # Adjusted space before sub-point for better alignment
        p_sub.space_after = Pt(10)  # Increased space after sub-point
        p_sub.font.size = Pt(20)
        p_sub.font.color.rgb = RGBColor(255, 255, 255)  # White color for sub-points
        p_sub.space_after = Pt(15)  # More space after sub-points for clarity
        
        # Make sub-points bold to match the format of the target image
        p_sub.font.bold = True

# Save presentation
presentation.save("render.pptx")