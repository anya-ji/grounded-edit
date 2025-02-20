from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a Presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color to white
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)

# Add title with a dark blue background
title_text = "Dynamics of Local and Global Culture"
title_box = slide.shapes.add_textbox(Inches(0), Inches(0), presentation.slide_width, Inches(1))
fill = title_box.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark blue

title_frame = title_box.text_frame
title_frame.word_wrap = True

p = title_frame.add_paragraph()
p.text = title_text
p.font.size = Pt(32)  # Changed font size to 32pt
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)  # White text

# Add content text
content_text = (
    "Global flows of culture tend to move more easily around the globe than ever before, "
    "especially through non-material digital forms.\n\n"
    "There are three perspectives on global cultural flows:\n\n"
    "These are cultural differentialism, hybridization, and convergence."
)

content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), presentation.slide_width - Inches(1), Inches(5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

p = content_frame.add_paragraph()
p.text = content_text
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Save the presentation
presentation.save("render.pptx")