from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide to the presentation
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add background image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, 0, 0, presentation.slide_width, presentation.slide_height)

# Add main title textbox
title_textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1))
title_textbox.fill.solid()
title_textbox.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow background
text_frame = title_textbox.text_frame
text_frame.clear()

p = text_frame.add_paragraph()
p.text = "THE CONTENT OF ART"
p.font.bold = True
p.font.size = Pt(36)  # Reduced font size from 44 to 36
p.font.color.rgb = RGBColor(0, 0, 0)  # Black text
text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Add subtitle textbox
subtitle_textbox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(0.7))
subtitle_textbox.fill.solid()
subtitle_textbox.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow background
text_frame_sub = subtitle_textbox.text_frame
text_frame_sub.clear()

p_sub = text_frame_sub.add_paragraph()
p_sub.text = "The Subject and Content of Art"
p_sub.font.bold = True
p_sub.font.size = Pt(20)
p_sub.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add content textbox
content_textbox = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(14), Inches(2))
content_textbox.fill.solid()
content_textbox.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background
text_frame_content = content_textbox.text_frame
text_frame_content.clear()

p_content = text_frame_content.add_paragraph()
p_content.text = "✓ The content of art is the meaning, message, and/or feeling imparted by a work of art."
p_content.font.size = Pt(24)
p_content.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Save the presentation
presentation.save("render.pptx")