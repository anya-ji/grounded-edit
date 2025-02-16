from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from PIL import Image
import io

# Helper function to create a circular image
def create_circular_image(image_path):
    with Image.open(image_path) as img:
        size = min(img.size)
        mask = Image.new('L', img.size, 0)
        draw = ImageDraw.Draw(mask)
        draw.ellipse((0, 0, size, size), fill=255)
        circular_img = Image.new('RGBA', img.size)
        circular_img.paste(img, (0, 0))
        circular_img.putalpha(mask)
        output = io.BytesIO()
        circular_img.thumbnail((Inches(2), Inches(2)), Image.ANTIALIAS)
        circular_img.save(output, format='PNG')
    return output.getvalue()

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]  # Using a blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(47, 47, 47)

# Add circular images to the slide
fruits_vegetables_image_path = "../../media/image_0.jpg"
hanging_clothes_image_path = "../../media/image_1.jpg"

# Prepare circular images
fruits_img = create_circular_image(fruits_vegetables_image_path)
clothes_img = create_circular_image(hanging_clothes_image_path)

# Add the images to slide
slide.shapes.add_picture(io.BytesIO(fruits_img), Inches(1), Inches(1), width=Inches(2), height=Inches(2))
slide.shapes.add_picture(io.BytesIO(clothes_img), Inches(12), Inches(2.5), width=Inches(2), height=Inches(2))

# Add text to the slide
text_box = slide.shapes.add_textbox(Inches(9), Inches(0.5), Inches(6), Inches(1))
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = "From where you would like to buy the following products:"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(255, 255, 255)  # White color for contrast
p.alignment = MSO_ANCHOR.MIDDLE

presentation.save("render.pptx")