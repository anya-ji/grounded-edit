from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(10, 40, 75)  # Dark blue background

# Add the main title
title_text = "Three perspectives on global cultural flows:"
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5))
title_frame = title_box.text_frame
title_frame.text = title_text

title_paragraph = title_frame.paragraphs[0]
title_paragraph.font.size = Pt(32)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
title_frame.word_wrap = True

# Add the first perspective title
perspective_title = "1. CULTURAL DIFFERENTIALISM"
perspective_title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(1))
perspective_title_frame = perspective_title_box.text_frame
perspective_title_frame.text = perspective_title

perspective_title_paragraph = perspective_title_frame.paragraphs[0]
perspective_title_paragraph.font.size = Pt(28)
perspective_title_paragraph.font.bold = True
perspective_title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
perspective_title_frame.word_wrap = True

# Add bullet points
bullet_points = [
    "✓ emphasizes the fact that cultures are essentially different and are only superficially affected by global flows.",
    "It also involves barriers that prevent flows that serve to make cultures more a line; cultures tend to remain stubbornly different from one another."
]

content_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

for point in bullet_points:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)  # White text
    p.space_after = Pt(10)  # Space after each bullet point

presentation.save("render.pptx")