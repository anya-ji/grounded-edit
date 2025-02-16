from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add title
title_text = "Globalization isn’t possible to occur without media."
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1.5))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = title_text
title.font.size = Pt(44)
title.font.bold = True
title.alignment = PP_ALIGN.CENTER

# Add bullet points
content_text = [
    "Electronic media allowed opportunities to spread all over the world.",
    "Radio and Television is a powerful mass medium in providing accessible information for people.",
    "Digital media through phone and computers allows people to access information from around the world."
]

left = Inches(1)
top = Inches(3)
width = Inches(9)
height = Inches(3)

content_box = slide.shapes.add_textbox(left, top, width, height)
content_frame = content_box.text_frame
content_frame.word_wrap = True

for point in content_text:
    p = content_frame.add_paragraph()
    p.text = point
    p.font.size = Pt(24)

# Add image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(10.5), Inches(2.5), Inches(4), Inches(4))

# Save presentation
presentation.save("render.pptx")