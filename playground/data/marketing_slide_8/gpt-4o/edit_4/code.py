from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsdecls
from pptx.oxml import parse_xml

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set background color gradient
background = slide.background
fill = background.fill
fill.solid()
fill_color1, fill_color2 = RGBColor(0, 128, 128), RGBColor(0, 100, 0)
fill.gradient()
fill.gradient_stops[0].position = 0.0
fill.gradient_stops[0].color.rgb = fill_color1
fill.gradient_stops[1].position = 1.0
fill.gradient_stops[1].color.rgb = fill_color2

# Add light bulb graphic
image_path = "../../media/image_0.jpg"
img_left = Inches(1)
img_top = Inches(2)
img_width = Inches(4)
img_height = Inches(4)
light_bulb = slide.shapes.add_picture(image_path, img_left, img_top, img_width, img_height)

# Adding a shape on top of the image to change its color
overlay_rectangle = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    img_left,
    img_top,
    img_width,
    img_height
)
overlay_rectangle.fill.solid()
overlay_rectangle.fill.fore_color.rgb = RGBColor(0, 0, 255)  # More vibrant blue
overlay_rectangle.opacity = 0.5  # Set semi-transparent to see the image underneath
overlay_rectangle.line.color.rgb = RGBColor(255, 255, 255)  # Optional: white border

# Add title text
title_box = slide.shapes.add_textbox(Inches(6), Inches(1), Inches(9), Inches(2))
title_tf = title_box.text_frame
p = title_tf.add_paragraph()
p.text = "THINK LINE:"
p.font.bold = True
p.font.size = Pt(48)
p.font.color.rgb = RGBColor(173, 216, 230)  # Light blue

# Add question text
question_box = slide.shapes.add_textbox(Inches(6), Inches(3), Inches(9), Inches(2))
question_tf = question_box.text_frame
p = question_tf.add_paragraph()
p.text = "Q. How products reach to us?"
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(173, 216, 230)  # Light blue

presentation.save("render.pptx")