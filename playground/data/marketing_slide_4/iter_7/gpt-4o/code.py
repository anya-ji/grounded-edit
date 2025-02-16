from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set background color to red
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red color

# Add black horizontal stripe
stripe_height = Inches(1.5)
stripe_top = (presentation.slide_height - stripe_height) / 2
stripe = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE,
    0,
    stripe_top,
    presentation.slide_width,
    stripe_height
)

stripe.fill.solid()
stripe.fill.fore_color.rgb = RGBColor(0, 0, 0)  # Black color

# Add slide title
title_text = "BIG QUESTION"
title_box = slide.shapes.add_textbox(
    Inches(0.5), stripe_top - Inches(0.5), presentation.slide_width - Inches(1), Inches(1)
)
title_frame = title_box.text_frame
p = title_frame.add_paragraph()
p.text = title_text
p.font.bold = True
p.font.size = Pt(36)
p.font.color.rgb = RGBColor(0, 191, 255)  # Brighter blue color
p.alignment = PP_ALIGN.CENTER

# Add main question centered in bold mixed white font
question_text = "HOW DO MARKETS HELP US?"
question_box = slide.shapes.add_textbox(
    0, stripe_top, presentation.slide_width, stripe_height
)
question_frame = question_box.text_frame
q = question_frame.add_paragraph()
q.text = question_text
q.font.bold = True
q.font.size = Pt(48)

# Change color for main question: all white
q.font.color.rgb = RGBColor(255, 255, 255)  # White color

# Add brown bookmark ribbon in top left corner
ribbon_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.2), Inches(0.2), Inches(1.5), Inches(0.5)
)
ribbon_shape.fill.solid()
ribbon_shape.fill.fore_color.rgb = RGBColor(139, 69, 19)  # Brown color

# Save the presentation
presentation.save("render.pptx")