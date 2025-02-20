from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(13.33)  # Widescreen 16:9
presentation.slide_height = Inches(7.5)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[6]
slide = presentation.slides.add_slide(slide_layout)

# Define background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 128, 128)

# Add title
title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(13.33), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Types of Fixed Retailers"
title.font.bold = True
title.font.size = Pt(44)
title.alignment = 1  # Center alignment
title.font.color.rgb = RGBColor(255, 255, 255)  # Changed to white

# Image file paths
image_paths = [
    "../../media/image_0.jpg",
    "../../media/image_1.jpg",
    "../../media/image_2.jpg"
]

# Labels for images
labels = ["General Stores", "Department Store", "Chain Store"]

# Position settings
image_lefts = [Inches(0), Inches(4.44), Inches(8.88)]
label_top = Inches(2)
image_top = Inches(2.5)
image_height = Inches(4)

# Add images and labels
for i in range(3):
    # Add image
    slide.shapes.add_picture(image_paths[i], image_lefts[i], image_top, height=image_height)
    
    # Add label above image
    label_box = slide.shapes.add_textbox(image_lefts[i], label_top, Inches(4.44), Inches(0.5))
    label_frame = label_box.text_frame
    label = label_frame.add_paragraph()
    label.text = labels[i]
    label.font.size = Pt(24)
    label.font.color.rgb = RGBColor(255, 255, 255)
    label.alignment = 1  # Center alignment

# Save the presentation
presentation.save("render.pptx")