from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Define colors
dark_blue = RGBColor(0, 0, 128)
white = RGBColor(255, 255, 255)

# Add title with dark blue background
title_box = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(16), Inches(1.5))
title_frame = title_box.text_frame
title_frame.clear()
p = title_frame.add_paragraph()
p.text = "5 Stages of Development of Media"
p.font.bold = True
p.font.size = Pt(44)
p.font.color.rgb = white
title_box.fill.solid()
title_box.fill.fore_color.rgb = dark_blue

# Adjusted position for the subheading
subheading_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(15), Inches(1))
subheading_frame = subheading_box.text_frame
subheading_frame.clear()
p = subheading_frame.add_paragraph()
p.text = "5. DIGITAL MEDIA"
p.font.size = Pt(36)
p.font.color.rgb = dark_blue

# Adjusted position for content with bullet points
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(4))
content_frame = content_box.text_frame
content_frame.clear()
p = content_frame.add_paragraph()
p.text = "• Digitized content that can be transmitted over the internet or computer networks."
p.font.size = Pt(24)
p.font.color.rgb = dark_blue

p = content_frame.add_paragraph()
p.text = "• It allows the advertisement of products and online business transactions."
p.font.size = Pt(24)
p.font.color.rgb = dark_blue

# Add image on the right
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(10.5), Inches(2.5), Inches(5), Inches(5))

# Save the presentation
presentation.save("render.pptx")