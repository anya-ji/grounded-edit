from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation and slide
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add black background
slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8), 0, Inches(8), Inches(9)).fill.solid()
slide.shapes[-1].fill.fore_color.rgb = RGBColor(0, 0, 0)

# Add image with a curved shape on the left half
img_path = "../../media/image_0.jpg"
left = Inches(0)
top = Inches(0)
pic = slide.shapes.add_picture(img_path, left, top, width=Inches(8), height=Inches(9))

# ** Removed the circular shape mask on the left side over the image **

# Add title
title_box = slide.shapes.add_textbox(Inches(9), Inches(7), Inches(6), Inches(1))
title_frame = title_box.text_frame
title_p = title_frame.add_paragraph()
title_p.text = "ELO's"
title_p.font.size = Pt(64)
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(255, 255, 255)

# Add subtitle
subtitle_p = title_frame.add_paragraph()
subtitle_p.text = "1. Types of Market"
subtitle_p.font.size = Pt(48)
subtitle_p.font.color.rgb = RGBColor(255, 255, 255)

# Add small red rectangle in the top right corner
rect_left = Inches(15)
rect_top = Inches(0.5)
rect_width = Inches(0.5)
rect_height = Inches(0.5)
rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, rect_left, rect_top, rect_width, rect_height)
rect.fill.solid()
rect.fill.fore_color.rgb = RGBColor(255, 0, 0)

# Add attribution text
attribution_box = slide.shapes.add_textbox(Inches(12), Inches(8.5), Inches(3.5), Inches(0.5))
attribution_frame = attribution_box.text_frame
attribution_p = attribution_frame.add_paragraph()
attribution_p.text = "This Photo by Unknown author is licensed under CC BY-SA"
attribution_p.font.size = Pt(10)
attribution_p.font.color.rgb = RGBColor(255, 255, 255)

# Save the presentation
presentation.save("render.pptx")