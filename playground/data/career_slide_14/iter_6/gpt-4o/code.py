from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]  
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color to dark
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(34, 34, 34)  # Dark grey color

# Title for the slide
title_text = "ACADEMIC WORLD AND PROFESSIONAL WORLD: DEMANDS AND CHARACTERISTICS"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_paragraph = title_frame.add_paragraph()
title_paragraph.text = title_text
title_paragraph.font.size = Pt(32)
title_paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White color
title_paragraph.font.bold = True

# Title for Academic World section
academic_title_text = "ACADEMIC WORLD"
academic_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(7), Inches(0.5))
academic_title_frame = academic_title_box.text_frame
academic_title_paragraph = academic_title_frame.add_paragraph()
academic_title_paragraph.text = academic_title_text
academic_title_paragraph.font.size = Pt(24)
academic_title_paragraph.font.color.rgb = RGBColor(0, 0, 255) # Blue color
academic_title_paragraph.font.bold = True

# Academic World: Goal and Focus
academic_goal_focus_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(7), Inches(3))
academic_goal_focus_frame = academic_goal_focus_box.text_frame

academic_goal_focus_header = academic_goal_focus_frame.add_paragraph()
academic_goal_focus_header.text = "Goal and Focus"
academic_goal_focus_header.font.size = Pt(22)  # Adjusted font size
academic_goal_focus_header.font.color.rgb = RGBColor(255, 255, 0)  # Yellow color
academic_goal_focus_header.font.bold = True

academic_goal = academic_goal_focus_frame.add_paragraph()
academic_goal.text = "Goal: The primary goal of the academic world is to facilitate learning, knowledge acquisition, and intellectual development."
academic_goal.space_after = Pt(10)
academic_goal.font.size = Pt(14)  # Ensure consistent font size
academic_goal.font.color.rgb = RGBColor(255, 255, 0)  # Yellow color for "Goal"

academic_focus = academic_goal_focus_frame.add_paragraph()
academic_focus.text = "Focus: The focus is on education, research, and the pursuit of new knowledge."
academic_focus.font.size = Pt(14)  # Ensure consistent font size
academic_focus.font.color.rgb = RGBColor(255, 255, 0)  # Yellow color for "Focus"

# Academic World: Environment
academic_environment_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(7), Inches(3))
academic_environment_frame = academic_environment_box.text_frame

academic_environment_header = academic_environment_frame.add_paragraph()
academic_environment_header.text = "Environment"
academic_environment_header.font.size = Pt(20)
academic_environment_header.font.color.rgb = RGBColor(255, 255, 0)  # Yellow color
academic_environment_header.font.bold = True

# Updated bullet points to have increased spacing for clarity
academic_setting = academic_environment_frame.add_paragraph()
academic_setting.text = "• Setting: Academics typically take place in educational institutions such as schools, colleges, and universities."
academic_setting.space_after = Pt(15)  # Increased space after
academic_setting.font.size = Pt(14)
academic_setting.font.color.rgb = RGBColor(255, 255, 255)  # White color

academic_atmosphere = academic_environment_frame.add_paragraph()
academic_atmosphere.text = "• Atmosphere: The atmosphere is often centered around learning, exploration, and academic freedom."
academic_atmosphere.space_after = Pt(15)  # Increased space after
academic_atmosphere.font.size = Pt(14)
academic_atmosphere.font.color.rgb = RGBColor(255, 255, 255)  # White color

# Title for Professional World section
professional_title_text = "PROFESSIONAL WORLD"
professional_title_box = slide.shapes.add_textbox(Inches(8.5), Inches(1.5), Inches(7), Inches(0.5))
professional_title_frame = professional_title_box.text_frame
professional_title_paragraph = professional_title_frame.add_paragraph()
professional_title_paragraph.text = professional_title_text
professional_title_paragraph.font.size = Pt(24)
professional_title_paragraph.font.color.rgb = RGBColor(0, 0, 255)  # Consistent Blue color
professional_title_paragraph.font.bold = True

# Professional World: Goal and Focus
professional_goal_focus_box = slide.shapes.add_textbox(Inches(8.5), Inches(2.5), Inches(7), Inches(3))
professional_goal_focus_frame = professional_goal_focus_box.text_frame

professional_goal_focus_header = professional_goal_focus_frame.add_paragraph()
professional_goal_focus_header.text = "Goal and Focus"
professional_goal_focus_header.font.size = Pt(22)  # Adjusted font size
professional_goal_focus_header.font.color.rgb = RGBColor(255, 255, 0)  # Yellow color
professional_goal_focus_header.font.bold = True

professional_goal = professional_goal_focus_frame.add_paragraph()
professional_goal.text = "Goal: The professional world is primarily focused on applying knowledge and skills to achieve specific outcomes or goals."
professional_goal.space_after = Pt(10)
professional_goal.font.size = Pt(14)  # Ensure consistent font size
professional_goal.font.color.rgb = RGBColor(255, 255, 0)  # Yellow color for "Goal"

professional_focus = professional_goal_focus_frame.add_paragraph()
professional_focus.text = "Focus: The focus is on productivity, problem-solving, and delivering results."
professional_focus.font.size = Pt(14)  # Ensure consistent font size
professional_focus.font.color.rgb = RGBColor(255, 255, 0)  # Yellow color for "Focus"

# Professional World: Environment
professional_environment_box = slide.shapes.add_textbox(Inches(8.5), Inches(4), Inches(7), Inches(3))
professional_environment_frame = professional_environment_box.text_frame

professional_environment_header = professional_environment_frame.add_paragraph()
professional_environment_header.text = "Environment"
professional_environment_header.font.size = Pt(20)
professional_environment_header.font.color.rgb = RGBColor(255, 255, 0)  # Yellow color
professional_environment_header.font.bold = True

professional_setting = professional_environment_frame.add_paragraph()
professional_setting.text = "• Setting: Professionals work in various settings, including businesses, government agencies, nonprofits, and industries."
professional_setting.space_after = Pt(15)  # Increased space after for clarity
professional_setting.font.size = Pt(14)
professional_setting.font.color.rgb = RGBColor(255, 255, 255)  # White color

professional_atmosphere = professional_environment_frame.add_paragraph()
professional_atmosphere.text = "• Atmosphere: The atmosphere is goal-oriented, with an emphasis on productivity and meeting objectives."
professional_atmosphere.space_after = Pt(15)  # Increased space after for clarity
professional_atmosphere.font.size = Pt(14)
professional_atmosphere.font.color.rgb = RGBColor(255, 255, 255)  # White color

# Save presentation
presentation.save("render.pptx")