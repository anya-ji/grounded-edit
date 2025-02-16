from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation with specified dimensions
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background color of the slide
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(35, 35, 35)  # dark background

# Add a title text box
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1))
title_frame = title.text_frame
title_frame.word_wrap = True

p = title_frame.add_paragraph()
p.text = "IMPORTANCE OF SKILLS BEYOND ACADEMICS"
p.font.size = Pt(44)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)  # white text

# Add a content text box for bullet points
content = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(15), Inches(6.5))
text_frame = content.text_frame
text_frame.word_wrap = True

# Define bullet points and details with formatting
bullet_points = [
    ("Health and Wellness", (
        "Skills related to physical and mental health, "
        "such as stress management, nutrition, and self-care, "
        "contribute to overall well-being and a higher quality of life."
    )),
    ("Digital Literacy", (
        "With the digitalization of many aspects of life, digital literacy skills are "
        "essential for navigating technology, staying safe online, and using digital "
        "tools for various purposes."
    )),
    ("Networking and Relationship Building", (
        "Building a strong professional and social network is vital for career advancement "
        "and personal growth.\n"
        "Effective networking skills can open doors to opportunities and resources."
    )),
    ("Conflict Resolution", (
        "Conflict is a natural part of life, and the ability to resolve conflicts "
        "constructively is valuable in both personal and professional relationships."
    )),
]

# Add bullet points with highlighted subtitles
for title, content in bullet_points:
    # Sub-title
    p = text_frame.add_paragraph()
    p.text = f"• {title}"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 223, 0)  # change all subtitles to warm yellow

    # Content without indentation
    p = text_frame.add_paragraph()
    p.text = content  # removed extra space before the text
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(211, 211, 211)  # light gray text

# Save the presentation
presentation.save("render.pptx")