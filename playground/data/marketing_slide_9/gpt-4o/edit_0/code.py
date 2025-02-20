from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add title
title_text = "From Factory to the Market"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(15), Inches(1))
title_frame = title_box.text_frame
title_paragraph = title_frame.add_paragraph()
title_paragraph.text = title_text
title_paragraph.font.size = Pt(40)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(34, 34, 34)  # Changed to a darker shade
title_frame.word_wrap = True

# Define step details
steps = [
    ("Producer", "• Goods are produced in large factories or farms.", RGBColor(255, 165, 0)),
    ("Wholesale", "• Wholesalers buy goods in large quantity from factories or farms.", RGBColor(255, 223, 0)),
    ("Distributors", "• They provide goods to retailers on their demand.", RGBColor(255, 223, 0)),
    ("Retailer", "• Sell goods to consumers.", RGBColor(255, 223, 0))
]

# Add rectangles and text for each step
for i, (title, text, color) in enumerate(steps):
    left = Inches(0.5 + i * 4)
    top = Inches(1.5)
    width = Inches(3.5)
    height = Inches(6.5)
    
    shape = slide.shapes.add_shape(
        6,  # Rectangle
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

    # Add title inside shape
    text_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), height - Inches(0.4))
    text_frame = text_box.text_frame
    title_paragraph = text_frame.add_paragraph()
    title_paragraph.text = title
    title_paragraph.font.size = Pt(20)
    title_paragraph.font.bold = True
    title_paragraph.font.color.rgb = RGBColor(0, 0, 0)

    # Add content inside shape
    content_paragraph = text_frame.add_paragraph()
    content_paragraph.text = text
    content_paragraph.font.size = Pt(16)
    content_paragraph.font.color.rgb = RGBColor(0, 0, 0)

# Add background image
background_image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(background_image_path, Inches(0), Inches(0), Inches(16), Inches(9))

# Save the presentation
presentation.save("render.pptx")