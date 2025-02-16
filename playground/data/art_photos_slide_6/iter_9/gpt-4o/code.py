from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]  # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set the swirling brushstroke background
background_fill = slide.background.fill
background_fill.solid()
background_fill.fore_color.rgb = RGBColor(255, 255, 255)  # White background

# Define light yellow color
light_yellow = RGBColor(255, 255, 204)  # Light yellow color

# Add the main title, adjusted position
title_box = slide.shapes.add_textbox(Inches(3.5), Inches(1.5), Inches(9), Inches(1))  # Lowered position
title_box.fill.solid()
title_box.fill.fore_color.rgb = light_yellow  # Change to light yellow
frame = title_box.text_frame
frame.clear()
p = frame.add_paragraph()
p.text = "Representational or Objective Art"
p.font.bold = True
p.font.size = Pt(50)  # Increased text size for visibility
p.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add the smaller title
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
subtitle_box.fill.solid()
subtitle_box.fill.fore_color.rgb = light_yellow  # Change to vibrant yellow
frame = subtitle_box.text_frame
frame.clear()
p = frame.add_paragraph()
p.text = "The Subject and Content of Art"
p.font.bold = True
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Insert the image of Van Gogh's self-portrait
image_path = "../../media/image_1.jpg"
slide.shapes.add_picture(image_path, Inches(10), Inches(1.5), Inches(5), Inches(5))

# Add a caption below the image
caption_box = slide.shapes.add_textbox(Inches(10), Inches(6.6), Inches(5), Inches(0.5))
caption_box.fill.solid()
caption_box.fill.fore_color.rgb = light_yellow  # Change to vibrant yellow
frame = caption_box.text_frame
frame.clear()
p = frame.add_paragraph()
p.text = "Self-portrait of Van Gogh"
p.font.size = Pt(18)
p.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Save the presentation
presentation.save("render.pptx")