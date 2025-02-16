from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from PIL import Image, ImageEnhance

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5]  # blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Set background color to light gray
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(211, 211, 211)  # Light gray color

# Title
title_box = slide.shapes.add_textbox(Inches(4), Inches(0.5), Inches(8), Inches(1))
title_tf = title_box.text_frame
title_tf.text = "Blockchain in Cryptocurrency"
title = title_tf.paragraphs[0]
title.font.bold = True
title.font.size = Pt(44)
title.font.color.rgb = RGBColor(255, 255, 255)  # White color for text
title.alignment = PP_ALIGN.CENTER

# Content
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(10), Inches(5))
content_tf = content_box.text_frame
content_tf.word_wrap = True

bullet_points = [
    "The term blockchain is often used to refer to cryptocurrency. Cryptocurrency is a medium of exchange such as US dollars.",
    "It is just an application in the form of e-currency using blockchain.",
    "It is not governed by any financial institution.",
    "The main difference between blockchain and cryptocurrency is that cryptocurrency is created and held electronically in forms such as a virtual wallet.",
    "It is decentralized and it is not governed by anyone whereas blockchain is an advanced record and it has all information related to cryptocurrency exchanges over a shared system.",
]

for point in bullet_points:
    p = content_tf.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(255, 255, 255)  # White color for text
    p.space_after = Pt(10)  # Add space between bullet points

    # Adjust line spacing for readability
    p.space_before = Pt(5)  # Add space before each bullet point

# Image processing
image_path = "../../media/image_0.jpg"
image = Image.open(image_path)

# Enhance brightness and contrast
enhancer_brightness = ImageEnhance.Brightness(image)
image = enhancer_brightness.enhance(1.2)  # Adjust brightness factor as needed
enhancer_contrast = ImageEnhance.Contrast(image)
image = enhancer_contrast.enhance(1.3)  # Adjust contrast factor as needed

# Save modified image temporarily
modified_image_path = "../../media/modified_image.jpg"
image.save(modified_image_path)

# Add modified image to slide
slide.shapes.add_picture(modified_image_path, Inches(11), Inches(2), Inches(4), Inches(4))

# Save presentation
presentation.save("render.pptx")