from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create presentation
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add slide
slide_layout = presentation.slide_layouts[5]  # blank slide layout
slide = presentation.slides.add_slide(slide_layout)

# Title
title_box = slide.shapes.add_textbox(Inches(4), Inches(0.5), Inches(8), Inches(1))
title_tf = title_box.text_frame
title_tf.text = "Blockchain in Cryptocurrency"
title = title_tf.paragraphs[0]
title.font.bold = True
title.font.size = Pt(44)
title.alignment = PP_ALIGN.CENTER

# Content
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(10), Inches(5))
content_tf = content_box.text_frame
content_tf.word_wrap = True

bullet_points = [
    "The term blockchain is often used to refer to cryptocurrency. Cryptocurrency is a medium of exchange such as US dollars.",
    "It is just an application in the form of e-currency using blockchain.",
    "It is not governed by any financial institution.",
    "The main difference between blockchain and cryptocurrency is that cryptocurrency is created and held electronically in forms such as a virtual wallet.",
    "It is decentralized and it is not governed by anyone whereas blockchain is an advanced record and it has all information related to cryptocurrency exchanges over a shared system.",
]

for point in bullet_points:
    p = content_tf.add_paragraph()
    p.text = point
    p.font.size = Pt(18)
    p.space_before = Pt(6)  # Add space before each bullet point
    p.space_after = Pt(6)   # Add space after each bullet point
    p.line_spacing = Pt(27)  # Set line spacing to 1.5 (18 * 1.5 = 27)

# Image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(11), Inches(2), Inches(4), Inches(4))

# Save presentation
presentation.save("render.pptx")