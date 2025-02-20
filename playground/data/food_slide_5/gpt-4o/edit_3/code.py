from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(245, 245, 245)  # Light background

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "PART 3"
title.font.size = Pt(36)
title.font.bold = True

# Add main instruction text
text_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(10), Inches(2))
text_frame = text_box.text_frame
text = text_frame.add_paragraph()
text.text = (
    "Imagine you went to a restaurant for a first date. Something went wrong involving French fries. "
    "Imagine you’re telling the story to a friend. Use the prompt below to start your story."
)
text.font.size = Pt(18)  # Font size increased to 18pt

# Insert man's picture
man_image_path = "../../media/image_1.jpg"
slide.shapes.add_picture(man_image_path, Inches(0.5), Inches(3), width=Inches(3), height=Inches(3))

# Insert speech bubble with text
speech_box = slide.shapes.add_textbox(Inches(3.5), Inches(3), Inches(8), Inches(2))
speech_frame = speech_box.text_frame
speech_frame.word_wrap = True
speech_text = speech_frame.add_paragraph()
speech_text.text = (
    "We were out to dinner, we were getting along, having a really nice time, "
    "I was thinking she was really cool, but then, out of nowhere, ..."
)
speech_text.font.size = Pt(16)

# Insert French fries image
fries_image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(fries_image_path, Inches(12), Inches(3), width=Inches(3), height=Inches(3))

# Add "PREVIEW ACTIVITY" button
button_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(3), Inches(1))
button_frame = button_box.text_frame
button_text = button_frame.add_paragraph()
button_text.text = "PREVIEW ACTIVITY"
button_text.font.size = Pt(20)
button_text.font.bold = True

# Save the presentation
presentation.save("render.pptx")