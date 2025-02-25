from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5] # blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set gradient background color
background = slide.background
fill = background.fill
fill.gradient()  # Set to gradient
fill.gradient_stops[0].color.rgb = RGBColor(255, 255, 255)  # white
fill.gradient_stops[1].color.rgb = RGBColor(211, 211, 211)  # light gray

# Insert the laptop with paint splashes image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(0.5), Inches(2), height=Inches(5))

# Add the text "Graphic Design"
text_box = slide.shapes.add_textbox(Inches(0), Inches(0.5), Inches(16), Inches(2))
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = "Graphic Design"
p.font.bold = True
p.font.size = Pt(54)
p.font.color.rgb = RGBColor(0, 0, 0)  # black

# Center the text horizontally
text_frame.paragraphs[0].alignment = 1  # Centered

# Save the presentation
presentation.save("render.pptx")