from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# Create a presentation object
presentation = Presentation()

# Set slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add a solid white background
background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(9))
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)  # White color

# Add a white layer for better text visibility
shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(16), Inches(9))
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 255)
shape.opacity = 0.7  # Semi-transparent

# Add the main title
title_text = "D. Mythology and religion, dreams and fantasies."
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1.5))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = title_text
title.font.size = Pt(44)
title.font.bold = True
title.font.color.rgb = RGBColor(0, 0, 0)  # Changed font color to black

# Add the smaller title at the top left corner
subtitle_text = "The Subject and Content of Art"
subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(7), Inches(0.75))
subtitle_frame = subtitle_box.text_frame
subtitle = subtitle_frame.add_paragraph()
subtitle.text = subtitle_text
subtitle.font.size = Pt(24)
subtitle.font.bold = True
subtitle.font.color.rgb = RGBColor(255, 0, 0)  # Contrasting color

# Add the images side by side
img1_path = "../../media/image_1.jpg"
img2_path = "../../media/image_2.jpg"
img_top = Inches(3)
img_left1 = Inches(1)
img_left2 = Inches(8.5)
img_width = Inches(6.5)

slide.shapes.add_picture(img1_path, img_left1, img_top, width=img_width)
slide.shapes.add_picture(img2_path, img_left2, img_top, width=img_width)

# Save the presentation
presentation.save("render.pptx")