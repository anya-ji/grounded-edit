from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the slide background color to dark blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(10, 24, 74)  # Dark blue

# Add the slide title
title_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(15), Inches(1))
title_frame = title_textbox.text_frame
title_frame.word_wrap = True

title_p = title_frame.add_paragraph()
title_p.text = "5 Stages of Development of Media"
title_p.font.size = Pt(36)  # Changed from 44 to 36
title_p.font.bold = True
title_p.font.color.rgb = RGBColor(255, 255, 255)
title_p.alignment = PP_ALIGN.LEFT

# Add the subtitle
subtitle_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(15), Inches(1))
subtitle_frame = subtitle_textbox.text_frame
subtitle_frame.word_wrap = True

subtitle_p = subtitle_frame.add_paragraph()
subtitle_p.text = "1. ORAL COMMUNICATION"
subtitle_p.font.size = Pt(32)
subtitle_p.font.bold = True
subtitle_p.font.color.rgb = RGBColor(10, 24, 74)  # Dark blue
subtitle_p.alignment = PP_ALIGN.LEFT

# Add bullet points
content_textbox = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(10), Inches(4))
content_frame = content_textbox.text_frame
content_frame.word_wrap = True

bullet1 = content_frame.add_paragraph()
bullet1.text = "Language allowed humans to communicate and share information."
bullet1.font.size = Pt(24)
bullet1.font.color.rgb = RGBColor(0, 0, 0)  # Black

bullet2 = content_frame.add_paragraph()
bullet2.text = "Language became the most important tool for exploring the world and different cultures."
bullet2.font.size = Pt(24)
bullet2.font.color.rgb = RGBColor(0, 0, 0)  # Black

# Insert the image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(10.5), Inches(2), Inches(5), Inches(5))

# Save the presentation
presentation.save("render.pptx")