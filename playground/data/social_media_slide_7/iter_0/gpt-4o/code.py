from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Define slide layout
slide_layout = presentation.slide_layouts[5]  # Use a blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add title with dark blue background
title_shape = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(16), Inches(1.5))
title_frame = title_shape.text_frame
title_frame.word_wrap = True
p = title_frame.add_paragraph()
p.text = "5 Stages of Development of Media"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
title_shape.fill.solid()
title_shape.fill.fore_color.rgb = RGBColor(0, 51, 102)  # Dark Blue

# Add subtitle
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(15), Inches(1))
subtitle_frame = subtitle_box.text_frame
p_subtitle = subtitle_frame.add_paragraph()
p_subtitle.text = "4. ELECTRONIC MEDIA"
p_subtitle.font.size = Pt(36)
p_subtitle.font.bold = True

# Add bullet points
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(10), Inches(5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

bullet_points = [
    "It includes the telegraphs, telephone, radio, film and television.",
    "The wide range of these media continue to open up new perspectives on economic, political and cultural processes of globalization."
]

for point in bullet_points:
    p_bullet = content_frame.add_paragraph()
    p_bullet.text = point
    p_bullet.font.size = Pt(24)
    p_bullet.space_after = Pt(10)

# Insert image
image_path = "../../media/image_0.jpg"  # Adjust path as needed
slide.shapes.add_picture(image_path, Inches(11), Inches(2.5), Inches(4), Inches(4))

# Add a shape to cover the watermark
watermark_cover = slide.shapes.add_shape(
    1,  # Shape type: rectangle
    Inches(11), Inches(6.3), Inches(4), Inches(0.5)  # Adjust the position and size to cover watermark
)
watermark_cover.fill.solid()
watermark_cover.fill.fore_color.rgb = RGBColor(255, 255, 255)  # White color

# Save the presentation
presentation.save("render.pptx")