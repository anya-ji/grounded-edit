from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide with a blank layout
slide_layout = presentation.slide_layouts[5] # blank slide
slide = presentation.slides.add_slide(slide_layout)

# Add a title placeholder
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1.5))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Types of Blockchain"
title.font.bold = True
title.font.size = Pt(44)
title.alignment = PP_ALIGN.CENTER

# Add a text box for bullet points
content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(6))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Add bullet points with checked checkboxes
bullet_points = [
    ("☑ Private Blockchain Networks: ", 
     "Private blockchains operate on closed networks, and tend to work well for private businesses and organizations."),
    ("☑ Public Blockchain Networks: ", 
     "Bitcoin and other cryptocurrencies originated from public blockchains, which also played a role in popularizing distributed ledger technology (DLT)."),
    ("☑ Permissioned Blockchain Networks: ", 
     "Also sometimes known as hybrid blockchains, permissioned blockchain networks are private blockchains that allow special access for authorized individuals."),
    ("☑ Hybrid Blockchains: ", 
     "Hybrid blockchains are the combination of both public and private blockchains. In a hybrid blockchain, some parts of the blockchain are public and transparent, while others are private and accessible only to authorized and specific participants."),
]

for bullet_title, bullet_text in bullet_points:
    p = content_frame.add_paragraph()
    p.text = bullet_title + bullet_text
    p.font.size = Pt(20)
    p.space_after = Pt(14) # Add spacing between bullet points

# Save the presentation
presentation.save("render.pptx")