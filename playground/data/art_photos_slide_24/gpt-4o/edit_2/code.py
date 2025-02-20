from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the background image
image_path = "../../media/image_0.jpg"
background = slide.shapes.add_picture(image_path, 0, 0, width=presentation.slide_width, height=presentation.slide_height)

# Add title textbox
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1.5))
title_frame = title_box.text_frame
title_frame.word_wrap = True

title_paragraph = title_frame.add_paragraph()
title_paragraph.text = "KEEPING ART"
title_paragraph.font.size = Pt(44)
title_paragraph.font.bold = True
title_paragraph.font.color.rgb = RGBColor(0, 0, 0)

# Set the background color of the title textbox
fill = title_box.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 0)

# Add subtitle textbox
subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(15), Inches(1.5))
subtitle_frame = subtitle_box.text_frame
subtitle_frame.word_wrap = True

subtitle_paragraph = subtitle_frame.add_paragraph()
subtitle_paragraph.text = "3. Preservation and Restoration"
subtitle_paragraph.font.size = Pt(36)
subtitle_paragraph.font.bold = True
subtitle_paragraph.font.color.rgb = RGBColor(0, 0, 0)

# Set the background color of the subtitle textbox
subtitle_fill = subtitle_box.fill
subtitle_fill.solid()
subtitle_fill.fore_color.rgb = RGBColor(255, 255, 0)

# Add main content textbox
content_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(15), Inches(4))
content_frame = content_box.text_frame
content_frame.word_wrap = True

content_paragraph = content_frame.add_paragraph()
content_paragraph.text = "✓ Because art is valuable, enormous human efforts and financial resources are devoted to preserving art from the ravages of time, the environment, industrial-by-product, and even any other human being."
content_paragraph.font.size = Pt(24)  # Font size increased to 24 pt
content_paragraph.font.color.rgb = RGBColor(0, 0, 0)

# Set the background color of the content textbox
content_fill = content_box.fill
content_fill.solid()
content_fill.fore_color.rgb = RGBColor(255, 255, 255)

# Save the presentation
presentation.save("render.pptx")