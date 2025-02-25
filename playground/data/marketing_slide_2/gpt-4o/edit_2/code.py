from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a slide
slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Set slide background color to blue
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(0, 0, 255)  # Blue

# Main question text (modified to title case)
textbox = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(10), Inches(1.5))
text_frame = textbox.text_frame
p = text_frame.add_paragraph()
p.text = "What is Filling Your Bucket Today and What’s Draining It?"  # Title case
p.font.size = Pt(40)
p.font.color.rgb = RGBColor(255, 255, 255)  # White

# Add image of bucket being filled
image_path = "../../media/image_1.jpg"
slide.shapes.add_picture(image_path, Inches(10.5), Inches(1.5), Inches(5), Inches(5))

# Add "Filling" label
filling_label = slide.shapes.add_textbox(Inches(14), Inches(1), Inches(2), Inches(0.5))
text_frame = filling_label.text_frame
p = text_frame.add_paragraph()
p.text = "Filling"
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(255, 255, 255)  # White

# Add "Draining" label
draining_label = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(2), Inches(0.5))
text_frame = draining_label.text_frame
p = text_frame.add_paragraph()
p.text = "Draining"
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(255, 255, 255)  # White

# Add the Pear Deck text and QR code image
image_path_peardeck = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path_peardeck, Inches(0.5), Inches(7.5), Inches(15), Inches(1.5))

# Save the presentation
presentation.save("render.pptx")