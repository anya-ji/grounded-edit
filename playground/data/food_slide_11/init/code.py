from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

# Create a presentation object
presentation = Presentation()

# Define slide layout and dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add title
title = "PART 3"
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(5), Inches(1))
title_frame = title_box.text_frame
title_frame.text = title
title_frame.paragraphs[0].font.size = Pt(36)
title_frame.paragraphs[0].font.bold = True

# Add instructions
instruction = "Choose true (T), false (F), or not given (N) according to the information in the video."
instruction_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(10), Inches(0.5))
instruction_frame = instruction_box.text_frame
instruction_frame.text = instruction
instruction_frame.paragraphs[0].font.size = Pt(18)

# Define table properties
rows, cols = 4, 4
left, top, width, height = Inches(0.5), Inches(2.5), Inches(9), Inches(2)
table = slide.shapes.add_table(rows, cols, left, top, width, height).table

# Set column widths
table.columns[0].width = Inches(5)
table.columns[1].width = Inches(1.33)
table.columns[2].width = Inches(1.33)
table.columns[3].width = Inches(1.33)

# Add table headers
headers = ["", "T", "F", "N"]
for col, header in enumerate(headers):
    cell = table.cell(0, col)
    cell.text = header
    for paragraph in cell.text_frame.paragraphs:
        paragraph.font.bold = True
        paragraph.alignment = PP_ALIGN.CENTER

# Add statements
statements = [
    "1. Phoebe sarcastically says that Sarah is a monster for her etiquette.",
    "2. Joey prefers eating French fries with his fingers.",
    "3. Rachel is surprised to hear about Joey’s food sharing rule."
]

for row, statement in enumerate(statements, start=1):
    table.cell(row, 0).text = statement

# Add checkbox image
checkbox_img_path = "../../media/image_1.jpg"

for row in range(1, rows):
    for col in range(1, cols):
        placeholder = slide.shapes.add_picture(checkbox_img_path, Inches(7 + (col-1)*1.33), Inches(2 + row*0.5), Inches(0.3), Inches(0.3))

# Add image of Joey
joey_img_path = "../../media/image_0.jpg"
slide.shapes.add_picture(joey_img_path, Inches(10), Inches(2.5), Inches(5), Inches(5))

# Add "VIEWING ACTIVITY" button
button_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(2.5), Inches(0.5))
button_frame = button_box.text_frame
button_p = button_frame.add_paragraph()
button_p.text = "VIEWING ACTIVITY"
button_p.font.size = Pt(16)
button_box.fill.solid()
button_box.fill.fore_color.rgb = RGBColor(192, 192, 192)
button_p.alignment = PP_ALIGN.CENTER

# Save presentation
presentation.save("render.pptx")