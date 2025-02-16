from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Add title
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(5), Inches(1))
title_frame = title_box.text_frame
title = title_frame.add_paragraph()
title.text = "Biography"
title.font.bold = True
title.font.size = Pt(44)

# Add details
details = [
    "Name: Elon Reeve Musk",
    "Born: 28 June 1971",
    "Origin: South African",
    "Education: University of Pennsylvania",
    "Spouses: Talulah Riley",
    "Residence: Los Angeles, USA",
    "Occupation: An Entrepreneur, Engineer, Investor"
]

top_margin = 1.5  # Start below the title

for detail in details:
    text_box = slide.shapes.add_textbox(Inches(0.5), Inches(top_margin), Inches(7.5), Inches(0.5))
    text_frame = text_box.text_frame
    p = text_frame.add_paragraph()
    p.text = detail
    p.font.size = Pt(20)

    # Add blue line separator
    line = slide.shapes.add_line(Inches(0.5), Inches(top_margin + 0.5), Inches(8), Inches(top_margin + 0.5))
    line.line.color.rgb = RGBColor(0, 0, 255)

    top_margin += 0.8  # Adjust for next line

# Insert portrait image within a white frame
image_path = "../../media/image_0.jpg"
pic = slide.shapes.add_picture(image_path, Inches(9), Inches(1), width=Inches(6))

# Add yellow vertical bar for emphasis
left = Inches(15.5)
top = Inches(0)
height = Inches(9)
width = Inches(0.5)

shape = slide.shapes.add_shape(
    1, left, top, width, height
)  # 1 is for rectangle MSO_SHAPE.MSO_SHAPE_RECTANGLE
fill = shape.fill
fill.solid()
fill.fore_color.rgb = RGBColor(255, 255, 0)

presentation.save("render.pptx")