from pptx import Presentation
from pptx.util import Inches, Pt, Cm
from pptx.dml.color import RGBColor

presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

slide_layout = presentation.slide_layouts[5]  # Blank layout
slide = presentation.slides.add_slide(slide_layout)

# Title on the left side
left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(7.5), Inches(2))
left_title_frame = left_title_box.text_frame
left_title_frame.word_wrap = True
left_title_box.fill.solid()
left_title_box.fill.fore_color.rgb = RGBColor(0, 0, 255)  # Blue background

left_title_p = left_title_frame.add_paragraph()
left_title_p.text = "principles of graphic design"  # Changed to lowercase
left_title_p.font.size = Pt(60)  # Increased font size
left_title_p.font.bold = True
left_title_p.font.color.rgb = RGBColor(255, 255, 255)  # Change text color to white

# Right side text box with subtitle and list
right_text_box = slide.shapes.add_textbox(Inches(8.5), Inches(0.5), Inches(7), Inches(7))
right_text_frame = right_text_box.text_frame
right_text_frame.word_wrap = True
right_text_box.fill.solid()
right_text_box.fill.fore_color.rgb = RGBColor(173, 216, 230)  # Light blue background

right_title_p = right_text_frame.add_paragraph()
right_title_p.text = "The principles of graphic design related to the areas are:"
right_title_p.font.size = Pt(24)

# List of items with varying background shades
items = [
    ("Arrangement", RGBColor(173, 216, 230)),  # Light blue
    ("Proximity", RGBColor(152, 251, 152)),   # Light green
    ("Repetition", RGBColor(135, 206, 235)),  # Sky blue
    ("Contrast", RGBColor(144, 238, 144)),    # Light green
    ("Balance", RGBColor(176, 224, 230))      # Powder blue
]

start_y = 2.0  # Start below the title
for i, (text, color) in enumerate(items):
    y_position = start_y + i * 1.2  # Increment for each item
    item_box = slide.shapes.add_textbox(Inches(8.5), Inches(y_position), Inches(7), Inches(1))
    item_box.fill.solid()
    item_box.fill.fore_color.rgb = color
    
    item_frame = item_box.text_frame
    item_p = item_frame.add_paragraph()
    item_p.text = text
    item_p.font.size = Pt(22)
    item_p.font.bold = True

# Save the presentation
presentation.save("render.pptx")