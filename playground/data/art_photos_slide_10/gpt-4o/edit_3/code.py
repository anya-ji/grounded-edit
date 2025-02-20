from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()

# Slide dimensions
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add background image
background_image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(background_image_path, 0, 0, presentation.slide_width, presentation.slide_height)

# Add a text box for the small title
small_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(10), Inches(1))
small_title_frame = small_title_box.text_frame
small_title_frame.text = "The Subject and Content of Art"
small_title_box.fill.solid()
small_title_box.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow background
small_title_box.text_frame.paragraphs[0].font.size = Pt(24)
small_title_box.text_frame.paragraphs[0].font.bold = True
small_title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add a text box for the main title "Portrait"
main_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.5), Inches(15), Inches(1))
main_title_frame = main_title_box.text_frame
main_title_frame.text = "Portrait"
main_title_box.fill.solid()
main_title_box.fill.fore_color.rgb = RGBColor(255, 255, 0)  # Yellow background
main_title_box.text_frame.paragraphs[0].font.size = Pt(36)  # Changed font size to 36pt
main_title_box.text_frame.paragraphs[0].font.bold = True
main_title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Add images
mona_lisa_image_path = "../../media/image_1.jpg"
child_portrait_image_path = "../../media/image_2.jpg"

left_image = slide.shapes.add_picture(mona_lisa_image_path, Inches(0.5), Inches(1.5), width=Inches(7))
right_image = slide.shapes.add_picture(child_portrait_image_path, Inches(8.5), Inches(1.5), width=Inches(7))

# Save the presentation
presentation.save("render.pptx")