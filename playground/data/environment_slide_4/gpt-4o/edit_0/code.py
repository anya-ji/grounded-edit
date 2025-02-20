from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a presentation with a custom slide size
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5] 
slide = presentation.slides.add_slide(slide_layout)

# Set a lighter green background color
background = slide.background
fill = background.fill
fill.solid()
fill.fore_color.rgb = RGBColor(204, 255, 204)  # light green

# Add a white textbox in the center
left = Inches(3)
top = Inches(1.5)
width = Inches(10)
height = Inches(6)
text_box = slide.shapes.add_textbox(left, top, width, height)
text_box.fill.solid()
text_box.fill.fore_color.rgb = RGBColor(255, 255, 255)  # white

# Add title text
text_frame = text_box.text_frame
p = text_frame.add_paragraph()
p.text = "NATURAL ENVIRONMENT"
p.font.bold = True
p.font.size = Pt(44)
p.font.color.rgb = RGBColor(0, 128, 0)  # green
p.alignment = PP_ALIGN.CENTER

# Add bullet point text
bullet_points = [
    "The natural environment means the non human made surroundings conditions in which all biotic and abiotic things exist on earth.",
    "Natural environment can further be divided into 4 parts."
]

for bullet in bullet_points:
    p = text_frame.add_paragraph()
    p.text = bullet
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(0, 0, 0)  # black
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(12)  # spacing between bullets

# Add image below the text
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(5.5), Inches(4.5), Inches(5), Inches(3.5))

# Save the presentation
presentation.save("render.pptx")