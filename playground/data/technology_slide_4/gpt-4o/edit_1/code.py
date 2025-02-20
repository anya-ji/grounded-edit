from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a PowerPoint presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Title text
title_text = "Types of Blockchain"
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame

# Title formatting
title_p = title_frame.add_paragraph()
title_p.text = title_text
title_p.font.size = Pt(50)  # Increased font size
title_p.font.bold = True
title_p.alignment = PP_ALIGN.CENTER

# Content text
content_text = (
    "1. Private Blockchain Networks: Private blockchains operate on closed networks, and tend to work well for private businesses and organizations.\n"
    "2. Public Blockchain Networks: Bitcoin and other cryptocurrencies originated from public blockchains, which also played a role in popularizing distributed ledger technology (DLT).\n"
    "3. Permissioned Blockchain Networks: Also sometimes known as hybrid blockchains, permissioned blockchain networks are private blockchains that allow special access for authorized individuals.\n"
    "4. Hybrid Blockchains: Hybrid blockchains are the combination of both public and private blockchains. In a hybrid blockchain, some parts of the blockchain are public and transparent, while others are private and accessible only to authorized and specific participants."
)

# Add content as bullet points
content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(6.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Split content into individual bullet points
bullet_points = content_text.split("\n")

# Add each bullet point
for bullet in bullet_points:
    bullet_p = content_frame.add_paragraph()
    bullet_p.text = bullet
    bullet_p.font.size = Pt(24)
    bullet_p.font.underline = True if ':' in bullet else False  # Underline topic
    bullet_p.space_after = Pt(12)

# Save the presentation
presentation.save("render.pptx")