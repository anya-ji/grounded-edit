from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Create a PowerPoint presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Title text
title_text = "Types of Blockchain"
title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(14), Inches(1))
title_frame = title_box.text_frame

# Title formatting
title_p = title_frame.add_paragraph()
title_p.text = title_text
title_p.font.size = Pt(24)  # Changed font size to match content
title_p.font.bold = True
title_p.alignment = PP_ALIGN.CENTER

# Content text with periods
content_text = (
    "Private Blockchain Networks: Private blockchains operate on closed networks, and tend to work well for private businesses and organizations.\n"
    "Public Blockchain Networks: Bitcoin and other cryptocurrencies originated from public blockchains, which also played a role in popularizing distributed ledger technology (DLT).\n"
    "Permissioned Blockchain Networks: Also sometimes known as hybrid blockchains, permissioned blockchain networks are private blockchains that allow special access for authorized individuals.\n"
    "Hybrid Blockchains: Hybrid blockchains are the combination of both public and private blockchains. In a hybrid blockchain, some parts of the blockchain are public and transparent, while others are private and accessible only to authorized and specific participants."
)

# Add content as bullet points
content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(14), Inches(6.5))
content_frame = content_box.text_frame
content_frame.word_wrap = True

# Split content into individual bullet points
bullet_points = content_text.split("\n")

# Add each bullet point as standard bullet points
for bullet in bullet_points:
    bullet_p = content_frame.add_paragraph()
    bullet_p.text = bullet  # Removed checkmark symbol
    bullet_p.font.size = Pt(24)
    bullet_p.space_after = Pt(12)
    bullet_p.level = 0  # Set bullet point level to 0

# Save the presentation
presentation.save("render.pptx")