from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Choose a blank layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Add title with light gray background
title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(15), Inches(1.0))
title_text_frame = title_box.text_frame
title_text_frame.word_wrap = True
title_paragraph = title_text_frame.add_paragraph()
title_paragraph.text = "Branding"
title_paragraph.font.bold = True
title_paragraph.font.size = Pt(44)
title_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Set background color to light gray
fill = title_box.fill
fill.solid()
fill.fore_color.rgb = RGBColor(211, 211, 211)  # Light gray

# Insert the image
image_path = "../../media/image_0.jpg"
slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), Inches(14), Inches(5))

# Add main content text (adjusted y-coordinate for the margin)
content_box = slide.shapes.add_textbox(Inches(1), Inches(7 + 1.5 + 5 + 0.020), Inches(14), Inches(2))
content_frame = content_box.text_frame
content_frame.word_wrap = True
content_paragraph = content_frame.add_paragraph()
content_paragraph.text = (
    "'Belong Anywhere’ became our compass. The ‘Belo’ symbol - a sort of heart/drop pin "
    "hybrid that can be drawn by anyone - transcends language and culture in an ode to "
    "Airbnb’s welcoming spirit, while a warm and welcoming new colour palette reflects the "
    "open heartedness of its people.\n\n"
    "Source: Airbnb and Design Studio, 2017"
)
content_paragraph.font.size = Pt(18)
content_paragraph.font.color.rgb = RGBColor(0, 0, 0)  # Black text

# Save the presentation
presentation.save("render.pptx")