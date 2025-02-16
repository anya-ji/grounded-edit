from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor

# Create a presentation object
presentation = Presentation()
presentation.slide_width = Inches(16)
presentation.slide_height = Inches(9)

# Add a blank slide layout
slide_layout = presentation.slide_layouts[5]
slide = presentation.slides.add_slide(slide_layout)

# Set the title text
title_text = "Home Fun:"
question_text = "Q. Compare the different types of Markets."

# Add background with a gradient fill
background = slide.background
fill = background.fill
fill.solid()

# Define the dark teal color for background
dark_teal = RGBColor(3, 83, 89)
fill.fore_color.rgb = dark_teal

# Create a centered text box for title
title_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(14), Inches(1.5))
title_box.fill.solid()
title_box.fill.fore_color.rgb = dark_teal
frame = title_box.text_frame
frame.word_wrap = True
p = frame.add_paragraph()
p.text = title_text
p.font.bold = True
p.font.size = Pt(54)
p.font.color.rgb = RGBColor(255, 0, 0)  # Red color
frame.alignment = MSO_SHAPE.CENTER

# Create a centered text box for question
question_box = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(14), Inches(1.5))
question_box.fill.solid()
question_box.fill.fore_color.rgb = dark_teal
q_frame = question_box.text_frame
q_frame.word_wrap = True
p = q_frame.add_paragraph()
p.text = question_text
p.font.size = Pt(30)
p.font.color.rgb = RGBColor(255, 255, 255)  # White color
q_frame.alignment = MSO_SHAPE.CENTER

# Add a thick red border around the content area
border_shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.5), Inches(15), Inches(8)
)
border_shape.line.color.rgb = RGBColor(255, 0, 0)  # Red color
border_shape.line.width = Pt(6)  # Thick border

# Overlap the border behind the text boxes
border_shape.fill.background()

# Save the presentation to the specified path
presentation.save("render.pptx")